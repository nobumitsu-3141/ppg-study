#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""研究1（VitalDB 862例）の結果報告スライドを川副式書式で生成する。

書式ルール（slide-format スキル）:
  - タイトル44pt太字・金BF9000・黒縁取り2.25pt・全スライド同位置・1行
  - タイトル直下に金色下線（y=1.52in・太さ8pt・全幅）
  - 右上に章ナビ帯（現在章のみティール00A8AA、他は灰D9D9D9）
  - 本文は22pt以上（出典16pt・章ナビ11ptのみ例外）
  - 図解優先・詳細はノートへ・図と文字を重ねない
  - 対比色はブルー0072B2 × バーミリオンD55E00（＋ティール）。金は構造色として予約

使い方: python3 build_slides_ja.py 出力.pptx
"""
from __future__ import annotations

import sys

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

GOLD, INK, WHITE = "BF9000", "1A1A1A", "FFFFFF"
BLUE, VERM, TEAL = "0072B2", "D55E00", "00A8AA"
RED, GREY, LGREY = "C00000", "808080", "D9D9D9"
FONT = "メイリオ"
SW, SH = 13.333, 7.5

CHAPTERS = ["背景", "方法", "結果", "考察"]


def _rgb(h):
    return RGBColor.from_string(h)


def _units(text: str) -> float:
    """lint と同じ文字幅ユニット計算（全角=1.0・ASCII=0.56）。"""
    u = 0.0
    for ch in text:
        o = ord(ch)
        if ch == " ":
            u += 0.30
        elif o <= 0x24F or 0x2080 <= o <= 0x208E:
            u += 0.56
        else:
            u += 1.0
    return u


def width_in(text: str, pt: float) -> float:
    """lint 準拠の推定描画幅（インチ）。Meiryo安全係数1.12込み。"""
    return _units(text) * (pt / 72.0) * 1.12


def outline_run(run, w_pt=2.25, color=INK):
    """run に文字の縁取り（a:ln）を付ける。rPr の先頭に入れる必要がある。"""
    rPr = run.font._rPr
    ln = rPr.makeelement(qn("a:ln"), {"w": str(int(w_pt * 12700))})
    fill = ln.makeelement(qn("a:solidFill"), {})
    clr = fill.makeelement(qn("a:srgbClr"), {"val": color})
    fill.append(clr)
    ln.append(fill)
    rPr.insert(0, ln)


def set_text(tf, lines, size=24, color=INK, bold=False, align=PP_ALIGN.LEFT,
             space_after=6, line_spacing=1.2):
    """テキストフレームに行を流し込む（各行=1段落・1行に収まる前提）。"""
    tf.word_wrap = True
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.line_spacing = line_spacing
        txt, sz, col, bd = ln if isinstance(ln, tuple) else (ln, size, color, bold)
        r = p.add_run()
        r.text = txt
        r.font.size = Pt(sz)
        r.font.bold = bd
        r.font.name = FONT
        r.font.color.rgb = _rgb(col)


def textbox(slide, x, y, w, h, lines, **kw):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.text_frame.margin_left = Inches(0.05)
    tb.text_frame.margin_right = Inches(0.05)
    tb.text_frame.margin_top = Inches(0.02)
    tb.text_frame.margin_bottom = Inches(0.02)
    set_text(tb.text_frame, lines, **kw)
    return tb


def box(slide, x, y, w, h, lines, fill=None, line=None, shape=MSO_SHAPE.ROUNDED_RECTANGLE,
        anchor=MSO_ANCHOR.MIDDLE, **kw):
    """塗り箱の中に直接テキストを入れる（別テキストボックスを重ねない）。"""
    sp = slide.shapes.add_shape(shape, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = _rgb(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = _rgb(line)
        sp.line.width = Pt(1.75)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = Inches(0.10)
    tf.margin_right = Inches(0.10)
    tf.margin_top = Inches(0.04)
    tf.margin_bottom = Inches(0.04)
    set_text(tf, lines, **kw)
    return sp


def arrow(slide, x, y, w, h, color=GREY):
    sp = slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = _rgb(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def down_arrow(slide, x, y, w, h, color=GREY):
    sp = slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    sp.fill.solid()
    sp.fill.fore_color.rgb = _rgb(color)
    sp.line.fill.background()
    sp.shadow.inherit = False
    return sp


def new_slide(prs, title, chapter=None, notes=""):
    """タイトル＋金下線＋章ナビ帯の共通ヘッダを持つスライドを作る。"""
    s = prs.slides.add_slide(prs.slide_layouts[5])   # Title Only
    t = s.shapes.title
    t.left, t.top, t.width, t.height = Inches(0.55), Inches(0.35), Inches(9.25), Inches(0.95)
    tf = t.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    r.font.size = Pt(44)
    r.font.bold = True
    r.font.name = FONT
    r.font.color.rgb = _rgb(GOLD)
    outline_run(r)

    ln = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.52), Inches(SW), Pt(0))
    ln.fill.solid()
    ln.fill.fore_color.rgb = _rgb(GOLD)
    ln.line.color.rgb = _rgb(GOLD)
    ln.line.width = Pt(8)
    ln.shadow.inherit = False

    if chapter:
        cw, gap = 0.68, 0.06
        n = len(CHAPTERS)
        x0 = SW - 0.30 - (n * cw + (n - 1) * gap)
        for i, ch in enumerate(CHAPTERS):
            cur = ch == chapter
            c = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   Inches(x0 + i * (cw + gap)), Inches(0.06),
                                   Inches(cw), Inches(0.64))
            c.fill.solid()
            c.fill.fore_color.rgb = _rgb(TEAL if cur else LGREY)
            c.line.fill.background()
            c.shadow.inherit = False
            ctf = c.text_frame
            ctf.word_wrap = False
            ctf.vertical_anchor = MSO_ANCHOR.MIDDLE
            ctf.margin_left = ctf.margin_right = Inches(0.01)
            cp = ctf.paragraphs[0]
            cp.alignment = PP_ALIGN.CENTER
            cr = cp.add_run()
            cr.text = ch
            cr.font.size = Pt(11)
            cr.font.bold = cur
            cr.font.name = FONT
            cr.font.color.rgb = _rgb(WHITE if cur else "595959")
    if notes:
        s.notes_slide.notes_text_frame.text = notes
    return s


def source(slide, text):
    textbox(slide, 0.5, 6.92, 12.3, 0.42, [text], size=16, color=GREY)


def check_title(title):
    w = width_in(title, 44)
    assert w <= 9.3, f"タイトルが幅超過 {w:.2f}in > 9.3in: {title}"


def build(out_path):
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(SW), Inches(SH)

    # ---------------------------------------------------------- 0 表紙
    s = prs.slides.add_slide(prs.slide_layouts[6])
    tb = textbox(s, 0.7, 1.9, 11.9, 1.5, [
        ("脈波伝播時間の変動は", 40, GOLD, True),
        ("血管指標では説明されない", 40, GOLD, True),
    ], line_spacing=1.15)
    for p in tb.text_frame.paragraphs:
        for r in p.runs:
            outline_run(r)
    box(s, 0.7, 3.6, 11.9, 0.62, [
        ("― esCCO較正定数の血管補正という発想の前提を検証する ―", 24, INK, False)],
        fill=None, line=None, align=PP_ALIGN.LEFT)
    box(s, 0.7, 4.45, 11.9, 0.95, [
        ("公開周術期データベース VitalDB 862例", 26, WHITE, True),
        ("事前登録・参照CO非依存の確認的解析", 24, WHITE, False)],
        fill=BLUE, align=PP_ALIGN.CENTER, space_after=2)
    textbox(s, 0.7, 5.75, 11.9, 0.9, [
        ("川副 靖晃", 26, INK, True),
        ("五島中央病院 麻酔科", 24, INK, False)], space_after=2)
    s.notes_slide.notes_text_frame.text = (
        "研究1の結果報告。VitalDB 874例から862例(98.6%)・161,737ウィンドウを解析した。\n"
        "統計解析計画は波形を見る前に凍結し、Zenodo(doi:10.5281/zenodo.22167118)に"
        "第三者タイムスタンプを取得している。")

    # ---------------------------------------------------------- 1 背景
    t = "esCCOとPWTT法"
    s = new_slide(prs, t, "背景", notes=(
        "esCCOは心電図と光電容積脈波だけで連続的に心拍出量を推定する実用化された方法。"
        "R波から脈波の足までの区間（PWTT）が1回拍出量と逆相関することを利用する。"
        "多施設検証では臨床的なトレンド追随として十分な性能を示し、開発者自身も"
        "絶対値ではなくトレンドモニタと位置づけている。"))
    box(s, 0.6, 2.0, 3.5, 1.15, [("心電図 R波", 26, WHITE, True)], fill=BLUE)
    arrow(s, 4.3, 2.35, 0.9, 0.45)
    box(s, 5.4, 2.0, 3.5, 1.15, [("脈波の足", 26, WHITE, True)], fill=VERM)
    box(s, 9.1, 2.0, 3.6, 1.15, [("PWTT", 26, WHITE, True)], fill=TEAL)
    down_arrow(s, 10.6, 3.3, 0.6, 0.6)
    box(s, 6.6, 4.1, 6.1, 1.0, [("1回拍出量 → 心拍出量", 26, INK, True)],
        fill="F2F2F2", line=TEAL)
    textbox(s, 0.6, 4.0, 5.6, 2.5, [
        ("必要なのは既に全例で記録済みの", 24, INK, False),
        ("心電図と脈波の2つだけ", 24, BLUE, True),
        ("", 24, INK, False),
        ("追加の侵襲・追加の機器が不要", 24, INK, False),
    ], space_after=4)
    source(s, "Ishihara 2004; Sugo 2010; Yamada 2012")

    # ---------------------------------------------------------- 2 背景
    t = "誤差は血管状態と相関"
    s = new_slide(prs, t, "背景", notes=(
        "esCCOの誤差はランダムではない。体血管抵抗や実効動脈エラスタンスに関連して"
        "系統的に悪化する。一方、較正定数Kは年齢・性別・身長・体重から決まり、"
        "いったん較正すると症例中固定される。つまり血管状態の情報を持たず、その変化を"
        "追えない。較正手順の改善だけではこの方法は直らないことも報告されている"
        "(Smetkin 2017)。"))
    box(s, 0.6, 2.0, 5.9, 2.0, [
        ("誤差は血管状態に依存", 26, WHITE, True),
        ("体血管抵抗・動脈エラスタンス", 24, WHITE, False),
    ], fill=VERM, space_after=6)
    box(s, 6.8, 2.0, 5.9, 2.0, [
        ("較正定数 K は固定", 26, WHITE, True),
        ("年齢・性別・身長・体重のみ", 24, WHITE, False),
    ], fill=BLUE, space_after=6)
    down_arrow(s, 6.35, 4.15, 0.65, 0.6)
    box(s, 1.8, 4.95, 9.7, 1.15, [
        ("K は血管情報を持たず、その変化を追えない", 28, WHITE, True)],
        fill=RED)
    source(s, "Ishihara & Tsutsui 2014; Biais 2015; Magliocca 2018; Smetkin 2017")

    # ---------------------------------------------------------- 3 背景（核心）
    t = "検証されていない前提"
    s = new_slide(prs, t, "背景", notes=(
        "血管指標でKを補正するという発想は、ひとつの検証されていない前提の上に立つ。"
        "すなわち『補正したいPWTTの変動が、その血管指標の測る血管状態に由来している』"
        "という前提である。R波起点のPWTTは純粋な血管の区間ではなく、前駆出期(PEP)を"
        "含む。PEPは前負荷・後負荷・収縮性で変動する。最も直接的には、メーカー自身の"
        "研究グループがPEPがPWTT変化の約半分を占めたと報告している(Sugo 2012)。\n"
        "本研究はこの前提そのものを問う。"))
    box(s, 0.6, 1.95, 12.1, 0.85, [
        ("PWTT ＝ 前駆出期（心臓側） ＋ 動脈通過時間（血管側）", 28, INK, True)],
        fill="F2F2F2", line=GREY)
    box(s, 0.6, 3.05, 5.9, 1.5, [
        ("前駆出期 PEP", 26, WHITE, True),
        ("前負荷・後負荷・収縮性", 24, WHITE, False)], fill=VERM, space_after=4)
    box(s, 6.8, 3.05, 5.9, 1.5, [
        ("動脈通過時間", 26, WHITE, True),
        ("血管指標が測るのはここだけ", 24, WHITE, False)], fill=BLUE, space_after=4)
    box(s, 0.6, 4.9, 12.1, 1.3, [
        ("PWTTの変動を心臓側が支配するなら", 26, INK, False),
        ("指標をどれだけ正確に測っても補正は成功しえない", 26, RED, True)],
        fill=None, line=RED, space_after=2)
    source(s, "Ochiai 1999; Payne 2006; Sugo 2012（PEPがPWTT変化の約半分）")

    # ---------------------------------------------------------- 4 方法
    t = "対象と解析デザイン"
    s = new_slide(prs, t, "方法", notes=(
        "VitalDB 6,388例から、脈波・心電図・動脈圧（いずれも500Hz）と連続COトラックを"
        "すべて持つ874例を抽出。有効ウィンドウ12以上（較正1＋評価11）を満たした862例"
        "（98.6%）を解析した。161,737個の60秒ウィンドウ。\n"
        "統計解析計画は、いかなる波形も見る前に確定・凍結した。指標定義・品質閾値・"
        "モデル仕様・解釈規準のすべてを含む。開発と検証には真値既知の合成信号のみを"
        "用いた。Zenodo doi:10.5281/zenodo.22167118 に第三者タイムスタンプ。"))
    box(s, 0.6, 1.95, 3.6, 1.05, [("VitalDB", 26, WHITE, True), ("6,388例", 26, WHITE, True)],
        fill=GREY, space_after=0)
    arrow(s, 4.4, 2.25, 0.7, 0.45)
    box(s, 5.3, 1.95, 3.6, 1.05, [("4波形すべて", 26, WHITE, True), ("874例", 26, WHITE, True)],
        fill=BLUE, space_after=0)
    arrow(s, 9.1, 2.25, 0.7, 0.45)
    box(s, 10.0, 1.95, 2.7, 1.05, [("解析", 26, WHITE, True), ("862例", 26, WHITE, True)],
        fill=TEAL, space_after=0)
    textbox(s, 0.6, 3.18, 12.1, 1.0, [
        ("4波形すべて＝脈波・心電図・動脈圧・連続CO", 24, INK, False),
        ("採用率 98.6％・161,737ウィンドウ", 24, INK, False)], space_after=0)
    box(s, 0.6, 4.25, 12.1, 2.2, [
        ("解析計画は波形を見る前に凍結した", 28, WHITE, True),
        ("指標定義・品質閾値・モデル・解釈規準をすべて事前確定", 24, WHITE, False),
        ("開発と検証には真値既知の合成信号のみを使用", 24, WHITE, False),
        ("Zenodo に第三者タイムスタンプを取得", 24, WHITE, False),
    ], fill=BLUE, space_after=5)
    source(s, "doi:10.5281/zenodo.22167118（タグ sap-v0.3・凍結日 2026年8月28日）")

    # ---------------------------------------------------------- 5 方法
    t = "参照COの限界"
    s = new_slide(prs, t, "方法", notes=(
        "解析862例のうち846例(98.1%)の参照COは動脈圧波形由来(FloTrac系)。これらは"
        "脈圧の標準偏差から1回拍出量を推定するため出力が血圧と連動し、信頼性が体血管"
        "抵抗に依存する。検証したい仮説も血管状態に関わるため、この参照との一致度が"
        "改善しても真のCO精度向上か共通依存かを区別できない。\n"
        "合成データでの事前検証では、精度比較は真の効果と見かけの改善を3回中3回とも"
        "区別できなかった。区別できたのは前提検証のみ(r² 0.658 対 0.001)。\n"
        "この制約こそが主解析を参照非依存に設計した理由である。"))
    box(s, 0.6, 1.95, 12.1, 1.15, [
        ("参照COの98％が動脈圧波形由来（FloTrac系）", 28, WHITE, True)], fill=VERM)
    box(s, 0.6, 3.35, 5.9, 1.5, [
        ("一致度が改善しても", 24, INK, False),
        ("真の精度向上か", 24, INK, True),
        ("共通依存かを区別できない", 24, INK, True)],
        fill="F2F2F2", line=VERM, space_after=2)
    box(s, 6.8, 3.35, 5.9, 1.5, [
        ("合成データでの事前検証", 24, INK, False),
        ("精度比較は判別力なし", 24, VERM, True),
        ("前提検証だけが判別できた", 24, BLUE, True)],
        fill="F2F2F2", line=BLUE, space_after=2)
    box(s, 0.6, 5.2, 12.1, 1.1, [
        ("→ 主解析は参照COを一切使わない設計にした", 28, WHITE, True)], fill=BLUE)
    source(s, "事前検証: 真の効果 r²=0.658 対 見かけの改善 r²=0.001（SAP §7.6）")

    # ---------------------------------------------------------- 6 方法
    t = "指標の定義"
    s = new_slide(prs, t, "方法", notes=(
        "アンサンブル平均した拍を歪みガウス2成分（前進波・反射波）に分解し、8初期値から"
        "非線形最小二乗で当てはめた。収束検算3つ（境界張り付き・振幅ゼロ・競合解）を"
        "すべて通過した当てはめのみ採用。\n"
        "SI・RIは文献間で定義が統一されていないため、実データを見る前に真値既知の合成"
        "脈波で候補を比較して凍結した。ΔTは候補5定義中最も頑健（誤差≤1.8ms）、RIは"
        "候補3定義中最も頑健（誤差≤0.9%）。RIは振幅パラメータの比ではなく成分ピーク"
        "高さの比である点に注意。"))
    box(s, 0.6, 1.95, 12.1, 0.8, [
        ("脈波を前進波と反射波の2成分に分解する", 26, INK, True)], fill="F2F2F2", line=GREY)
    box(s, 0.6, 3.0, 5.9, 1.75, [
        ("ΔT（スティフネス指標）", 26, WHITE, True),
        ("2成分のピーク時刻の差", 24, WHITE, False),
        ("＝タイミングの指標", 24, WHITE, True)], fill=BLUE, space_after=3)
    box(s, 6.8, 3.0, 5.9, 1.75, [
        ("RI（反射係数）", 26, WHITE, True),
        ("2成分のピーク高さの比", 24, WHITE, False),
        ("＝振幅の指標", 24, WHITE, True)], fill=VERM, space_after=3)
    box(s, 0.6, 5.1, 12.1, 1.15, [
        ("定義は実データを見る前に合成脈波で選定し凍結した", 26, INK, True)],
        fill=None, line=TEAL)
    source(s, "候補5定義からΔTを選定（誤差≤1.8ms）、候補3定義からRIを選定（誤差≤0.9%）")

    # ---------------------------------------------------------- 7 方法
    t = "主解析＝前提検証"
    s = new_slide(prs, t, "方法", notes=(
        "各症例で、初回ウィンドウからのPWTTの変化量を、同時点のΔTとRIの相対変化に"
        "切片つきで回帰した。参照COを一切使わない。\n"
        "前提の弱さと測定ノイズを切り分けるため、連続ウィンドウ間の1次自己相関を事前"
        "指定した。測定ノイズ支配なら自己相関は0近傍、生理を追う系列なら高い。\n"
        "さらに陽性対照として、加齢と動脈スティフネスの既知の関係（加齢でΔT短縮）を"
        "指標が再現するかを確認した。これは解析計画の凍結後に追加した探索的解析で、"
        "参照COも症例内変化も使わないため主要エンドポイントと独立である。"))
    box(s, 0.6, 1.95, 12.1, 1.0, [
        ("ΔPWTT ＝ b₀ ＋ b₁・ΔΔT％ ＋ b₂・ΔRI％", 28, WHITE, True)], fill=BLUE)
    textbox(s, 0.6, 3.15, 12.1, 0.5, [
        ("症例内の変化量だけを使う。参照COは登場しない", 24, INK, True)])
    box(s, 0.6, 3.9, 3.85, 2.3, [
        ("① 説明割合", 26, WHITE, True),
        ("血管指標で説明される", 24, WHITE, False),
        ("PWTT変動の割合", 24, WHITE, False)], fill=TEAL, space_after=3)
    box(s, 4.72, 3.9, 3.85, 2.3, [
        ("② 自己相関", 26, WHITE, True),
        ("測定ノイズによる", 24, WHITE, False),
        ("希釈でないかを確認", 24, WHITE, False)], fill=TEAL, space_after=3)
    box(s, 8.85, 3.9, 3.85, 2.3, [
        ("③ 陽性対照", 26, WHITE, True),
        ("加齢とΔTの既知の", 24, WHITE, False),
        ("関係を再現するか", 24, WHITE, False)], fill=TEAL, space_after=3)
    source(s, "③は解析計画の凍結後に追加した探索的解析（主解析の定義は不変）")

    # ---------------------------------------------------------- 8 結果
    t = "対象862例の特性"
    s = new_slide(prs, t, "結果", notes=(
        "年齢61歳(51-70)、男性60%、BMI 23.0。ASA 1-2が80%。一般外科71%・胸部外科25%、"
        "全例全身麻酔。主な術式は肝17%・移植15%・胆膵14%・大切除13%。緊急手術9%。"
        "記録長は中央値257分、有効ウィンドウ177個/例。\n"
        "血行動態が大きく動く長時間手術が中心であり、PWTTの変動を観察するには適した"
        "集団である。"))
    rows = [
        ("年齢", "61 歳（51–70）"),
        ("男性", "517 例（60％）"),
        ("ASA 1–2", "689 例（80％）"),
        ("高血圧の既往", "328 例（38％）"),
        ("記録長", "257 分（182–338）"),
        ("有効ウィンドウ／例", "177 個（104–255）"),
    ]
    for i, (k, v) in enumerate(rows):
        yy = 1.95 + i * 0.72
        box(s, 0.6, yy, 3.9, 0.68, [(k, 24, INK, False)], fill="F2F2F2", line=None,
            align=PP_ALIGN.LEFT)
        box(s, 4.6, yy, 3.9, 0.68, [(v, 24, INK, True)], fill=None, line=None,
            align=PP_ALIGN.LEFT)
    box(s, 8.7, 1.95, 4.0, 2.06, [
        ("一般外科 71％", 24, WHITE, False),
        ("胸部外科 25％", 24, WHITE, False),
        ("全例 全身麻酔", 24, WHITE, True)], fill=BLUE, space_after=3)
    box(s, 8.7, 4.19, 4.0, 2.06, [
        ("肝 17％／移植 15％", 24, WHITE, False),
        ("胆膵 14％", 24, WHITE, False),
        ("血行動態が動く手術", 24, WHITE, True)], fill=TEAL, space_after=3)
    source(s, "値は中央値（四分位範囲）または n（％）")

    # ---------------------------------------------------------- 9 結果（核心）
    t = "主結果：説明されない"
    s = new_slide(prs, t, "結果", notes=(
        "862例・161,737ウィンドウにおいて、血管指標はPWTTの症例内変動をまったく説明"
        "しなかった。プールr² = 0.000。係数はΔSI%あたり-0.027、ΔRI%あたり-0.003。\n"
        "症例内のΔSI%係数は78%の症例で血管仮説の予測どおりの符号を持ち、862例では"
        "偶然では説明できない。しかしその大きさは無視できる。症例内r²の中央値は0.101。"))
    box(s, 0.6, 1.95, 6.0, 2.5, [
        ("血管指標で説明された", 26, WHITE, False),
        ("PWTT変動の割合", 26, WHITE, False),
        ("r² ＝ 0.000", 54, WHITE, True)], fill=RED, space_after=2)
    box(s, 7.0, 1.95, 5.7, 2.5, [
        ("862 例", 26, INK, True),
        ("161,737 ウィンドウ", 26, INK, True),
        ("係数 ΔSI％ −0.027", 24, INK, False),
        ("係数 ΔRI％ −0.003", 24, INK, False)], fill="F2F2F2", line=GREY, space_after=3)
    box(s, 0.6, 4.7, 12.1, 1.55, [
        ("ただし係数の符号は78％の症例で仮説どおりの向き", 26, INK, True),
        ("＝関係は「無い」のではなく「小さすぎる」", 26, BLUE, True)],
        fill=None, line=BLUE, space_after=4)
    source(s, "症例内 r² 中央値 0.101。符号の一致は862例では偶然では説明できない")

    # ---------------------------------------------------------- 10 結果
    t = "測定は健全である"
    s = new_slide(prs, t, "結果", notes=(
        "陰性が測定失敗でないことを3つの独立した根拠で示す。\n"
        "① 陽性対照（成人849例）: ΔTは加齢とともに短縮 ρ=-0.197 [95%CI -0.261, -0.131] "
        "p<0.0001。高血圧既往例で259 対 267 ms と短い。陰性対照(症例IDとΔT)は分散に"
        "して1桁小さい。\n"
        "② 自己相関: PWTT +0.75、ΔT +0.50。雑音ではなく再現性のある生理を追う系列。\n"
        "③ 同定可能性: 患者データを見る前に真値既知の合成脈波で確立済み。"))
    box(s, 0.6, 1.95, 3.85, 2.6, [
        ("① 陽性対照", 26, WHITE, True),
        ("加齢でΔTが短縮", 24, WHITE, False),
        ("ρ ＝ −0.197", 28, WHITE, True),
        ("p ＜ 0.0001", 24, WHITE, False)], fill=BLUE, space_after=3)
    box(s, 4.72, 1.95, 3.85, 2.6, [
        ("② 自己相関", 26, WHITE, True),
        ("隣接ウィンドウ間", 24, WHITE, False),
        ("PWTT ＋0.75", 28, WHITE, True),
        ("ΔT ＋0.50", 24, WHITE, False)], fill=BLUE, space_after=3)
    box(s, 8.85, 1.95, 3.85, 2.6, [
        ("③ 同定可能性", 26, WHITE, True),
        ("合成脈波で事前に", 24, WHITE, False),
        ("確立済み", 28, WHITE, True),
        ("（真値既知）", 24, WHITE, False)], fill=BLUE, space_after=3)
    box(s, 0.6, 4.85, 12.1, 1.4, [
        ("ΔTについて、関係は覆い隠されたのではなく存在しない", 28, WHITE, True)],
        fill=TEAL)
    source(s, "陽性対照は成人849例。高血圧既往例でΔT 259 対 267 ms と期待どおり短い")

    # ---------------------------------------------------------- 11 結果
    t = "RIは判定できない"
    s = new_slide(prs, t, "結果", notes=(
        "ΔTが通過した同じ陽性対照を、RIは通過しなかった。成人849例で年齢とまったく"
        "関連しない ρ=+0.041, p=0.23。症例内でも有界な比なのに変動係数0.70とΔTの3倍"
        "動きながら、血圧との結びつきはΔTより弱い。生理ではなく雑音の振る舞い。\n"
        "合成波形で機序を絞り込んだ結果、一律の利得正規化ではRIは壊れない（同一拍内の"
        "比なので拍全体のスケーリングに不変）。壊すのは拍の内部で時定数0.25秒程度で"
        "変化する利得のみ（RI +61%、ΔTは-9msのみ）。\n"
        "この判断は陽性対照のみに基づき、前提検証の結果を見る前に行った。"))
    box(s, 0.6, 1.95, 5.9, 1.9, [
        ("ΔT（タイミング）", 26, WHITE, True),
        ("陽性対照を通過", 24, WHITE, False),
        ("ρ ＝ −0.197（p＜0.0001）", 24, WHITE, True)], fill=BLUE, space_after=3)
    box(s, 6.8, 1.95, 5.9, 1.9, [
        ("RI（振幅）", 26, WHITE, True),
        ("陽性対照を通過せず", 24, WHITE, False),
        ("ρ ＝ ＋0.041（p ＝ 0.23）", 24, WHITE, True)], fill=VERM, space_after=3)
    box(s, 0.6, 4.15, 12.1, 1.1, [
        ("RIの陰性は「関係が無い」ではなく「判定できない」", 28, WHITE, True)], fill=RED)
    textbox(s, 0.6, 5.45, 12.1, 1.0, [
        ("拍内で時定数0.25秒の利得変動だけがRIを選択的に壊す", 24, INK, True),
        ("（合成波形で機序を特定。一律の利得正規化では壊れない）", 24, INK, False)],
        space_after=2)
    source(s, "Couceiro 2015（振幅比はすべての血管参照に対して失敗）の独立した再現")

    # ---------------------------------------------------------- 12 結果
    t = "精度も改善しない"
    s = new_slide(prs, t, "結果", notes=(
        "血管補正は参照との一致度を改善しなかった。誤差率の中央値は対照26.9%、補正"
        "27.2%、差は+0.2ポイント[95%CI +0.1, +0.4]。すなわち補正はわずかだが統計的に"
        "有意な悪化をもたらしており、情報ではなくノイズを加えたことと整合する。\n"
        "Bland-Altmanバイアス -0.07 L/min（一致限界 -3.28〜+3.14）、4象限一致率0.56。\n"
        "誤差率は調整の組み合わせによらず平坦（対照26.9／血圧27.0／血管指標27.2／"
        "両方27.1）。事前指定の無益性基準に該当した。"))
    box(s, 0.6, 1.95, 5.9, 1.6, [
        ("対照（PWTT型）", 26, WHITE, False),
        ("誤差率 26.9％", 32, WHITE, True)], fill=GREY, space_after=3)
    box(s, 6.8, 1.95, 5.9, 1.6, [
        ("提案（K を血管補正）", 26, WHITE, False),
        ("誤差率 27.2％", 32, WHITE, True)], fill=VERM, space_after=3)
    box(s, 0.6, 3.85, 12.1, 1.15, [
        ("差 ＋0.2 ポイント［95％CI ＋0.1, ＋0.4］", 28, INK, True)],
        fill="F2F2F2", line=RED)
    textbox(s, 0.6, 5.2, 12.1, 1.1, [
        ("改善しないどころか、わずかだが有意に悪化した", 26, RED, True),
        ("＝補正が加えたのは情報ではなくノイズ", 26, INK, True)], space_after=2)
    source(s, "事前指定の無益性基準に該当。Bland-Altman バイアス −0.07 L/min")

    # ---------------------------------------------------------- 13 結果
    t = "感度解析：結論不変"
    s = new_slide(prs, t, "結果", notes=(
        "パイロット15例を除外しても全結論が不変（847例: r²=0.005、符号一致78%、"
        "差+0.2[+0.1,+0.3]）。\n"
        "ウィンドウを5分・20分に集約すると両群とも誤差率は低下したが（26.9→23.9→"
        "21.5%）、どの集約レベルでも補正は精度を改善しなかった。したがって精度の陰性は"
        "60秒ウィンドウの産物ではない。\n"
        "前提回帰に心拍数を加えると説明割合は0.000から0.077に上昇したが、血管指標の"
        "係数は実質不変。血管側の陰性は心拍数交絡の産物ではない。"))
    items = [
        ("パイロット15例を除外", "結論すべて不変（847例）"),
        ("ウィンドウ 5分 集約", "改善なし（−0.1）"),
        ("ウィンドウ 20分 集約", "改善なし（−0.2）"),
        ("心拍数を回帰に追加", "血管係数は不変"),
    ]
    for i, (k, v) in enumerate(items):
        yy = 1.95 + i * 1.0
        box(s, 0.6, yy, 5.9, 0.85, [(k, 24, INK, True)], fill="F2F2F2", line=None,
            align=PP_ALIGN.LEFT)
        box(s, 6.8, yy, 5.9, 0.85, [(v, 24, BLUE, True)], fill=None, line=BLUE,
            align=PP_ALIGN.LEFT)
    box(s, 0.6, 6.0, 12.1, 0.75, [
        ("結論は指標定義・前処理・時間スケールに依存しない", 26, INK, True)],
        fill=None, line=TEAL)

    # ---------------------------------------------------------- 14 考察（核心）
    t = "方向は正しいが小さい"
    s = new_slide(prs, t, "考察", notes=(
        "これは有意水準に達しなかったのではなく、精密に推定された無効果である。"
        "862例・161,737ウィンドウで血管説明割合は0.000。\n"
        "スティフネス指標が10%変化してもPWTTは0.27%しか変化しない。一方、症例内の"
        "PWTTは数%の幅で動いている。つまり血管成分は方向としては検出できるが、補正に"
        "用いるには約2桁小さい。\n"
        "「関係が無い」ではなく「2桁小さい」と書けることが、この研究の情報量である。"))
    box(s, 0.6, 1.95, 5.9, 1.9, [
        ("スティフネス指標が", 24, WHITE, False),
        ("10％ 変化しても", 26, WHITE, True),
        ("PWTT は 0.27％ しか動かない", 24, WHITE, True)], fill=BLUE, space_after=3)
    box(s, 6.8, 1.95, 5.9, 1.9, [
        ("一方 症例内の", 24, WHITE, False),
        ("PWTT は数％ の幅で", 26, WHITE, True),
        ("実際に動いている", 24, WHITE, True)], fill=VERM, space_after=3)
    box(s, 0.6, 4.15, 12.1, 1.25, [
        ("血管成分は「無い」のではなく 約2桁小さい", 30, WHITE, True)], fill=RED)
    textbox(s, 0.6, 5.6, 12.1, 0.9, [
        ("有意差が出なかったのではなく、無効果を精密に推定した", 24, INK, True),
        ("862例・161,737ウィンドウでの推定である", 24, INK, False)], space_after=2)
    source(s, "Altman & Bland 1995（absence of evidence は evidence of absence ではない）")

    # ---------------------------------------------------------- 15 考察
    t = "PWTTは心拍数を追う"
    s = new_slide(prs, t, "考察", notes=(
        "探索的解析として前提回帰に心拍数を加えると、説明割合は0.000から0.077に上昇"
        "した。血管指標が全く説明しなかった症例内PWTT変動の7.7%を心拍数だけで説明する。"
        "しかも血管指標の係数はその投入によって変わらない。\n"
        "向きも生理に合う（心拍数上昇でPWTT短縮＝交感神経賦活によるPEP短縮と整合）。\n"
        "この研究は「血管指標を追わない」という陰性だけでなく、「では何を追うのか」に"
        "対する陽性の答えも持ったことになる。"))
    box(s, 0.6, 1.95, 5.9, 2.2, [
        ("血管指標のみ", 26, WHITE, False),
        ("r² ＝ 0.000", 40, WHITE, True)], fill=VERM, space_after=4)
    box(s, 6.8, 1.95, 5.9, 2.2, [
        ("心拍数を加えると", 26, WHITE, False),
        ("r² ＝ 0.077", 40, WHITE, True)], fill=BLUE, space_after=4)
    box(s, 0.6, 4.45, 12.1, 1.05, [
        ("血管指標の係数は心拍数を入れても変わらない", 26, INK, True)],
        fill="F2F2F2", line=GREY)
    textbox(s, 0.6, 5.7, 12.1, 0.9, [
        ("PWTTが追っているのは血管ではなく心拍数・自律神経状態", 26, INK, True)],
        space_after=2)
    source(s, "心拍数上昇でPWTT短縮。交感神経賦活による前駆出期の短縮と整合する")

    # ---------------------------------------------------------- 16 考察
    t = "前駆出期が支配する"
    s = new_slide(prs, t, "考察", notes=(
        "本結果は異常ではなく、機序に関する文献が予測するとおりである。\n"
        "・Ochiai 1999: 犬でm-PWTT = PEP + PWTT と明示的に分解\n"
        "・Payne 2006: 心電図起点の通過時間は純粋な血管機能の指標として不適\n"
        "・Djupedal 2022: 同一被験者内でPEPと血管通過時間が逆方向に動きうる\n"
        "・Pilz 2023: ストレス下でPEPが最大約50ms変動\n"
        "・Sugo 2012（メーカー自身）: PEPがPWTT変化の約半分を占めた\n"
        "本研究はその機序の、ヒト・術中・参照非依存での確認である。"))
    box(s, 0.6, 1.95, 12.1, 1.35, [
        ("メーカー自身の研究グループの報告", 24, WHITE, False),
        ("「前駆出期がPWTT変化の約半分を占めた」", 28, WHITE, True)], fill=BLUE, space_after=3)
    lit = [
        ("Payne 2006", "心電図起点は血管指標として不適"),
        ("Djupedal 2022", "PEPと通過時間が逆向きに動く"),
        ("Pilz 2023", "ストレスでPEPが約50ms変動"),
    ]
    for i, (k, v) in enumerate(lit):
        yy = 3.55 + i * 0.85
        box(s, 0.6, yy, 3.5, 0.72, [(k, 24, INK, True)], fill="F2F2F2", line=None,
            align=PP_ALIGN.LEFT)
        box(s, 4.35, yy, 8.35, 0.72, [(v, 24, INK, False)], fill=None, line=None,
            align=PP_ALIGN.LEFT)
    box(s, 0.6, 6.1, 12.1, 0.7, [
        ("本研究はこの機序のヒト・術中・参照非依存での確認", 24, INK, True)],
        fill=None, line=TEAL)

    # ---------------------------------------------------------- 17 考察
    t = "では何が有望か"
    s = new_slide(prs, t, "考察", notes=(
        "補正すべき量が補正変数と共変しないなら、指標をどれほど正確に測っても補正は"
        "機能しえない。単一部位の光電容積脈波による較正補正には上限がある。\n"
        "より有望なのは、①心臓側の項を直接測って除く方向（心音図・インピーダンス"
        "心拍出量計・バイオリアクタンスによる大動脈弁開放の計時）、②ドリフトをモデル化"
        "するのではなく再較正の頻度を上げる方向、である。"))
    box(s, 0.6, 1.95, 12.1, 1.05, [
        ("補正すべき量が補正変数と共変しないなら補正は効かない", 26, WHITE, True)],
        fill=RED)
    box(s, 0.6, 3.3, 5.9, 2.3, [
        ("① 心臓側を直接測る", 26, WHITE, True),
        ("心音図", 24, WHITE, False),
        ("インピーダンス", 24, WHITE, False),
        ("大動脈弁開放の計時", 24, WHITE, False)], fill=BLUE, space_after=3)
    box(s, 6.8, 3.3, 5.9, 2.3, [
        ("② 再較正を増やす", 26, WHITE, True),
        ("ドリフトをモデル化", 24, WHITE, False),
        ("するのではなく", 24, WHITE, False),
        ("較正の頻度で対処", 24, WHITE, False)], fill=TEAL, space_after=3)
    textbox(s, 0.6, 5.85, 12.1, 0.65, [
        ("血管成分をモデル化するより心臓側成分を分離する方が有望", 24, INK, True)])

    # ---------------------------------------------------------- 18 考察
    t = "副産物：装置遅延"
    s = new_slide(prs, t, "考察", notes=(
        "本データベースの脈波チャネルは心電図に対する固定処理遅延を持つ。本コホートでは"
        "症例レベル中央値660ms、四分位範囲644-676msという狭さで、生理ではなく装置の"
        "処理定数であることと整合する。\n"
        "遅延が高心拍数下で心周期を超えるため、これを考慮せずに計算した通過時間は"
        "単にずれるのではなく折り返す（エイリアシング）。\n"
        "この現象自体は他所で大規模に記録済みで新しくない(Ruffolo 2025)。我々が加える"
        "のは、波形同期は妥当と仮定され、その仮定が派生ベンチマークに伝播している"
        "(PulseDB 2022)この特定の資源における定量と補正である。"))
    box(s, 0.6, 1.95, 6.0, 1.9, [
        ("脈波チャネルの装置遅延", 26, WHITE, False),
        ("660 ms", 44, WHITE, True),
        ("（四分位 644–676）", 24, WHITE, False)], fill=BLUE, space_after=2)
    box(s, 7.0, 1.95, 5.7, 1.9, [
        ("極めて狭い分布", 26, INK, True),
        ("＝生理ではなく", 24, INK, False),
        ("装置の処理定数", 24, INK, True)], fill="F2F2F2", line=BLUE, space_after=3)
    box(s, 0.6, 4.15, 12.1, 1.1, [
        ("高心拍では心周期を超え、通過時間が折り返す", 28, WHITE, True)], fill=RED)
    textbox(s, 0.6, 5.45, 12.1, 1.0, [
        ("公開データベースを使う他の研究者への実用的な警告", 24, INK, True),
        ("既存ベンチマークは同期を問題なしと仮定している", 24, INK, False)], space_after=2)
    source(s, "現象自体は既報（Ruffolo 2025）。本研究の貢献はVitalDBでの定量と補正")

    # ---------------------------------------------------------- 19 考察
    t = "限界"
    s = new_slide(prs, t, "考察", notes=(
        "① 参照COは主に動脈圧波形由来（846/862例）で独立標準ではない。副次解析にのみ"
        "影響し、主解析には影響しない。独立参照は熱希釈5例・食道ドプラ11例。\n"
        "② 単一データベース・単一施設の後ろ向き解析。統制された血管運動負荷を伴わない。\n"
        "③ 実機の係数は非公開のため、対照推定器は公表式の再現であり実機ではない。\n"
        "④ 2成分分解の成分帰属は物理的に異なる波に対応する保証がない。\n"
        "⑤ SI・RIは主に安静時指標として開発されており術中への外挿は仮定である。\n"
        "⑥ 脈波は処理済みモニタ出力。タイミングは生き延びるが振幅は生き延びない。"))
    lims = [
        "参照COの98％が動脈圧由来（主解析には影響しない）",
        "単一データベース・単一施設・後ろ向き",
        "実機の係数は非公開のため公表式の再現である",
        "成分の物理的帰属は保証されない",
        "安静時指標を術中の急性変化へ外挿している",
    ]
    for i, txt in enumerate(lims):
        yy = 1.92 + i * 0.88
        box(s, 0.6, yy, 12.1, 0.76, [(f"{i + 1}．{txt}", 24, INK, False)],
            fill="F2F2F2", line=None, align=PP_ALIGN.LEFT)
    box(s, 0.6, 6.42, 12.1, 0.7, [
        ("主解析は参照COを使わないため①の影響を受けない", 24, INK, True)],
        fill=None, line=TEAL)

    # ---------------------------------------------------------- 20 結論
    t = "結論"
    s = new_slide(prs, t, "考察", notes=(
        "大規模周術期波形データベースにおいて、拍ごとの脈波伝播時間の変動は、既知の"
        "血管シグナルを検出できるだけ十分に測定されているにもかかわらず、光電容積脈波"
        "由来の動脈スティフネス指標ではほとんど説明されなかった。対応する振幅由来の"
        "指標はこの信号源では妥当性を確立できず未検証のままである。\n"
        "したがって、通過時間型心拍出量推定の較正定数をこれらの指標で動的に補正する"
        "方策には、機能する余地がほとんどない。通過時間の血管成分をモデル化するのでは"
        "なく、その心臓側の成分を分離する方法こそが、より有望な方向である。"))
    box(s, 0.6, 2.0, 12.1, 1.3, [
        ("PWTTの術中変動は血管指標では説明されない", 30, WHITE, True)], fill=BLUE)
    box(s, 0.6, 3.5, 12.1, 1.3, [
        ("測定は健全。血管成分は方向は正しいが約2桁小さい", 28, INK, True)],
        fill="F2F2F2", line=BLUE)
    box(s, 0.6, 5.0, 12.1, 1.3, [
        ("較正定数の血管補正には機能する余地がほとんどない", 28, WHITE, True)], fill=RED)
    source(s, "VitalDB 862例・161,737ウィンドウ・事前登録解析（doi:10.5281/zenodo.22167118）")

    prs.save(out_path)
    return len(prs.slides.__iter__.__self__._sldIdLst)


if __name__ == "__main__":
    out = sys.argv[1] if len(sys.argv) > 1 else "研究1_結果報告.pptx"
    n = build(out)
    print(f"生成: {out}")
