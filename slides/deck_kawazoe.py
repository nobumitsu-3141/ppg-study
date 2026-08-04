# -*- coding: utf-8 -*-
"""川副式スライド組版ヘルパ（Web 版・python-pptx 実装）.

slide-format スキルの 6 つの絶対ルールを機械的に満たすための最小ライブラリ。
  1. タイトル 44pt 太字・金 BF9000・黒縁取り 2.25pt・直下に金の全幅下線（8pt, y=1.52")
  2. 右上に章ナビ帯（チップ列）。現在章のみティール、他は灰
  3. 本文 22pt 以上（出典 16pt / チップ 10.5pt / ページ番号 12pt のみ例外）
  4. 文章はすべてネイティブのテキスト（画像に焼き込まない）
  5. 図解優先・図と文字を重ねない・詳述はノートへ
  6. 単語の途中で改行しない（1 段落 1 行に収める）
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

# ---------------------------------------------------------------- 配色（役割固定）
GOLD = "BF9000"      # 構造色（タイトル・下線）。対比の一色には使わない
TEAL = "00A8AA"      # 数学モデル・推定／現在章チップ
BLUE = "0072B2"      # 時間の軸（ΔT・SI・大動脈スティフネス）＋前進波
VERM = "D55E00"      # 高さの軸（P2/P1・RI・末梢トーヌス）＋反射波
RED = "C00000"       # 注意・限界（単独の意味色）
GRAY = "D9D9D9"      # 非現在章チップ
INK = "222222"       # 本文
SUB = "5A5A5A"       # 補足・出典
WAVE = "3A3A3A"      # 実測波形

FONT = "メイリオ"

# ---------------------------------------------------------------- 版面座標（16:9）
SW_IN, SH_IN = 13.333, 7.5
TITLE_X, TITLE_Y, TITLE_W, TITLE_H = 0.53, 0.35, 8.35, 0.95
RULE_Y = 1.52
CHIP_W, CHIP_H, CHIP_Y = 0.68, 0.64, 0.06
CHIP_RIGHT = 13.02
BODY_TOP = 1.80
BODY_BOTTOM = 6.62
SRC_Y = 6.86
A = "xmlns:a=\"http://schemas.openxmlformats.org/drawingml/2006/main\""


# ---------------------------------------------------------------- 幅推定（lint と同一式）
JP_W = 1.12


def _units(text):
    u = 0.0
    for ch in text:
        o = ord(ch)
        if ch == " ":
            u += 0.30
        elif o <= 0x24F:
            u += 0.56
        elif 0x2080 <= o <= 0x208E:
            u += 0.56
        else:
            u += 1.0
    return u


def text_w_in(text, pt):
    """lint と同じ推定式でテキスト幅（インチ）を返す。"""
    return _units(text) * (pt / 72.0) * JP_W


WARNINGS = []


def _check_line(text, pt, box_w_in, where):
    if pt >= 20 and text.strip():
        w = text_w_in(text, pt)
        if w > box_w_in + 0.059:      # lint の +0.15cm 相当
            WARNINGS.append(f"[折返し] {where}: {pt}pt 幅{w:.2f}\" > 枠{box_w_in:.2f}\" 「{text[:22]}」")


def line_h_in(pt, ls=1.0):
    """lint の est_overflow と同じ行高（インチ）。JP_H=1.18 を含む。"""
    return (pt / 72.0) * ls * 1.18


def _check_overflow(items, box_w, box_h, space_after, ls, where):
    """lint の est_overflow と同じ推定で、枠からのはみ出しを事前警告する。"""
    total = 0.0
    for t, sz in items:
        if not t.strip():
            total += 0.2 / 2.54
            continue
        lines = max(1, int(-(-text_w_in(t, sz) * 1000 // (box_w * 1000))))
        total += lines * line_h_in(sz, ls) + space_after / 72.0
    if total - box_h > 0.3 / 2.54:
        WARNINGS.append(f"[あふれ] {where}: 推定{total:.2f}\" > 枠{box_h:.2f}\"")


# ---------------------------------------------------------------- 低レベル部品
def _srgb(el_tag, color):
    return parse_xml(f'<a:{el_tag} {A}><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:{el_tag}>')


def outline_run(run, color="000000", pt=2.25):
    """文字の縁取り（a:ln）。a:rPr の先頭に入れる必要がある。"""
    rPr = run._r.get_or_add_rPr()
    w = int(pt * 12700)
    ln = parse_xml(f'<a:ln {A} w="{w}"><a:solidFill><a:srgbClr val="{color}"/></a:solidFill></a:ln>')
    rPr.insert(0, ln)


SUBS = {"₀": "0", "₁": "1", "₂": "2", "₃": "3", "₄": "4",
        "₅": "5", "₆": "6", "₇": "7", "₈": "8", "₉": "9"}


def emit_runs(para, text, size, color, bold=False):
    """下付き文字（U+2080-2089）を、通常の数字＋baseline 下げの run に変換して流す。

    メイリオに下付き数字のグリフが無い環境でも確実に描画され、
    PowerPoint 上では通常どおり編集できる。
    """
    buf = ""
    for ch in text:
        if ch in SUBS:
            if buf:
                set_run(para.add_run(), buf, size, color, bold)
                buf = ""
            r = para.add_run()
            set_run(r, SUBS[ch], size, color, bold)
            r.font._rPr.set("baseline", "-25000")
        else:
            buf += ch
    if buf:
        set_run(para.add_run(), buf, size, color, bold)


def set_run(run, text, size, color, bold=False, italic=False):
    run.text = text
    f = run.font
    f.name = FONT
    f.size = Pt(size)
    f.bold = bold
    f.italic = italic
    f.color.rgb = RGBColor.from_string(color)
    # 東アジアフォントも明示（Meiryo で描画させる）
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:ea", "a:cs"):
        e = parse_xml(f'<{tag} {A} typeface="{FONT}"/>')
        rPr.append(e)


def textbox(slide, x, y, w, h, paras, size=24, color=INK, bold=False,
            align=PP_ALIGN.LEFT, space_after=8, line_spacing=1.0, anchor=MSO_ANCHOR.TOP,
            where=""):
    """paras: 文字列 or (text, size, color, bold) タプルのリスト。1 段落 1 行前提。"""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(0.02)
    tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    first = True
    seen = []
    for p in paras:
        if isinstance(p, str):
            t, sz, col, bd = p, size, color, bold
        else:
            t = p[0]
            sz = p[1] if len(p) > 1 and p[1] else size
            col = p[2] if len(p) > 2 and p[2] else color
            bd = p[3] if len(p) > 3 else bold
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.alignment = align
        para.line_spacing = line_spacing
        para.space_after = Pt(space_after)
        para.space_before = Pt(0)
        seen.append((t, sz))
        if t:
            emit_runs(para, t, sz, col, bd)
            _check_line(t, sz, w - 0.04, where or f"({x:.1f},{y:.1f})")
    _check_overflow(seen, w, h, space_after, line_spacing,
                    where or f"txt({x:.1f},{y:.1f})")
    return box


def box_h_for(paras, size=24, space_after=8, line_spacing=1.0):
    """textbox の推定描画高さ（インチ）＋余裕。"""
    total = 0.0
    for p in paras:
        sz = size if isinstance(p, str) else (p[1] if len(p) > 1 and p[1] else size)
        total += line_h_in(sz, line_spacing) + space_after / 72.0
    return total + 0.16


def panel(slide, x, y, w, h, fill=None, line=None, line_pt=1.75, radius=0.08,
          paras=None, size=24, color=INK, bold=False, align=PP_ALIGN.CENTER,
          anchor=MSO_ANCHOR.MIDDLE, space_after=4, where=""):
    """色分けボックス。文字は図形自身のテキストフレームに入れる（重なり判定を回避）。"""
    sp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y),
                                Inches(w), Inches(h))
    try:
        sp.adjustments[0] = radius
    except Exception:
        pass
    if fill:
        sp.fill.solid()
        sp.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        sp.fill.background()
    if line:
        sp.line.color.rgb = RGBColor.from_string(line)
        sp.line.width = Pt(line_pt)
    else:
        sp.line.fill.background()
    sp.shadow.inherit = False
    tf = sp.text_frame
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.margin_left = tf.margin_right = Inches(0.06)
    tf.margin_top = tf.margin_bottom = Inches(0.02)
    tf.vertical_anchor = anchor
    if paras:
        first = True
        seen = []
        for p in paras:
            if isinstance(p, str):
                t, sz, col, bd = p, size, color, bold
            else:
                t = p[0]
                sz = p[1] if len(p) > 1 and p[1] else size
                col = p[2] if len(p) > 2 and p[2] else color
                bd = p[3] if len(p) > 3 else bold
            para = tf.paragraphs[0] if first else tf.add_paragraph()
            first = False
            para.alignment = align
            para.line_spacing = 1.0
            para.space_after = Pt(space_after)
            para.space_before = Pt(0)
            seen.append((t, sz))
            if t:
                emit_runs(para, t, sz, col, bd)
                _check_line(t, sz, w - 0.12, where or f"panel({x:.1f},{y:.1f})")
        _check_overflow(seen, w, h, space_after, 1.0,
                        where or f"panel({x:.1f},{y:.1f})")
    return sp


def line(slide, x1, y1, x2, y2, color=WAVE, pt=1.5, dash=None, arrow=False, head=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1),
                                   Inches(x2), Inches(y2))
    c.line.color.rgb = RGBColor.from_string(color)
    c.line.width = Pt(pt)
    if dash:
        c.line.dash_style = dash
    ln = c.line._get_or_add_ln()
    if head:
        ln.append(parse_xml(f'<a:headEnd {A} type="triangle" w="med" len="med"/>'))
    if arrow:
        ln.append(parse_xml(f'<a:tailEnd {A} type="triangle" w="med" len="med"/>'))
    c.shadow.inherit = False
    return c


def curve(slide, pts, color=WAVE, pt=2.5, dash=None):
    """pts: [(x_in, y_in), ...] の折れ線をフリーフォームで描く。"""
    ex = [Emu(int(px * 914400)) for px, _ in pts]
    ey = [Emu(int(py * 914400)) for _, py in pts]
    b = slide.shapes.build_freeform(ex[0], ey[0])
    b.add_line_segments(list(zip(ex[1:], ey[1:])), close=False)
    sp = b.convert_to_shape()
    sp.fill.background()
    sp.line.color.rgb = RGBColor.from_string(color)
    sp.line.width = Pt(pt)
    if dash:
        sp.line.dash_style = dash
    sp.shadow.inherit = False
    return sp


# ---------------------------------------------------------------- デッキ
class Deck:
    def __init__(self, chapters):
        self.prs = Presentation()
        self.prs.slide_width = Inches(SW_IN)
        self.prs.slide_height = Inches(SH_IN)
        self.chapters = chapters
        self.page = 0

    def _title(self, slide, text):
        ph = slide.shapes.title
        ph.left, ph.top = Inches(TITLE_X), Inches(TITLE_Y)
        ph.width, ph.height = Inches(TITLE_W), Inches(TITLE_H)
        tf = ph.text_frame
        tf.word_wrap = False
        tf.auto_size = MSO_AUTO_SIZE.NONE
        tf.margin_left = tf.margin_right = 0
        tf.margin_top = tf.margin_bottom = 0
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        p.line_spacing = 1.0
        r = p.add_run()
        set_run(r, text, 44, GOLD, bold=True)
        outline_run(r, "000000", 2.25)
        w = text_w_in(text, 44)
        if w > TITLE_W:
            WARNINGS.append(f"[タイトル幅] 「{text}」 推定{w:.2f}\" > {TITLE_W}\"")
        return ph

    def _rule(self, slide):
        line(slide, 0, RULE_Y, SW_IN, RULE_Y, GOLD, 8)

    def _chips(self, slide, active):
        n = len(self.chapters)
        x0 = CHIP_RIGHT - n * CHIP_W
        for i, name in enumerate(self.chapters):
            cur = (active == i)
            panel(slide, x0 + i * CHIP_W + 0.02, CHIP_Y, CHIP_W - 0.04, CHIP_H,
                  fill=(TEAL if cur else GRAY), line=None, radius=0.14,
                  paras=[(name, 10.5, "FFFFFF" if cur else "666666", cur)],
                  align=PP_ALIGN.CENTER, space_after=0, where="chip")

    def add(self, title, chapter=None, source=None, notes=None, chips=True):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._title(slide, title)
        self._rule(slide)
        if chips:
            self._chips(slide, chapter)
        self.page += 1
        textbox(slide, 12.35, SRC_Y, 0.65, 0.34, [str(self.page)], size=12, color=SUB,
                align=PP_ALIGN.RIGHT, space_after=0, where="pageno")
        if source:
            textbox(slide, 0.5, SRC_Y - 0.04, 11.4, 0.58, [source], size=16, color=SUB,
                    space_after=0, where="source")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return slide

    def cover(self, title, lines, notes=None):
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[5])
        self._title(slide, title)
        self._rule(slide)
        textbox(slide, 0.5, 2.1, 12.3, 2.4, lines, size=28, color=INK, space_after=16,
                where="cover")
        if notes:
            slide.notes_slide.notes_text_frame.text = notes
        return slide

    def save(self, path):
        self.prs.save(path)
