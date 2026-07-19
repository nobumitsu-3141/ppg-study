# -*- coding: utf-8 -*-
"""川副式スライド ビルダー（再利用可能ヘルパー）.
テンプレート _template_kawazoe.pptx の上に、金色タイトル＋下線・パンくず・ページ番号・
出典を自動付与しながらスライドを組み立てる。ルールは スライド作成ルール.md を参照。"""
import os
import sys
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ---- palette (deck spec) ----
GOLD   = RGBColor(0xBF,0x90,0x00)   # accent4 darker 25% : title / key / emphasis
TEAL   = RGBColor(0x00,0xA8,0xAA)   # breadcrumb active
TEALD  = RGBColor(0x00,0x80,0x7F)
INK    = RGBColor(0x26,0x26,0x26)
GRAY   = RGBColor(0x80,0x80,0x80)
LGRAY  = RGBColor(0xD9,0xD9,0xD9)
MGRAY  = RGBColor(0xA6,0xA6,0xA6)
BLUE   = RGBColor(0x1F,0x4E,0x79); BLUEL = RGBColor(0xDD,0xEB,0xF7)
ORANGE = RGBColor(0x8A,0x3B,0x00); ORANGEL=RGBColor(0xFB,0xE5,0xD6)
GOLDBG = RGBColor(0xF7,0xEC,0xCF)
RED    = RGBColor(0xC0,0x00,0x00)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
JP = "Meiryo"

# 8-chapter breadcrumb labels (EtCO2 deck)
CHAPTERS = ["歴史", "原理", "生理", "正常", "術中", "異常", "対応", "特殊"]

SW, SH = Cm(33.867), Cm(19.05)   # 16:9

# ---- title safety floors (スライド作成ルール.md) ----
TITLE_LEFT_MIN = Cm(1.6)    # left margin must be >= 1.6cm (template default 2.33cm already clears this)
TITLE_BOTTOM_MAX = Cm(3.60) # keep the title's rendered bottom above the gold underline (fixed at y=3.85cm)
BODY_MIN_PT = 22            # no body text below this floor anywhere in the deck (except exempt small sizes)

def new_deck(template):
    return Presentation(template)

def _set_font(run, size=None, bold=None, color=None, name=JP):
    f = run.font
    if size is not None: f.size = Pt(size)
    if bold is not None: f.bold = bold
    if color is not None: f.color.rgb = color
    if name:
        f.name = name
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin","a:ea","a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {}); rPr.append(el)
            el.set("typeface", name)

def _txbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, 0)
    return tb, tf

def add_para(tf, runs, align=PP_ALIGN.LEFT, space_after=6, level=0, line=None, first=False):
    """runs: list of (text, dict-of-font-kwargs)."""
    p = tf.paragraphs[0] if first and tf.paragraphs[0].text=="" and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.level = level
    if space_after is not None: p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    for txt, kw in runs:
        r = p.add_run(); r.text = txt; _set_font(r, **kw)
    return p

# ---------- title (gold, from layout) ----------
# The gold underline lives at a FIXED y in the layout, so the title MUST stay on
# one line. Auto-shrink the font (default 46pt) to fit the content width.
def _title_units(text):
    u = 0.0
    for ch in text:
        o = ord(ch)
        if ch == " ": u += 0.30
        elif o <= 0x24F: u += 0.56           # latin / digits
        elif 0x2080 <= o <= 0x208E: u += 0.56  # subscripts (₂ …)
        else: u += 1.0                        # CJK / full-width punct / arrows
    return u

TITLE_BC_RIGHT = Cm(18.8)  # breadcrumb-bearing slides: title box right edge must stay <=this,
                            # clear of the pennant strip (starts ~19.2cm) — see スライド作成ルール.md

def _set_title(slide, text, cap=54, maxw=20.5, breadcrumb=False):
    # Title font is FIXED at 44pt — never auto-shrunk. The primary fix for overflow is
    # shortening the title text (warn to stderr when it doesn't fit at 44pt within maxw).
    # As a defense-in-depth backstop, breadcrumb-bearing slides get their title box
    # width hard-capped at TITLE_BC_RIGHT with word_wrap=True, so any title that still
    # doesn't fit wraps to a 2nd line inside the box instead of spilling across the
    # breadcrumb pennants (2行化 is then caught by slide_lint and must be fixed by
    # shortening the text further — wrapping is a safety net, not the goal).
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            u = _title_units(text)
            size = 44.0
            w_at_size = u * (size/72*2.54)          # width in cm at fixed size
            if u > 0 and w_at_size > maxw:
                print(f"[title-overflow] {w_at_size:.1f}cm>{maxw} : {text}", file=sys.stderr)
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    _set_font(r, size=size, name=JP)   # keep layout gold color
            tf = ph.text_frame
            tf.word_wrap = bool(breadcrumb)    # non-breadcrumb slides: single line only, never wrap
            tf.auto_size = MSO_AUTO_SIZE.NONE  # never let the renderer auto-shrink the fixed 44pt
            tf.vertical_anchor = MSO_ANCHOR.BOTTOM
            # Read the (possibly layout-inherited) geometry BEFORE writing anything:
            # placeholders on a freshly-added slide have no explicit <a:xfrm> yet, so
            # left/top/width resolve via inheritance from the layout. Writing just one
            # dimension (e.g. height) forces python-pptx to create a bare <a:xfrm> whose
            # other numbers default to 0 (width=0 → invisible title) — so all four must
            # be captured and re-written together.
            left, top, width = ph.left, ph.top, ph.width
            if left < TITLE_LEFT_MIN:           # left-margin safety floor
                left = TITLE_LEFT_MIN
            if breadcrumb:
                width = max(Cm(4.0), TITLE_BC_RIGHT - left)
            # bottom-anchor within a box short enough that the rendered line always
            # stays above the gold underline (fixed at y=3.85cm), regardless of
            # template tweaks to the placeholder's own height.
            height = max(Cm(1.0), TITLE_BOTTOM_MAX - top)
            ph.left, ph.top, ph.width, ph.height = left, top, width, height
            return ph
    return None

def _strip_body(slide):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 1:
            ph._element.getparent().remove(ph._element)

# ---------- breadcrumb (numbered downward pennants, packed top-right) ----------
# Number and label BOTH sit inside each pennant, stacked (number on top, label
# below) — no separate badge. Strip is packed into the top-right corner.
def add_breadcrumb(slide, active):
    n = len(CHAPTERS)
    w = Cm(1.72); penH = Cm(1.55); gap = Cm(0.10)
    total = n*w + (n-1)*gap
    x0 = SW - total - Cm(0.16); y0 = Cm(0.14)
    for i, lab in enumerate(CHAPTERS):
        x = x0 + i*(w+gap)
        act = (i == active)
        col = WHITE if act else GRAY
        pen = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR, x, y0, w, penH)
        pen.fill.solid(); pen.fill.fore_color.rgb = TEAL if act else LGRAY
        pen.line.fill.background(); pen.shadow.inherit = False
        tf = pen.text_frame; tf.word_wrap = True
        tf.margin_left=0; tf.margin_right=0; tf.margin_top=Cm(0.12); tf.margin_bottom=0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; p1.line_spacing = 1.0
        r1 = p1.add_run(); r1.text = str(i+1); _set_font(r1, size=10.5, bold=True, color=col)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(1); p2.line_spacing = 1.0
        r2 = p2.add_run(); r2.text = lab; _set_font(r2, size=12, bold=True, color=col)

# ---------- citation (small gray, optional hyperlink) ----------
def add_citation(slide, text, url=None, x=Cm(2.0), y=Cm(18.35), w=Cm(24), align=PP_ALIGN.LEFT):
    tb, tf = _txbox(slide, x, y, w, Cm(0.6))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; _set_font(r, size=10.5, color=GRAY)
    if url:
        r.hyperlink.address = url          # clickable link to the paper
        r.font.underline = False           # keep it looking like plain gray text
        r.font.color.rgb = GRAY            # gray, not the blue link colour

# ---------- figure fit ----------
def add_figure(slide, path, x=Cm(2.0), y=Cm(5.0), w=Cm(29.8), h=Cm(12.7), valign="mid"):
    iw, ih = Image.open(path).size
    ar = iw/ih; box_ar = (w/h)
    if ar > box_ar:
        nw = w; nh = int(w/ar)
    else:
        nh = h; nw = int(h*ar)
    px = x + (w-nw)//2
    if valign=="top": py = y
    elif valign=="bottom": py = y + (h-nh)
    else: py = y + (h-nh)//2
    return slide.shapes.add_picture(path, px, py, nw, nh)

# ---------- rounded box + text ----------
def add_box(slide, x, y, w, h, fill, line=None, round=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp

# ---------- slide factories ----------
def content_slide(prs, title, active=None, cite=None, cite_url=None):
    s = prs.slides.add_slide(prs.slide_layouts[2])
    # no breadcrumb (active is None) → title may use more width; else keep clear of pennants
    bc = active is not None
    _set_title(s, title, maxw=(26.0 if not bc else 16.2), breadcrumb=bc); _strip_body(s)
    if active is not None: add_breadcrumb(s, active)
    if cite: add_citation(s, cite, url=cite_url)
    return s

def figure_slide(prs, title, fig, active=None, cite=None, cite_url=None, caption=None):
    s = content_slide(prs, title, active, cite, cite_url)
    yb = Cm(5.0); hb = Cm(12.6 if not cite else 12.0)
    if caption:
        tb, tf = _txbox(s, Cm(2.0), Cm(4.8), Cm(29.8), Cm(1.0), anchor=MSO_ANCHOR.TOP)
        add_para(tf, [(caption, dict(size=16, bold=True, color=INK))], align=PP_ALIGN.CENTER, first=True)
        yb = Cm(6.0); hb = Cm(11.4 if not cite else 10.9)
    add_figure(s, fig, y=yb, h=hb)
    return s

def title_slide(prs, main, sub, hero=None):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            add_para(ph.text_frame, [(main, dict(size=40, bold=True, color=INK))], align=PP_ALIGN.CENTER, first=True)
        elif ph.placeholder_format.idx == 1:
            ph.text = ""
            for i,(t,kw) in enumerate(sub):
                add_para(ph.text_frame, [(t,kw)], align=PP_ALIGN.CENTER, first=(i==0), space_after=4)
    if hero:
        add_figure(s, hero, x=Cm(3.5), y=Cm(14.6), w=Cm(26.8), h=Cm(3.6), valign="mid")
    return s

def bullets(prs, title, items, active=None, cite=None, cite_url=None, size=22, lead=None):
    """items: list of (text, dict-of-font-kwargs). kwargs may include sa(space_after pt), lv(level)."""
    s = content_slide(prs, title, active=active, cite=cite, cite_url=cite_url)
    y = Cm(5.6)
    if lead:
        tb, tf = _txbox(s, Cm(2.4), y, Cm(29.0), Cm(1.4))
        add_para(tf, [(lead, dict(size=18, bold=True, color=GOLD))], first=True)
        y = Cm(7.4)
    tb, tf = _txbox(s, Cm(2.6), y, Cm(28.6), Cm(12.4 if not cite else 11.6))
    for i, (txt, kw) in enumerate(items):
        kw = dict(kw); sa = kw.pop("sa", 14); lv = kw.pop("lv", 0); kw.setdefault("size", size)
        add_para(tf, [(txt, kw)], first=(i == 0), space_after=sa, level=lv, line=1.15)
    return s

def bullets_figure(prs, title, items, fig, active=None, cite=None, cite_url=None, size=22, fig_w=Cm(12.5)):
    """Left: bullet list. Right: a supporting figure (e.g. 4.5 checklist + annotated waveform)."""
    s = content_slide(prs, title, active=active, cite=cite, cite_url=cite_url)
    tb, tf = _txbox(s, Cm(2.4), Cm(6.2), Cm(15.6), Cm(11.0), anchor=MSO_ANCHOR.MIDDLE)
    for i, (txt, kw) in enumerate(items):
        kw = dict(kw); sa = kw.pop("sa", 16); lv = kw.pop("lv", 0); kw.setdefault("size", size)
        add_para(tf, [(txt, kw)], first=(i == 0), space_after=sa, level=lv, line=1.2)
    add_figure(s, fig, x=Cm(18.6), y=Cm(5.4), w=fig_w, h=Cm(11.6 if not cite else 11.0))
    return s

def column(prs, title, heading, paras, active=None, cite=None, cite_url=None, accent=TEAL):
    """Single accent-bordered box: heading + list-of-paragraph strings."""
    s = content_slide(prs, title, active=active, cite=cite, cite_url=cite_url)
    add_box(s, Cm(2.2), Cm(5.3), Cm(29.5), Cm(12.2 if not cite else 11.5), RGBColor(0xF5, 0xF7, 0xF9), line=accent)
    tb, tf = _txbox(s, Cm(3.0), Cm(5.9), Cm(27.9), Cm(11.0 if not cite else 10.3))
    add_para(tf, [(heading, dict(size=26, bold=True, color=accent))], first=True, space_after=16, line=1.1)
    for para in paras:
        add_para(tf, [(para, dict(size=22, color=INK))], space_after=14, line=1.25)
    return s

_ROW_ACCENTS = [BLUE, ORANGE, RGBColor(0x70, 0xAD, 0x47), GOLD, TEAL]

def column_row(prs, title, items, active=None, cite=None, cite_url=None):
    """items: list of (heading, [paras], accent). Boxes laid out side by side.
    Cap at 2-3 boxes: at the 22pt body-text floor, a 4th column no longer has
    enough width per card — split into two slides instead (see スライド作成ルール.md
    "入りきらないときの正解は内容を削る／スライドを分けること")."""
    s = content_slide(prs, title, active=active, cite=cite, cite_url=cite_url)
    n = len(items)
    assert n <= 3, f"column_row: {n} boxes won't fit the 22pt floor — split into slides instead"
    top = Cm(5.3); h = Cm(12.2 if not cite else 11.5)
    gap = Cm(0.6); total_w = Cm(29.5)
    bw = Emu(int((total_w - gap * (n - 1)) / n))
    x = Cm(2.2)
    for i, (heading, paras, accent) in enumerate(items):
        accent = accent or _ROW_ACCENTS[i % len(_ROW_ACCENTS)]
        add_box(s, x, top, bw, h, RGBColor(0xF5, 0xF7, 0xF9), line=accent)
        tb, tf = _txbox(s, x + Cm(0.35), top + Cm(0.5), bw - Emu(int(Cm(0.7))), h - Cm(1.0))
        add_para(tf, [(heading, dict(size=22, bold=True, color=accent))], first=True, space_after=10, line=1.1)
        for para in paras:
            add_para(tf, [(para, dict(size=22, color=INK))], space_after=8, line=1.15)
        x = Emu(int(x + bw + gap))
    return s

def column_grid(prs, title, items, active=None, cite=None, cite_url=None, cols=2):
    """items: list of (heading, [paras], accent). Boxes laid out in a `cols`-wide grid.
    Cap at 2-3 boxes total: at the 22pt body-text floor a 2x2+ grid no longer has
    enough height per card — split into two slides (column_row, 2-3 cards each) instead."""
    s = content_slide(prs, title, active=active, cite=cite, cite_url=cite_url)
    n = len(items)
    assert n <= 3, f"column_grid: {n} boxes won't fit the 22pt floor — split into slides instead"
    rows = (n + cols - 1) // cols
    top0 = Cm(5.3); area_h = Cm(12.2 if not cite else 11.5); area_w = Cm(29.5)
    gap_x = Cm(0.5); gap_y = Cm(0.45)
    bw = Emu(int((area_w - gap_x * (cols - 1)) / cols))
    bh = Emu(int((area_h - gap_y * (rows - 1)) / rows))
    for i, (heading, paras, accent) in enumerate(items):
        accent = accent or _ROW_ACCENTS[i % len(_ROW_ACCENTS)]
        r, c = divmod(i, cols)
        x = Emu(int(Cm(2.2) + c * (bw + gap_x)))
        y = Emu(int(top0 + r * (bh + gap_y)))
        add_box(s, x, y, bw, bh, RGBColor(0xF5, 0xF7, 0xF9), line=accent)
        tb, tf = _txbox(s, x + Cm(0.35), y + Cm(0.3), bw - Emu(int(Cm(0.7))), bh - Cm(0.6))
        add_para(tf, [(heading, dict(size=22, bold=True, color=accent))], first=True, space_after=6, line=1.05)
        for para in paras:
            add_para(tf, [(para, dict(size=22, color=INK))], space_after=5, line=1.1)
    return s

def menu_slide(prs, title, full_list, sub_list=None):
    """2-column x 4-row layout (chapter name only, no sublines) — this is a manually
    positioned textbox (no placeholder autofit), so the height budget must be safe even
    if a long chapter name wraps to 2 lines. At 22pt/2-col-width, worst case (all 8 wrap)
    is ~9.1cm, well inside the ~12.9cm available — a single long line never overflows."""
    s = content_slide(prs, title, active=-1)
    n = len(full_list)
    half = (n + 1) // 2
    col_w = Cm(14.0)
    for c, chunk_start in enumerate([0, half]):
        chunk = full_list[chunk_start:chunk_start + half]
        tb, tf = _txbox(s, Cm(2.4) + c * Cm(14.8), Cm(5.6), col_w, Cm(12.2), anchor=MSO_ANCHOR.MIDDLE)
        for i, name in enumerate(chunk):
            idx = chunk_start + i
            add_para(tf, [(f"{idx+1}.  ", dict(size=22, bold=True, color=GOLD)),
                          (name, dict(size=22, bold=True, color=INK))],
                     first=(i == 0), space_after=22, line=1.15)
    return s

def divider(prs, idx, full_list, short_list, subs_map):
    """8-chapter section divider: left = all 8 chapter headings (active=gold), right =
    current chapter's numbered sub-sections in gold. subs_map[idx] = [(code,label), ...].
    Both columns are manually positioned textboxes (no placeholder autofit) — sized so
    that even worst-case line-wraps of the long FULL chapter names stay within the box."""
    s = content_slide(prs, f"{idx+1}. {short_list[idx]}", active=idx)
    # left: 8 chapter headings, single font size (22pt, meets body-text minimum) so
    # active/inactive differ by color+weight only, not size — keeps height predictable.
    tbL, tfL = _txbox(s, Cm(2.2), Cm(5.4), Cm(15.0), Cm(12.9), anchor=MSO_ANCHOR.MIDDLE)
    for i, name in enumerate(full_list):
        on = (i == idx)
        add_para(tfL, [(f"{i+1}. ", dict(size=22, bold=True, color=GOLD if on else LGRAY)),
                       (name, dict(size=22, bold=on, color=INK if on else MGRAY))],
                 first=(i == 0), space_after=9, line=1.12)
    # right: current chapter subsections, gold code + label.
    subs = subs_map[idx]
    tbR, tfR = _txbox(s, Cm(17.7), Cm(5.7), Cm(13.9), Cm(12.6), anchor=MSO_ANCHOR.MIDDLE)
    for i, (code, label) in enumerate(subs):
        add_para(tfR, [(code + "  ", dict(size=22, bold=True, color=GOLD)),
                       (label, dict(size=22, color=INK))], first=(i == 0), space_after=18, line=1.15)
    return s

def table_figure_slide(prs, title, fig, active=None, cite=None, cite_url=None):
    """Wide table image (no caption strip needed)."""
    s = content_slide(prs, title, active, cite, cite_url)
    add_figure(s, fig, x=Cm(1.4), y=Cm(5.1), w=Cm(31.0), h=Cm(12.6 if not cite else 12.0))
    return s

def references(prs, refs):
    """refs: list of (text, url_or_None). 2-column, gray, hyperlinked, 10.5pt."""
    s = content_slide(prs, "参考文献", active=None)
    col_w = Cm(14.6)
    half = (len(refs) + 1) // 2
    for x, chunk in [(Cm(2.2), refs[:half]), (Cm(17.4), refs[half:])]:
        tb, tf = _txbox(s, x, Cm(5.3), col_w, Cm(13.0))
        for i, (txt, url) in enumerate(chunk):
            p = add_para(tf, [(txt, dict(size=10.5, color=GRAY))], first=(i == 0), space_after=8, line=1.15)
            if url and p.runs:
                p.runs[-1].hyperlink.address = url
                p.runs[-1].font.underline = False
                p.runs[-1].font.color.rgb = GRAY
    return s

def set_notes(slide, text):
    """Attach speaker notes. Used to preserve the full original wording of any
    on-slide text that was shortened to meet the 22pt/no-overflow rules — the
    cut detail must survive verbatim here, never be discarded."""
    slide.notes_slide.notes_text_frame.text = text

def add_pagenumbers(prs):
    slides = list(prs.slides)
    numbered = slides[1:]  # skip title
    total = len(numbered)
    for i, s in enumerate(numbered, start=1):
        tb, tf = _txbox(s, Cm(30.7), Cm(18.35), Cm(2.8), Cm(0.6))
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"{i}/{total}"; _set_font(r, size=12, color=GRAY)
