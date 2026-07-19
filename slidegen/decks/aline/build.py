# -*- coding: utf-8 -*-
"""観血的動脈圧測定 講義スライド 組版（濃紺×ゴールド / Meiryo）。
Opus本文(aline_content.json)＋図版(figs/)から .pptx を生成。"""
import json, os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE=os.path.dirname(os.path.abspath(__file__)); FIG=os.path.join(HERE,"figs")
# ---- palette ----
NAVY=RGBColor(0x1E,0x2A,0x38); GOLD=RGBColor(0xD8,0xAE,0x3C); GOLD_D=RGBColor(0xB8,0x86,0x0B)
INK=RGBColor(0x26,0x26,0x26); WHITE=RGBColor(0xFF,0xFF,0xFF); SLATE=RGBColor(0x8A,0x94,0xA0)
SLATE2=RGBColor(0x6B,0x74,0x80); HAIR=RGBColor(0xC7,0xCE,0xD6); GRAY=RGBColor(0x80,0x80,0x80)
BLUE=RGBColor(0x2E,0x6F,0xBF); ORANGE=RGBColor(0xD8,0x70,0x2B); RED=RGBColor(0xC0,0x00,0x00)
TEAL=RGBColor(0x00,0xA0,0xA2); GREEN=RGBColor(0x5E,0x9E,0x3E)
NAVY_CHIP=RGBColor(0x24,0x33,0x45); LGRAY=RGBColor(0xED,0xF0,0xF3); CARDBG=RGBColor(0xF6,0xF7,0xF9)
F="Meiryo"

SECTIONS=[("歴史","直接測定はどう生まれ、臨床に来たか"),
          ("成り立ち","測定系の構成と動特性（自然周波数・減衰）"),
          ("基礎","動脈圧波形の解剖と前進波＋反射波"),
          ("正常","正常値・正常波形と正しい校正"),
          ("見どころ","麻酔中：拍動ごとの情報をどう使うか"),
          ("異常","系のアーチファクトと病態を映す波形"),
          ("対応","トラブルシューティングと合併症・安全"),
          ("特殊","パルスコンター・GDT・特殊状況")]
CHIPS=["歴史","成り立ち","基礎","正常","見どころ","異常","対応","特殊"]

try: CONTENT={s["id"]:s for s in json.load(open(os.path.join(HERE,"aline_content.json")))["slides"]}
except Exception: CONTENT={}
try: REFS=json.load(open(os.path.join(HERE,"aline_content.json"))).get("references",[])
except Exception: REFS=[]

def C(sid,key,default=None):
    d=CONTENT.get(sid,{})
    return d.get(key,default)

prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
BLANK=prs.slide_layouts[6]
EMW=prs.slide_width; EMH=prs.slide_height

def slide(): return prs.slides.add_slide(BLANK)

def _spacing(para,pts):
    pPr=para._p.get_or_add_pPr(); ln=pPr.makeelement(qn('a:lnSpc'),{}); pct=pPr.makeelement(qn('a:spcPct'),{'val':str(int(pts*1000))})
    ln.append(pct); pPr.insert(0,ln)

def rect(sl,x,y,w,h,fill=None,line=None,lw=1.0,shape=MSO_SHAPE.RECTANGLE,shadow=False):
    sp=sl.shapes.add_shape(shape,Inches(x),Inches(y),Inches(w),Inches(h))
    if fill is None: sp.fill.background()
    else: sp.fill.solid(); sp.fill.fore_color.rgb=fill
    if line is None: sp.line.fill.background()
    else: sp.line.color.rgb=line; sp.line.width=Pt(lw)
    sp.shadow.inherit=False
    if shadow:
        el=sp._element.spPr; from pptx.oxml.ns import qn as _q
    return sp

def bg(sl,color):
    r=rect(sl,-0.06,-0.06,13.46,7.62,fill=color); sl.shapes._spTree.remove(r._element); sl.shapes._spTree.insert(2,r._element); return r

def txt(sl,x,y,w,h,runs,align=PP_ALIGN.LEFT,anchor=MSO_ANCHOR.TOP,wrap=True,line_pts=None,space_after=2):
    """runs: list of paragraphs; each para = list of (text,size,color,bold,italic) or a single tuple."""
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame
    tf.word_wrap=wrap; tf.vertical_anchor=anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    if isinstance(runs,tuple): runs=[[runs]]
    first=True
    for para in runs:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.alignment=align; p.space_after=Pt(space_after); p.space_before=Pt(0)
        if line_pts: _spacing(p,line_pts)
        if isinstance(para,tuple): para=[para]
        for r in para:
            text,size,color,bold,ital=(list(r)+[False,False])[:5]
            run=p.add_run(); run.text=text; f=run.font
            f.size=Pt(size); f.name=F; f.bold=bold; f.italic=ital; f.color.rgb=color
    return tb

def kicker(sl,x,y,text,color=GOLD,size=12,square=True):
    if square: rect(sl,x,y+0.03,0.12,0.12,fill=color)
    txt(sl,x+(0.22 if square else 0),y-0.03,6,0.3,[(text,size,color,True,False)])

def page_no(sl,n):
    txt(sl,12.0,7.03,1.1,0.3,[(f"{n} / {TOTAL}",10,GRAY,False,False)],align=PP_ALIGN.RIGHT)

def citation(sl,cites):
    if not cites: return
    s=" ／ ".join(cites)
    txt(sl,0.55,7.03,10.5,0.3,[(s,10,GRAY,False,False)])

def chips(sl,active):  # active:1..8
    n=len(CHIPS); cw=0.70; gap=0.055; total=n*cw+(n-1)*gap
    x0=13.333-0.55-total; y=0.42; h=0.70
    for i,lab in enumerate(CHIPS):
        x=x0+i*(cw+gap); on=(i+1==active)
        rect(sl,x,y,cw,h,fill=(NAVY if on else LGRAY),shape=MSO_SHAPE.ROUNDED_RECTANGLE,
             line=(None if on else HAIR),lw=0.75)
        txt(sl,x,y+0.06,cw,0.28,[(str(i+1),13,(WHITE if on else SLATE2),True,False)],align=PP_ALIGN.CENTER)
        txt(sl,x-0.02,y+0.36,cw+0.04,0.28,[(lab,8.5,(WHITE if on else SLATE2),on,False)],align=PP_ALIGN.CENTER)

def content_head(sl,sec_no,title,active):
    kicker(sl,0.6,0.46,CHIPS[sec_no-1],color=GOLD_D,size=12)
    txt(sl,0.55,0.70,6.05,0.7,[(title,26,INK,True,False)])
    chips(sl,active)

def bullets_box(sl,x,y,w,h,items,size=15,dot=GOLD):
    tb=sl.shapes.add_textbox(Inches(x),Inches(y),Inches(w),Inches(h)); tf=tb.text_frame; tf.word_wrap=True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    first=True
    for it in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_after=Pt(9); p.space_before=Pt(0); _spacing(p,118)
        r1=p.add_run(); r1.text="● "; r1.font.size=Pt(size-3); r1.font.color.rgb=dot; r1.font.name=F; r1.font.bold=True
        r2=p.add_run(); r2.text=it; r2.font.size=Pt(size); r2.font.color.rgb=INK; r2.font.name=F
    return tb

def img_fit(sl,path,x,y,w,h):
    """箱(w,h)にアスペクト維持で内接配置(中央)。"""
    from PIL import Image
    iw,ih=Image.open(path).size; ar=iw/ih; box=w/h
    if ar>box: nw=w; nh=w/ar
    else: nh=h; nw=h*ar
    sl.shapes.add_picture(path,Inches(x+(w-nw)/2),Inches(y+(h-nh)/2),Inches(nw),Inches(nh))

# ---------------- slides ----------------
def cover():
    sl=slide(); bg(sl,NAVY)
    sl.shapes.add_picture(os.path.join(FIG,"deco_gold.png"),Inches(0),Inches(5.7),Inches(13.333),Inches(1.7))
    rect(sl,0.9,5.62,11.5,0.006,fill=SLATE2)
    kicker(sl,0.95,1.55,"A R T E R I A L   P R E S S U R E   M O N I T O R I N G",color=GOLD,size=14,square=False)
    txt(sl,0.9,2.05,11.5,1.2,[("観血的動脈圧測定",50,WHITE,True,False)])
    txt(sl,0.95,3.35,11.5,0.6,[("― 歴史・成り立ちから、波形の読みと応用まで ―",21,SLATE,False,True)])
    txt(sl,0.95,4.5,8,0.5,[("川副 靖晃",20,WHITE,False,False)])

def contents():
    sl=slide(); bg(sl,WHITE)
    kicker(sl,0.6,0.55,"CONTENTS",color=GOLD_D,size=13)
    txt(sl,0.55,0.85,8,1.0,[("メニュー",40,INK,True,False)])
    y=2.15
    for i,(name,sub) in enumerate(SECTIONS):
        txt(sl,1.0,y,0.8,0.5,[(f"{i+1}",22,GOLD,True,False)])
        txt(sl,1.65,y+0.02,10.5,0.5,[[(name,20,INK,True,False),("　—　"+sub,14,SLATE2,False,False)]])
        y+=0.585

def divider(no):
    sl=slide(); bg(sl,NAVY)
    name,sub=SECTIONS[no-1]
    sl.shapes.add_picture(os.path.join(FIG,"deco_gold_faint.png"),Inches(6.8),Inches(5.9),Inches(6.5),Inches(1.4))
    txt(sl,0.95,1.55,6,0.4,[("S E C T I O N",14,SLATE,False,False)])
    txt(sl,0.7,1.95,3.4,2.6,[(f"0{no}",150,GOLD,True,False)])
    rect(sl,4.45,2.35,0.014,2.7,fill=SLATE2)
    txt(sl,4.8,2.45,8.0,1.1,[(name,54,WHITE,True,False)])
    txt(sl,4.83,3.75,8.0,0.6,[(sub,20,SLATE,False,False)])
    subs=SUBITEMS.get(no,[])
    y=4.55
    for code,label in subs:
        rect(sl,4.86,y+0.07,0.11,0.11,fill=GOLD)
        txt(sl,5.12,y-0.02,7.5,0.4,[[(code+"　",15,GOLD,True,False),(label,15,RGBColor(0xE7,0xEA,0xEE),False,False)]])
        y+=0.5

SUBITEMS={
 1:[("1.1","直接測定の夜明け"),("1.2","臨床モニタリングへ")],
 2:[("2.1","測定系の構成"),("2.2","圧→電気信号"),("2.3","動特性：自然周波数と減衰"),("2.4","fast-flush テスト")],
 3:[("3.1","動脈圧波形の解剖"),("3.2","前進波＋反射波"),("3.3","中枢→末梢の変化"),("3.4","波形から得る指標")],
 4:[("4.1","正常値と正常波形"),("4.2","正しい校正と最適ダンピング")],
 5:[("5.1","拍動ごとのリアルタイムBP"),("5.2","PPV/SVVと輸液反応性"),("5.3","収縮期圧変動 SPV"),("5.4","波形の質的読み")],
 6:[("6.1","オーバーダンピング"),("6.2","アンダーダンピング・共振"),("6.3","奇脈・交互脈"),("6.4","弁膜症などの特徴波形")],
 7:[("7.1","ダンピング異常への対応"),("7.2","合併症と安全")],
 8:[("8.1","パルスコンター心拍出量"),("8.2","目標指向的循環管理 GDT"),("8.3","特殊状況")],
}

def content_fig(sid,sec,title,fig,active,n,fig_w=7.05,tall=False):
    sl=slide(); bg(sl,WHITE); content_head(sl,sec,title,active)
    hl=C(sid,"headline")
    if hl: txt(sl,0.6,1.42,12.2,0.5,[(hl,16,GOLD_D,True,False)])
    by=1.98
    if tall:  # figure full width, bullets below in 2 cols
        img_fit(sl,os.path.join(FIG,fig),0.6,by,12.1,3.35)
        items=C(sid,"bullets",[])
        half=(len(items)+1)//2
        bullets_box(sl,0.7,5.5,6.0,1.4,items[:half],size=13.5)
        bullets_box(sl,6.9,5.5,6.0,1.4,items[half:],size=13.5)
    else:
        img_fit(sl,os.path.join(FIG,fig),0.5,by,fig_w,4.35)
        bullets_box(sl,fig_w+0.65,2.1,12.78-(fig_w+0.65),4.2,C(sid,"bullets",[]),size=14.5)
    citation(sl,C(sid,"citations",[])); page_no(sl,n); return sl

# ---------- native-shape content slides ----------
def arrow_between(sl,x1,y,x2,color=SLATE2):
    ln=sl.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,Inches(x1),Inches(y),Inches(x2-x1),Inches(0.16))
    ln.fill.solid(); ln.fill.fore_color.rgb=color; ln.line.fill.background(); ln.shadow.inherit=False

def s1_1(n):  # 歴史 timeline
    sid="s1_1"; sl=slide(); bg(sl,WHITE); content_head(sl,1,"1.1  直接測定の夜明け",1)
    hl=C(sid,"headline","動脈圧を“直接”測る発想は 18〜19 世紀に確立した")
    txt(sl,0.6,1.42,12.2,0.5,[(hl,16,GOLD_D,True,False)])
    nodes=[("1733","Stephen Hales","ウマの頸動脈にガラス管を挿入し、\n血柱の高さで動脈圧を初めて可視化"),
           ("1828","J. Poiseuille","水銀圧力計 (hémodynamomètre) で\n小型化し mmHg 表示を可能に"),
           ("1847","Carl Ludwig","キモグラフで圧を連続的に\n記録（波形の誕生）")]
    y=4.15; rect(sl,1.2,y,11.0,0.02,fill=GOLD)
    xs=[2.4,6.6,10.5]
    for (yr,name,desc),x in zip(nodes,xs):
        rect(sl,x-0.13,y-0.11,0.26,0.26,fill=GOLD,shape=MSO_SHAPE.OVAL)
        txt(sl,x-1.6,y-1.15,3.2,0.9,[[(yr+"　",22,GOLD_D,True,False)],[(name,16,INK,True,False)]],align=PP_ALIGN.CENTER)
        txt(sl,x-1.75,y+0.35,3.5,1.4,[(desc,13,INK,False,False)],align=PP_ALIGN.CENTER,line_pts=120)
    citation(sl,C(sid,"citations",["Hales 1733","Poiseuille 1828","Ludwig 1847"])); page_no(sl,n)

def s1_2(n):  # 非観血 vs 観血
    sid="s1_2"; sl=slide(); bg(sl,WHITE); content_head(sl,1,"1.2  臨床モニタリングへ",1)
    hl=C(sid,"headline","非観血法との対比で、観血モニタの価値と適応を捉える")
    txt(sl,0.6,1.42,12.2,0.5,[(hl,16,GOLD_D,True,False)])
    rows=[("測定様式","間欠（数分ごと）","連続・拍動ごと"),
          ("波形情報","なし（数値のみ）","波形が得られる（前負荷・収縮性の手がかり）"),
          ("精度・不整脈/低灌流","動揺時に不正確になりやすい","基準に近く安定"),
          ("採血","不可","動脈血ガス・検体を反復採取"),
          ("侵襲・合併症","非侵襲","カニュレーション；虚血・感染などの合併症")]
    x0,y0=0.9,2.15; wl,wa,wb=3.0,3.6,5.0; rh=0.78
    # headers
    rect(sl,x0+wl,y0,wa,0.55,fill=LGRAY,line=HAIR,lw=0.75)
    rect(sl,x0+wl+wa,y0,wb,0.55,fill=NAVY)
    txt(sl,x0+wl,y0+0.11,wa,0.4,[("非観血 (NIBP)",15,SLATE2,True,False)],align=PP_ALIGN.CENTER)
    txt(sl,x0+wl+wa,y0+0.11,wb,0.4,[("観血 (A-line)",15,WHITE,True,False)],align=PP_ALIGN.CENTER)
    y=y0+0.55
    for i,(lab,a,b) in enumerate(rows):
        fillc=CARDBG if i%2 else WHITE
        rect(sl,x0,y,wl,rh,fill=LGRAY,line=HAIR,lw=0.5)
        rect(sl,x0+wl,y,wa,rh,fill=fillc,line=HAIR,lw=0.5)
        rect(sl,x0+wl+wa,y,wb,rh,fill=RGBColor(0xF3,0xF0,0xE6),line=HAIR,lw=0.5)
        txt(sl,x0+0.12,y,wl-0.2,rh,[(lab,12.5,INK,True,False)],anchor=MSO_ANCHOR.MIDDLE)
        txt(sl,x0+wl+0.12,y,wa-0.24,rh,[(a,12.5,SLATE2,False,False)],anchor=MSO_ANCHOR.MIDDLE)
        txt(sl,x0+wl+wa+0.14,y,wb-0.28,rh,[(b,12.5,INK,False,False)],anchor=MSO_ANCHOR.MIDDLE)
        y+=rh
    citation(sl,C(sid,"citations",["Riva-Rocci 1896","Lambert & Wood 1947"])); page_no(sl,n)

def s2_1(n):  # 測定系の構成 chain
    sid="s2_1"; sl=slide(); bg(sl,WHITE); content_head(sl,2,"2.1  測定系の構成",2)
    hl=C(sid,"headline","カニューレからモニタまで、系全体を一続きの測定チェーンとして見る")
    txt(sl,0.6,1.42,12.6,0.5,[(hl,16,GOLD_D,True,False)])
    steps=["動脈\nカニューレ","硬い\n非伸展チューブ","三方活栓","加圧バッグ\n持続フラッシュ\n(約3mL/h)","トランス\nデューサ","モニタ"]
    x=0.75; y=2.5; bw=1.8; bh=1.15; gap=0.28
    for i,s in enumerate(steps):
        c=NAVY if i in (4,) else CARDBG
        tc=WHITE if i in (4,) else INK
        rect(sl,x,y,bw,bh,fill=c,line=HAIR,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(sl,x,y,bw,bh,[(s,12.5,tc,True,False)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line_pts=112)
        if i<len(steps)-1: arrow_between(sl,x+bw+0.02,y+bh/2-0.08,x+bw+gap-0.02)
        x+=bw+gap
    # note
    rect(sl,0.75,4.2,11.8,0.02,fill=HAIR)
    bullets_box(sl,0.9,4.55,11.6,2.0,C(sid,"bullets",[
        "ゼロ校正：三方活栓を大気開放にして 0 mmHg を較正",
        "レベリング：トランスデューサを phlebostatic axis（第4肋間・中腋窩線＝右房レベル）に合わせる",
        "高さ誤差 1 cmH₂O ≒ 0.74 mmHg：体位変換のたびに再レベリングが必要",
    ]),size=14.5)
    citation(sl,C(sid,"citations",["Gardner 1981"])); page_no(sl,n)

def s7_1(n):  # troubleshooting flow
    sid="s7_1"; sl=slide(); bg(sl,WHITE); content_head(sl,7,"7.1  ダンピング異常への対応",7)
    hl=C(sid,"headline","まず fast-flush で系を評価し、原因に応じて是正する")
    txt(sl,0.6,1.42,12.6,0.5,[(hl,16,GOLD_D,True,False)])
    cols=[("オーバーダンプ",BLUE,["気泡・凝血・フィブリンを除去","チューブを短く硬く／活栓・延長を減らす","カテーテルの位置・キンクを確認"]),
          ("アンダーダンプ（共振）",RED,["延長チューブを短縮","必要ならダンピングデバイス\n(レゾネータ/ROSE)","fast-flush で ζ を再評価"]),
          ("NIBP と乖離",GOLD_D,["波形品質と fast-flush で信頼性判断","末梢増幅・末梢循環不全も考慮","再ゼロ・再レベリングを実施"])]
    x=0.75; w=3.9; gap=0.28; y=2.35
    for name,col,items in cols:
        rect(sl,x,y,w,0.6,fill=col,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(sl,x,y+0.13,w,0.4,[(name,15,WHITE,True,False)],align=PP_ALIGN.CENTER)
        rect(sl,x,y+0.75,w,3.15,fill=CARDBG,line=HAIR,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        bullets_box(sl,x+0.22,y+0.98,w-0.4,2.8,items,size=13.5,dot=col)
        x+=w+gap
    citation(sl,C(sid,"citations",["Gardner 1981"])); page_no(sl,n)

def s7_2(n):  # complications
    sid="s7_2"; sl=slide(); bg(sl,WHITE); content_head(sl,7,"7.2  合併症と安全",7)
    hl=C(sid,"headline","頻度は低いが重篤化しうる合併症を、部位選択と手技で予防する")
    txt(sl,0.6,1.42,12.6,0.5,[(hl,16,GOLD_D,True,False)])
    comps=[("血栓・遠位虚血","細径カテ・橈骨選択・灌流の観察"),
           ("感染・血流感染","無菌操作・閉鎖式・留置期間の管理"),
           ("出血・血腫","抜去後の十分な圧迫止血"),
           ("薬剤の動脈内誤注入","三方活栓の識別・ラベル；重篤な遠位虚血に注意"),
           ("神経障害・空気塞栓","穿刺部位の解剖把握・気泡除去")]
    x0,y0=0.75,2.3; w=3.85; h=1.15; gx=0.28; gy=0.28
    for i,(t,d) in enumerate(comps):
        r,c=divmod(i,3); x=x0+c*(w+gx); y=y0+r*(h+gy)
        rect(sl,x,y,w,h,fill=CARDBG,line=HAIR,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(sl,x+0.18,y+0.2,0.14,h-0.4,fill=RED)  # small marker (not full stripe)
        txt(sl,x+0.45,y+0.16,w-0.6,0.4,[(t,14,INK,True,False)])
        txt(sl,x+0.45,y+0.56,w-0.6,0.5,[(d,11.5,SLATE2,False,False)],line_pts=115)
    # Allen callout
    x=x0+2*(w+gx); y=y0+1*(h+gy)
    rect(sl,x,y,w,h,fill=RGBColor(0xF3,0xF0,0xE6),line=GOLD_D,lw=1.2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(sl,x+0.2,y+0.14,w-0.4,h-0.3,[[("Allen テストの限界　",13,GOLD_D,True,False)],
        [("陰性でも虚血予測能は限定的。側副の確認は参考程度に。",11.5,INK,False,False)]],line_pts=118)
    citation(sl,C(sid,"citations",["Slogoff 1983","Scheer 2002"])); page_no(sl,n)

def s8_2(n):  # GDT flow
    sid="s8_2"; sl=slide(); bg(sl,WHITE); content_head(sl,8,"8.2  目標指向的循環管理 (GDT)",8)
    hl=C(sid,"headline","動的指標＋SV最適化で輸液を、Ea_dyn で昇圧薬反応性を予測する")
    txt(sl,0.6,1.42,12.6,0.5,[(hl,16,GOLD_D,True,False)])
    boxes=[("動的指標が高い\n(PPV/SVV 高)",NAVY,WHITE),("輸液チャレンジ\n(晶質液など)",CARDBG,INK),
           ("SV が 10–15% 増加？",CARDBG,INK)]
    x=0.8; y=2.6; w=2.9; h=1.1; gap=0.55
    for i,(t,fc,tc) in enumerate(boxes):
        rect(sl,x,y,w,h,fill=fc,line=HAIR,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        txt(sl,x,y,w,h,[(t,13.5,tc,True,False)],align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE,line_pts=112)
        if i<len(boxes)-1: arrow_between(sl,x+w+0.06,y+h/2-0.08,x+w+gap-0.06)
        x+=w+gap
    # branch yes/no
    xe=0.8+3*(w)+2*gap-w
    rect(sl,xe+0.2,4.15,2.7,0.7,fill=GREEN,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(sl,xe+0.2,4.28,2.7,0.5,[("反応性あり → 追加輸液",12.5,WHITE,True,False)],align=PP_ALIGN.CENTER)
    rect(sl,xe-3.1,4.15,2.7,0.7,fill=SLATE2,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    txt(sl,xe-3.1,4.28,2.7,0.5,[("非反応 → 輸液中止・他手段",12,WHITE,True,False)],align=PP_ALIGN.CENTER)
    # Eadyn box（フロー図と重複しない2点に絞る＝はみ出し防止）
    bl=C(sid,"bullets",[
        "動的動脈エラスタンス Ea_dyn = PPV / SVV：輸液反応例で昇圧薬による血圧上昇（後負荷応答性）を予測",
        "周術期 GDT は一部で合併症・在院日数を減らすが、集団・アルゴリズム依存でエビデンスは一様でない",
    ])
    bl=bl[-2:] if len(bl)>2 else bl
    rect(sl,0.8,5.12,11.7,1.68,fill=RGBColor(0xF3,0xF0,0xE6),line=GOLD_D,lw=1.0,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
    bullets_box(sl,1.05,5.36,11.15,1.3,bl,size=13.5,dot=GOLD_D)
    citation(sl,C(sid,"citations",["Cannesson 2011","Monge García 2011","Pearse 2014"])); page_no(sl,n)

def s8_3(n):  # special cards
    sid="s8_3"; sl=slide(); bg(sl,WHITE); content_head(sl,8,"8.3  特殊状況",8)
    hl=C(sid,"headline","部位・病態に応じて測定と解釈を調整する")
    txt(sl,0.6,1.42,12.6,0.5,[(hl,16,GOLD_D,True,False)])
    cards=[("人工心肺 (CPB) 後の乖離","一過性に中枢＞橈骨。大腿/上腕での測定を考慮"),
           ("厳密な血圧管理","脳動脈瘤・頸動脈・大血管手術で拍動ごとの制御"),
           ("頻回の動脈採血","血ガス・電解質・乳酸を連続評価"),
           ("小児・特殊集団","細径・部位選択・合併症への配慮")]
    x0,y0=0.75,2.4; w=5.9; h=1.7; gx=0.3; gy=0.3
    for i,(t,d) in enumerate(cards):
        r,c=divmod(i,2); x=x0+c*(w+gx); y=y0+r*(h+gy)
        rect(sl,x,y,w,h,fill=CARDBG,line=HAIR,lw=0.75,shape=MSO_SHAPE.ROUNDED_RECTANGLE)
        rect(sl,x+0.25,y+0.28,0.4,0.4,fill=GOLD,shape=MSO_SHAPE.OVAL)
        txt(sl,x+0.25,y+0.3,0.4,0.4,[(str(i+1),15,WHITE,True,False)],align=PP_ALIGN.CENTER)
        txt(sl,x+0.85,y+0.3,w-1.0,0.5,[(t,16,INK,True,False)])
        txt(sl,x+0.85,y+0.82,w-1.05,0.7,[(d,13,SLATE2,False,False)],line_pts=118)
    citation(sl,C(sid,"citations",["Baba 1976"])); page_no(sl,n)

def summary(n):
    sid="summary"; sl=slide(); bg(sl,NAVY)
    sl.shapes.add_picture(os.path.join(FIG,"deco_gold_faint.png"),Inches(7.0),Inches(6.0),Inches(6.3),Inches(1.3))
    kicker(sl,0.9,0.85,"TAKE-HOME",color=GOLD,size=13,square=False)
    txt(sl,0.85,1.2,11,0.9,[("まとめ",40,WHITE,True,False)])
    items=C(sid,"bullets",[
        "正しくゼロ・レベリングし、最適ダンピング（fast-flush で 1–2 振動）にして初めて数値を信用できる",
        "動脈圧波形＝前進波＋反射波。末梢では増幅し、SBP・脈圧は中枢と異なる",
        "動的指標 PPV/SVV は前提（調節換気・洞調律・十分な一回換気量）を満たす場面でのみ有効",
        "オーバーダンプ＝SBP過小/DBP過大、アンダーダンプ＝SBP過大/DBP過小。fast-flush で鑑別",
        "合併症は稀だが重篤化しうる。部位選択・手技・薬剤誤注入回避で予防",
        "応用：パルスコンター CO・目標指向的循環管理（GDT）へ展開できる",
    ])
    tb=sl.shapes.add_textbox(Inches(0.9),Inches(2.3),Inches(11.6),Inches(4.4)); tf=tb.text_frame; tf.word_wrap=True
    for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
    first=True
    for it in items:
        p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
        p.space_after=Pt(11); _spacing(p,120)
        r1=p.add_run(); r1.text="●  "; r1.font.size=Pt(13); r1.font.color.rgb=GOLD; r1.font.name=F; r1.font.bold=True
        r2=p.add_run(); r2.text=it; r2.font.size=Pt(16); r2.font.color.rgb=RGBColor(0xEC,0xEF,0xF2); r2.font.name=F
    page_no(sl,n)

def references(n):
    sl=slide(); bg(sl,WHITE)
    kicker(sl,0.6,0.55,"REFERENCES",color=GOLD_D,size=13)
    txt(sl,0.55,0.85,8,1.0,[("主要参考文献",34,INK,True,False)])
    refs=REFS if REFS else ["（参考文献はOpusのJSONから取り込まれます）"]
    third=(len(refs)+2)//3
    def col(items,x):
        tb=sl.shapes.add_textbox(Inches(x),Inches(1.95),Inches(4.05),Inches(4.95)); tf=tb.text_frame; tf.word_wrap=True
        for m in ("margin_left","margin_right","margin_top","margin_bottom"): setattr(tf,m,0)
        first=True
        for it in items:
            p=tf.paragraphs[0] if first else tf.add_paragraph(); first=False
            p.space_after=Pt(6); _spacing(p,112)
            r=p.add_run(); r.text=it; r.font.size=Pt(9); r.font.color.rgb=INK; r.font.name=F
    col(refs[:third],0.7); col(refs[third:2*third],4.95); col(refs[2*third:],9.2)
    page_no(sl,n)

# ---------------- assemble ----------------
PLAN=[]  # ordered build actions
def build():
    global TOTAL
    # first pass to compute TOTAL
    order=[]
    order.append(("cover",)); order.append(("contents",))
    figmap={
      "s2_2":("s2_2_transducer.png",7.4,False),"s2_3":("s2_3_fnzeta.png",7.2,False),"s2_4":("s2_4_flush.png",None,True),
      "s3_1":("s3_1_anatomy.png",7.4,False),"s3_2":("s3_2_decomp.png",7.3,False),"s3_3":("s3_3_cenper.png",None,True),"s3_4":("s3_4_indices.png",7.2,False),
      "s4_1":("s4_1_normal.png",7.3,False),"s4_2":("s4_2_optflush.png",7.3,False),
      "s5_1":("s5_1_events.png",None,True),"s5_2":("s5_2_ppv.png",None,True),"s5_3":("s5_3_spv.png",None,True),"s5_4":("s5_4_svr.png",7.6,False),
      "s6_1":("s6_1_overdamp.png",7.4,False),"s6_2":("s6_2_underdamp.png",7.4,False),"s6_3":("s6_3_paradoxus.png",None,True),"s6_4":("s6_4_valve.png",None,True),
      "s8_1":("s8_1_pulsecontour.png",None,True),
    }
    titles={"s2_2":"2.2  圧→電気信号","s2_3":"2.3  動特性：fnとζ","s2_4":"2.4  fast-flush テスト",
     "s3_1":"3.1  動脈圧波形の解剖","s3_2":"3.2  前進波＋反射波","s3_3":"3.3  中枢→末梢の変化","s3_4":"3.4  波形から得る指標",
     "s4_1":"4.1  正常値と正常波形","s4_2":"4.2  校正と最適ダンピング",
     "s5_1":"5.1  リアルタイム BP と MAP","s5_2":"5.2  PPV/SVV と輸液反応性","s5_3":"5.3  収縮期圧変動 (SPV)","s5_4":"5.4  波形の質的読み",
     "s6_1":"6.1  オーバーダンピング","s6_2":"6.2  アンダーダンピング","s6_3":"6.3  奇脈・交互脈","s6_4":"6.4  弁膜症の特徴波形",
     "s8_1":"8.1  パルスコンター CO"}
    secmap={"s2":2,"s3":3,"s4":4,"s5":5,"s6":6,"s8":8}
    layout=[
      ("div",1),("native","s1_1"),("native","s1_2"),
      ("div",2),("native","s2_1"),("fig","s2_2"),("fig","s2_3"),("fig","s2_4"),
      ("div",3),("fig","s3_1"),("fig","s3_2"),("fig","s3_3"),("fig","s3_4"),
      ("div",4),("fig","s4_1"),("fig","s4_2"),
      ("div",5),("fig","s5_1"),("fig","s5_2"),("fig","s5_3"),("fig","s5_4"),
      ("div",6),("fig","s6_1"),("fig","s6_2"),("fig","s6_3"),("fig","s6_4"),
      ("div",7),("native","s7_1"),("native","s7_2"),
      ("div",8),("fig","s8_1"),("native","s8_2"),("native","s8_3"),
      ("summary",),("references",),
    ]
    TOTAL=2+len(layout)
    n=0
    cover(); n+=1
    contents(); n+=1
    natives={"s1_1":s1_1,"s1_2":s1_2,"s2_1":s2_1,"s7_1":s7_1,"s7_2":s7_2,"s8_2":s8_2,"s8_3":s8_3}
    for item in layout:
        n+=1
        if item[0]=="div": divider(item[1])
        elif item[0]=="native": natives[item[1]](n)
        elif item[0]=="summary": summary(n)
        elif item[0]=="references": references(n)
        elif item[0]=="fig":
            sid=item[1]; fig,fw,tall=figmap[sid]; sec=secmap[sid[:2]]
            content_fig(sid,sec,titles[sid],fig,int(sid[1]) if sid[1].isdigit() else sec,n,
                        fig_w=(fw or 7.05),tall=tall)
    prs.save(os.path.join(HERE,"aline.pptx"))
    print("saved aline.pptx  slides=",len(prs.slides._sldIdLst))

if __name__=="__main__":
    build()
