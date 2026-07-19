# -*- coding: utf-8 -*-
"""川副式スライド ビルダー（再利用可能ヘルパー）.
テンプレート _template_kawazoe.pptx の上に、金色タイトル＋下線・パンくず・ページ番号・
出典を自動付与しながらスライドを組み立てる。ルールは スライド作成ルール.md を参照。"""
import os
import sys
from pptx import Presentation
from pptx.util import Cm, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from PIL import Image

# ---- palette (deck spec) ----
GOLD   = RGBColor(0xBF,0x90,0x00)   # accent4 darker 25% : title / key / emphasis
TEAL   = RGBColor(0x00,0xA8,0xAA)   # breadcrumb active
TEALD  = RGBColor(0x00,0x80,0x7F)
INK    = RGBColor(0x26,0x26,0x26)
GRAY   = RGBColor(0x80,0x80,0x80)
LGRAY  = RGBColor(0xD9,0xD9,0xD9)
MGRAY  = RGBColor(0xA6,0xA6,0xA6)
BLUE   = RGBColor(0x1F,0x4E,0x79); BLUEL = RGBColor(0xDD,0xEB,0xF7)
ORANGE = RGBColor(0x8A,0x3B,0x00); ORANGEL=RGBColor(0xFB,0xE5,0xD6)
GOLDBG = RGBColor(0xF7,0xEC,0xCF)
RED    = RGBColor(0xC0,0x00,0x00)
WHITE  = RGBColor(0xFF,0xFF,0xFF)
JP = "Meiryo"

# 8-chapter breadcrumb labels (モニター心電図)
CHAPTERS = ["歴史", "生成", "基礎", "正常", "術中", "異常", "対応", "応用"]

SW, SH = Cm(33.867), Cm(19.05)   # 16:9

def new_deck(template):
    return Presentation(template)

def _set_font(run, size=None, bold=None, color=None, name=JP):
    f = run.font
    if size is not None: f.size = Pt(size)
    if bold is not None: f.bold = bold
    if color is not None: f.color.rgb = color
    if name:
        f.name = name
        rPr = run._r.get_or_add_rPr()
        for tag in ("a:latin","a:ea","a:cs"):
            el = rPr.find(qn(tag))
            if el is None:
                el = rPr.makeelement(qn(tag), {}); rPr.append(el)
            el.set("typeface", name)

def _txbox(slide, x, y, w, h, anchor=MSO_ANCHOR.TOP, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame; tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    for m in ("margin_left","margin_right","margin_top","margin_bottom"):
        setattr(tf, m, 0)
    return tb, tf

def add_para(tf, runs, align=PP_ALIGN.LEFT, space_after=6, level=0, line=None, first=False):
    """runs: list of (text, dict-of-font-kwargs)."""
    p = tf.paragraphs[0] if first and tf.paragraphs[0].text=="" and not tf.paragraphs[0].runs else tf.add_paragraph()
    p.alignment = align; p.level = level
    if space_after is not None: p.space_after = Pt(space_after)
    if line is not None:
        p.line_spacing = line
    for txt, kw in runs:
        r = p.add_run(); r.text = txt; _set_font(r, **kw)
    return p

# ---------- title (gold, from layout) ----------
# The gold underline lives at a FIXED y in the layout, so the title MUST stay on
# one line. Auto-shrink the font (default 46pt) to fit the content width.
def _title_units(text):
    u = 0.0
    for ch in text:
        o = ord(ch)
        if ch == " ": u += 0.30
        elif o <= 0x24F: u += 0.56           # latin / digits
        elif 0x2080 <= o <= 0x208E: u += 0.56  # subscripts (₂ …)
        else: u += 1.0                        # CJK / full-width punct / arrows
    return u

TITLE_MIN_LEFT = Cm(1.6)   # タイトル左マージンの下限（規約 §1）
TITLE_RIGHT_LIMIT_BC = Cm(18.8)  # パンくずがある面: タイトル枠の右端はここまで（安全マージン込み）
JP_W = 1.12  # 本番PowerPointのMeiryoはLibreOffice/Hiraginoより横幅が広い（slide_lint.py と同係数）

def _set_title(slide, text, cap=54, maxw=20.5, breadcrumb=False):
    # Title font is fixed at 44pt (never auto-shrunk). If the text would overflow
    # the usable width at 44pt, warn to stderr — the fix is fewer characters, not
    # a smaller font.
    # breadcrumb=True の面（右上にパンくずがある）では、タイトル枠自体の右端を
    # TITLE_RIGHT_LIMIT_BC 以下に物理的に狭め word_wrap=True にする。こうすると
    # 万一44ptで収まらない題名でも「パンくずに重なる」のではなく「2行化」として
    # 目視・lintで検出可能になる（再発防止。スライド作成ルール.md 参照）。
    size = 44.0
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = text
            # 44pt固定を強制: 自動縮小を無効化（レイアウトのTEXT_TO_FIT_SHAPEに依存しない）
            tf = ph.text_frame
            tf.auto_size = MSO_AUTO_SIZE.NONE
            if ph.left is not None and ph.left < TITLE_MIN_LEFT:
                ph.left = TITLE_MIN_LEFT
            if breadcrumb and ph.left is not None:
                new_w = TITLE_RIGHT_LIMIT_BC - ph.left
                if new_w > 0:
                    ph.width = int(new_w)
                tf.word_wrap = True
            else:
                tf.word_wrap = False
            u = _title_units(text)
            w_at_size = u * (size/72*2.54) * JP_W    # width in cm at fixed size (Meiryo安全係数込み)
            if u > 0 and w_at_size > maxw:
                print(f"[title-overflow] {w_at_size:.1f}cm>{maxw} : {text}", file=sys.stderr)
            for p in ph.text_frame.paragraphs:
                for r in p.runs:
                    _set_font(r, size=size, name=JP)   # keep layout gold color
            return ph
    return None

# ---------- speaker notes (削った詳細の保管先) ----------
def set_notes(slide, text):
    """スピーカーノートに原文をそのまま設定する。
    オンスライド短文を抜粋・短縮したときは、削った分をここに全文残す。"""
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = text

def _strip_body(slide):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx == 1:
            ph._element.getparent().remove(ph._element)

# ---------- breadcrumb (numbered downward pennants, packed top-right) ----------
# Number and label BOTH sit inside each pennant, stacked (number on top, label
# below) — no separate badge. Strip is packed into the top-right corner.
def add_breadcrumb(slide, active):
    n = len(CHAPTERS)
    w = Cm(1.50); penH = Cm(1.46); gap = Cm(0.10)   # 8章に合わせて縮小
    total = n*w + (n-1)*gap
    x0 = SW - total - Cm(0.14); y0 = Cm(0.14)
    for i, lab in enumerate(CHAPTERS):
        x = x0 + i*(w+gap)
        act = (i == active)
        col = WHITE if act else GRAY
        pen = slide.shapes.add_shape(MSO_SHAPE.FLOWCHART_OFFPAGE_CONNECTOR, x, y0, w, penH)
        pen.fill.solid(); pen.fill.fore_color.rgb = TEAL if act else LGRAY
        pen.line.fill.background(); pen.shadow.inherit = False
        tf = pen.text_frame; tf.word_wrap = True
        tf.margin_left=0; tf.margin_right=0; tf.margin_top=Cm(0.11); tf.margin_bottom=0
        tf.vertical_anchor = MSO_ANCHOR.TOP
        p1 = tf.paragraphs[0]; p1.alignment = PP_ALIGN.CENTER; p1.line_spacing = 1.0
        r1 = p1.add_run(); r1.text = str(i+1); _set_font(r1, size=9.5, bold=True, color=col)
        p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.CENTER; p2.space_before = Pt(1); p2.line_spacing = 1.0
        r2 = p2.add_run(); r2.text = lab; _set_font(r2, size=11, bold=True, color=col)

# ---------- citation (small gray, optional hyperlink) ----------
def add_citation(slide, text, url=None, x=Cm(2.0), y=Cm(18.35), w=Cm(24), align=PP_ALIGN.LEFT):
    tb, tf = _txbox(slide, x, y, w, Cm(0.6))
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; _set_font(r, size=10.5, color=GRAY)
    if url:
        r.hyperlink.address = url          # clickable link to the paper
        r.font.underline = False           # keep it looking like plain gray text
        r.font.color.rgb = GRAY            # gray, not the blue link colour

# ---------- figure fit ----------
def add_figure(slide, path, x=Cm(2.0), y=Cm(5.0), w=Cm(29.8), h=Cm(12.7), valign="mid"):
    iw, ih = Image.open(path).size
    ar = iw/ih; box_ar = (w/h)
    if ar > box_ar:
        nw = w; nh = int(w/ar)
    else:
        nh = h; nw = int(h*ar)
    px = x + (w-nw)//2
    if valign=="top": py = y
    elif valign=="bottom": py = y + (h-nh)
    else: py = y + (h-nh)//2
    return slide.shapes.add_picture(path, px, py, nw, nh)

# ---------- rounded box + text ----------
def add_box(slide, x, y, w, h, fill, line=None, round=True):
    shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE, x, y, w, h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb = fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb = line; shp.line.width = Pt(1.25)
    shp.shadow.inherit = False
    return shp

# ---------- slide factories ----------
def content_slide(prs, title, active=None, cite=None, cite_url=None):
    s = prs.slides.add_slide(prs.slide_layouts[2])
    bc = active is not None  # breadcrumb present → title box is physically narrowed in _set_title
    # no breadcrumb (active is None) → title may use more width; else warn near the box's real limit
    _set_title(s, title, maxw=(26.0 if not bc else 16.47), breadcrumb=bc); _strip_body(s)
    if active is not None: add_breadcrumb(s, active)
    if cite: add_citation(s, cite, url=cite_url)
    return s

def figure_slide(prs, title, fig, active=None, cite=None, cite_url=None, caption=None):
    s = content_slide(prs, title, active, cite, cite_url)
    yb = Cm(5.0); hb = Cm(12.6 if not cite else 12.0)
    if caption:
        tb, tf = _txbox(s, Cm(2.0), Cm(4.8), Cm(29.8), Cm(1.0), anchor=MSO_ANCHOR.TOP)
        add_para(tf, [(caption, dict(size=16, bold=True, color=INK))], align=PP_ALIGN.CENTER, first=True)
        yb = Cm(6.0); hb = Cm(11.4 if not cite else 10.9)
    add_figure(s, fig, y=yb, h=hb)
    return s

def title_slide(prs, main, sub, hero=None):
    s = prs.slides.add_slide(prs.slide_layouts[0])
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 0:
            ph.text = ""
            add_para(ph.text_frame, [(main, dict(size=40, bold=True, color=INK))], align=PP_ALIGN.CENTER, first=True)
        elif ph.placeholder_format.idx == 1:
            ph.text = ""
            for i,(t,kw) in enumerate(sub):
                add_para(ph.text_frame, [(t,kw)], align=PP_ALIGN.CENTER, first=(i==0), space_after=4)
    if hero:
        add_figure(s, hero, x=Cm(3.5), y=Cm(14.6), w=Cm(26.8), h=Cm(3.6), valign="mid")
    return s

def add_pagenumbers(prs):
    slides = list(prs.slides)
    numbered = slides[1:]  # skip title
    total = len(numbered)
    for i, s in enumerate(numbered, start=1):
        tb, tf = _txbox(s, Cm(30.7), Cm(18.35), Cm(2.8), Cm(0.6))
        p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT
        r = p.add_run(); r.text = f"{i}/{total}"; _set_font(r, size=12, color=GRAY)
