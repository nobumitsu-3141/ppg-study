#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
マンシェット血圧測定スライド ― pptx組版(python-pptx)
- content/scaffold.json + content/chNN.json を読み、design-spec.md の
  「インク＆ゴールド」体系で 47 枚の pptx を生成する。
- 中身(医学的内容)はここでは一切創作しない。JSON の文面をそのまま流し込むだけ。
"""
import json
import math
import os
import sys

from PIL import Image, ImageFont
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

HERE = os.path.dirname(os.path.abspath(__file__))
CONTENT = os.path.join(HERE, "content")
ASSETS = os.path.join(HERE, "assets")
OUT_PATH = os.path.join(HERE, "マンシェット血圧測定_v0.1.pptx")

FONT = "Meiryo"
TOTAL_SLIDES = 47

# ---- カラートークン (design-spec.md 準拠) ----
INK = RGBColor(0x26, 0x26, 0x26)
GOLD = RGBColor(0xBF, 0x90, 0x00)
TEAL = RGBColor(0x00, 0xA8, 0xAA)
MUTED = RGBColor(0x80, 0x80, 0x80)
NAV_OFF = RGBColor(0xD9, 0xD9, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
CALLOUT_FILL = RGBColor(0xF2, 0xE2, 0xB3)

NAV_SHORT = ["歴史", "原理", "基礎", "正常", "術中", "異常", "対応", "特殊"]

# ---- レイアウト定数(in) ----
BODY_LEFT = 0.5
BODY_RIGHT = 12.8
RULE_Y = 1.58
TITLE_Y = 0.72
TITLE_H = 0.80
FOOTER_Y = 7.06

# ---- タイトル自動縮小用のフォント幅測定 ----
# 実行環境に Meiryo 実体がないため、太字の和文ゴシック(Hiragino Sans W6)で幅を近似測定する。
# (本番の PowerPoint/Meiryo とグリフ幅は完全一致しないため、安全マージンを持たせる。)
_TITLE_METRIC_FONT = "/System/Library/Fonts/ヒラギノ角ゴシック W6.ttc"
_title_font_cache = {}


def _measure_text_width_in(text, size_pt):
    font = _title_font_cache.get(size_pt)
    if font is None:
        font = ImageFont.truetype(_TITLE_METRIC_FONT, size=int(round(size_pt)))
        _title_font_cache[size_pt] = font
    return font.getlength(text) / 72.0  # 72pt(px)/inch 換算


def fit_title_size(text, max_width_in, base_size=46, min_size=26):
    size = base_size
    while size > min_size and _measure_text_width_in(text, size) > max_width_in:
        size -= 1
    return size


def _wrap_line_count(text, size_pt, width_in):
    """CJK主体の文はほぼ任意の文字位置で折り返せるため、
    総幅/ボックス幅の切り上げで概算の行数を見積もる。"""
    w = _measure_text_width_in(text, size_pt)
    return max(1, math.ceil(w / width_in - 1e-6))


# 素のフォントサイズに対する「シングル行間」の目安倍率(一般的な和文ゴシックの行送り)。
# paragraph.line_spacing に渡す倍率(例 1.55)は、この基準行高に対する倍数として扱う。
_BASE_LINE_FACTOR = 1.2
# 実測フォントが本番の Meiryo と完全一致しないための安全マージン。
_HEIGHT_SAFETY_MARGIN = 1.08


def stack_height_in(texts, size_pt, width_in, line_mult=1.2, space_after_pt=0):
    """段落を縦に積んだときに必要な高さ(in)を概算する。
    ブレット数や文の長さがスライドごとに違うため、固定高さだと重なりが起きうる ―
    実測フォント幅で折り返し行数を見積もり、コールアウト位置を動的に決めるために使う。"""
    total_pt = 0.0
    for t in texts:
        lines = _wrap_line_count(t, size_pt, width_in)
        total_pt += lines * size_pt * _BASE_LINE_FACTOR * line_mult + space_after_pt
    return total_pt / 72.0 * _HEIGHT_SAFETY_MARGIN


# ======================================================================
# フォント・共通ヘルパー
# ======================================================================
def set_east_asian_font(run, name=FONT):
    """run.font.name は a:latin にしか効かないため a:ea/a:cs にも明示的に typeface を差す。"""
    rPr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        el = rPr.find(qn(tag))
        if el is None:
            el = rPr.makeelement(qn(tag), {})
            rPr.append(el)
        el.set("typeface", name)


def set_run(run, text=None, size=11, color=INK, bold=False, italic=False, font=FONT):
    if text is not None:
        run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = font
    set_east_asian_font(run, font)
    return run


def set_para_indent(paragraph, inches):
    pPr = paragraph._p.get_or_add_pPr()
    marL = int(Inches(inches))
    pPr.set("marL", str(marL))
    pPr.set("indent", "0")


def new_textbox(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    # 既定の spAutoFit (シェイプをテキストに合わせて自動リサイズ) を無効化する。
    # wrap=False と spAutoFit の組み合わせは LibreOffice で意図せぬ折り返しを招くため、
    # 固定サイズのボックスに統一して見た目を安定させる。
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.word_wrap = wrap
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = Emu(0)
    tf.vertical_anchor = anchor
    return box, tf


def add_plain_text(slide, x, y, w, h, text, size=11, color=INK, bold=False,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, wrap=True, italic=False):
    box, tf = new_textbox(slide, x, y, w, h, wrap=wrap, anchor=anchor)
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    set_run(r, text=text, size=size, color=color, bold=bold, italic=italic)
    return box


# ======================================================================
# ナビ帯(右上8タブ) / タイトル+罫線 / フッター / コールアウト / 画像
# ======================================================================
def add_nav(slide, current_chapter):
    """current_chapter: 1-8 で該当章をTEAL表示。None なら全タブ非活性。"""
    tab_w, tab_h, gap = 0.615, 0.50, 0.055
    total_w = 8 * tab_w + 7 * gap
    left_start = BODY_RIGHT - total_w
    top = 0.15
    for i in range(8):
        cnum = i + 1
        x = left_start + i * (tab_w + gap)
        is_current = (current_chapter == cnum)
        fill = TEAL if is_current else NAV_OFF
        tcolor = WHITE if is_current else MUTED
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                        Inches(x), Inches(top), Inches(tab_w), Inches(tab_h))
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
        shape.line.fill.background()
        try:
            shape.adjustments[0] = 0.16
        except Exception:
            pass
        shape.shadow.inherit = False
        tf = shape.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = Emu(0)
        tf.margin_top = tf.margin_bottom = Emu(0)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p0 = tf.paragraphs[0]
        p0.alignment = PP_ALIGN.CENTER
        r0 = p0.add_run()
        set_run(r0, text=str(cnum), size=7.5, color=tcolor, bold=True)
        p1 = tf.add_paragraph()
        p1.alignment = PP_ALIGN.CENTER
        r1 = p1.add_run()
        set_run(r1, text=NAV_SHORT[i], size=8.5, color=tcolor, bold=True)


def add_title_rule(slide, text, size=44):
    max_w = BODY_RIGHT - BODY_LEFT - 0.15
    # タイトルは自動縮小しない(44pt固定ルール)。幅超過は警告のみ出し、監督が文言短縮を判断する。
    if _measure_text_width_in(text, size) > max_w:
        print(f"[title-overflow] : {text}", file=sys.stderr)
    add_plain_text(slide, BODY_LEFT, TITLE_Y, BODY_RIGHT - BODY_LEFT, TITLE_H,
                    text, size=size, color=GOLD, bold=True,
                    align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(BODY_LEFT), Inches(RULE_Y),
                                   Inches(BODY_RIGHT - BODY_LEFT), Pt(3))
    rule.fill.solid()
    rule.fill.fore_color.rgb = GOLD
    rule.line.fill.background()
    rule.shadow.inherit = False


def add_footer(slide, sources, page_num, total=TOTAL_SLIDES):
    if sources:
        src_text = " ／ ".join(sources)
        add_plain_text(slide, BODY_LEFT, FOOTER_Y, 9.3, 0.34, src_text,
                        size=9.5, color=MUTED, anchor=MSO_ANCHOR.MIDDLE)
    add_plain_text(slide, 11.3, FOOTER_Y, 1.5, 0.34, f"{page_num}/{total}",
                    size=10.5, color=MUTED, align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)


def add_callout(slide, x, y, w, h, text, size=13):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = CALLOUT_FILL
    shape.line.color.rgb = GOLD
    shape.line.width = Pt(1.25)
    try:
        shape.adjustments[0] = 0.10
    except Exception:
        pass
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Inches(0.18)
    tf.margin_top = tf.margin_bottom = Inches(0.10)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    set_run(r, text=text, size=size, color=INK, bold=True)


def add_bullets(slide, x, y, w, h, bullets, size=14, color=INK, line_spacing=1.25,
                space_after=8, bold=False):
    box, tf = new_textbox(slide, x, y, w, h, wrap=True, anchor=MSO_ANCHOR.TOP)
    first = True
    for b in bullets:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = line_spacing
        p.space_after = Pt(space_after)
        r = p.add_run()
        set_run(r, text="・" + b, size=size, color=color, bold=bold)
    return box


def add_picture_fit(slide, img_path, x, y, max_w, max_h):
    with Image.open(img_path) as im:
        iw, ih = im.size
    ratio = iw / ih
    box_ratio = max_w / max_h
    if ratio > box_ratio:
        w = max_w
        h = w / ratio
    else:
        h = max_h
        w = h * ratio
    left = x + (max_w - w) / 2
    top = y + (max_h - h) / 2
    slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(w), height=Inches(h))


def add_notes(slide, notes_text):
    if not notes_text:
        return
    notes_slide = slide.notes_slide
    notes_slide.notes_text_frame.text = notes_text


# ======================================================================
# データ読み込み
# ======================================================================
def load_data():
    with open(os.path.join(CONTENT, "scaffold.json"), encoding="utf-8") as f:
        scaffold = json.load(f)
    chapters = {}
    for n in range(1, 9):
        path = os.path.join(CONTENT, f"ch{n:02d}.json")
        with open(path, encoding="utf-8") as f:
            chapters[n] = json.load(f)
    return scaffold, chapters


# ======================================================================
# スライド生成: 表紙
# ======================================================================
def build_cover(prs, blank_layout, scaffold):
    slide = prs.slides.add_slide(blank_layout)
    w = 13.333

    add_plain_text(slide, 0, 3.05, w, 0.95, scaffold["deck_title"],
                    size=40, color=INK, bold=True, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    add_plain_text(slide, 0, 4.02, w, 0.55, scaffold["deck_subtitle"],
                    size=20, color=GOLD, bold=True, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    add_plain_text(slide, 0, 4.58, w, 0.45, scaffold["deck_tertiary"],
                    size=13.5, color=MUTED, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE, wrap=False)
    add_plain_text(slide, 0, 5.35, w, 0.45, scaffold["author"],
                    size=15, color=GOLD, bold=True, align=PP_ALIGN.CENTER,
                    anchor=MSO_ANCHOR.MIDDLE, wrap=False)

    # 最下部中央の金の装飾(カフ+圧計モチーフ: 細線+ダイヤル)
    cx = w / 2
    dial_y = 6.78
    dial_r = 0.16
    dial = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(cx - dial_r), Inches(dial_y - dial_r),
                                   Inches(dial_r * 2), Inches(dial_r * 2))
    dial.fill.background()
    dial.line.color.rgb = GOLD
    dial.line.width = Pt(1.5)
    dial.shadow.inherit = False
    needle = slide.shapes.add_connector(1, Inches(cx), Inches(dial_y), Inches(cx + 0.10), Inches(dial_y - 0.10))
    needle.line.color.rgb = GOLD
    needle.line.width = Pt(1.5)
    for (x0, x1) in [(cx - 2.0, cx - dial_r - 0.08), (cx + dial_r + 0.08, cx + 2.0)]:
        line = slide.shapes.add_connector(1, Inches(x0), Inches(dial_y), Inches(x1), Inches(dial_y))
        line.line.color.rgb = GOLD
        line.line.width = Pt(1.5)

    return slide


# ======================================================================
# スライド生成: 目次(メニュー)
# ======================================================================
def build_toc(prs, blank_layout, scaffold, page_num):
    slide = prs.slides.add_slide(blank_layout)
    add_nav(slide, None)
    add_title_rule(slide, "メニュー", size=32)

    # メニューは各行1行厳守のため word_wrap 無効(タイトルと同じ方針)。
    # 幅は scratchpad/measure.py で全8章の全角文字幅を実測し22ptで1行に収まる字数を確認済み。
    # space_after=9pt: Meiryo安全係数込みの推定(JP_W/JP_H)で8項目の総高がbox高(4.6in)を
    # わずかに超える(+1.06cm)ため、行間だけを詰めて余白を作る(文言・項目数は不変)。
    box, tf = new_textbox(slide, 0.8, 2.0, 10.8, 4.6, wrap=False, anchor=MSO_ANCHOR.TOP)
    first = True
    for ch in scaffold["chapters"]:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_after = Pt(9)
        r1 = p.add_run()
        set_run(r1, text=f"{ch['n']}.  ", size=22, color=GOLD, bold=True)
        r2 = p.add_run()
        set_run(r2, text=ch["full_title"], size=22, color=INK, bold=False)

    add_footer(slide, [], page_num)
    return slide


# ======================================================================
# スライド生成: 章オープナー(節目次)
# ======================================================================
def build_chapter_opener(prs, blank_layout, chapter_num, scaffold, chapters_data, page_num):
    slide = prs.slides.add_slide(blank_layout)
    add_nav(slide, chapter_num)
    add_title_rule(slide, f"{chapter_num}. {NAV_SHORT[chapter_num - 1]}")

    # 章によって現在章の節数が3〜5件と変わり、かつ現在章がリストの最後(8章)に来ると
    # 累積の行高が最大化して footer に迫る/超える。全パターン(最大: 8main行+5節=13行)が
    # 本文22pt下限でも確実に収まるよう、top_y を早め・行間を控えめに固定する。
    # 各行1行厳守のため word_wrap 無効。字数は scratchpad/measure.py で実測済み(全パターン最大35字幅以内)。
    box, tf = new_textbox(slide, 0.8, 1.68, 11.2, 5.30, wrap=False, anchor=MSO_ANCHOR.TOP)
    first = True
    for meta in scaffold["chapters"]:
        n = meta["n"]
        is_current = (n == chapter_num)
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.line_spacing = 1.0
        p.space_after = Pt(3)
        r = p.add_run()
        set_run(r, text=f"{n}. {meta['full_title']}", size=22,
                color=(INK if is_current else MUTED), bold=is_current)
        if is_current:
            for s in chapters_data[n]["slides"]:
                p2 = tf.add_paragraph()
                set_para_indent(p2, 0.55)
                p2.line_spacing = 1.0
                p2.space_after = Pt(2)
                r1 = p2.add_run()
                set_run(r1, text=s["id"] + "   ", size=22, color=GOLD, bold=True)
                r2 = p2.add_run()
                set_run(r2, text=s["title"], size=22, color=INK, bold=False)

    add_footer(slide, [], page_num)
    return slide


# ======================================================================
# スライド生成: 内容ページ
# ======================================================================
def build_content_slide(prs, blank_layout, chapter_num, slide_data, page_num):
    slide = prs.slides.add_slide(blank_layout)
    add_nav(slide, chapter_num)
    add_title_rule(slide, f"{slide_data['id']}  {slide_data['title']}")

    has_fig = slide_data.get("figure") is not None
    callout = slide_data.get("callout")
    bullets = slide_data["bullets"]

    # ブレット数・文の長さはスライドごとに異なる。本文22pt下限のもとで折り返し無く
    # 収まるよう content/*.json 側で1行あたりの字数を実測済み(scratchpad/measure.py)。
    # 固定高さのブロックだとコールアウト箱と重なる事故が起きるため、実測フォント幅で
    # 折り返し行数を見積もり、コールアウトの Y 座標・高さを本文の実高さに追従させる。
    if has_fig:
        left_x, left_w = 0.6, 5.65
        right_x, right_w = 6.6, 6.2
        top_y = 1.90
        b_size, b_line_mult, b_space_after = 22, 1.15, 10
        bullets_h = stack_height_in(["・" + b for b in bullets], b_size, left_w,
                                     line_mult=b_line_mult, space_after_pt=b_space_after)
        add_bullets(slide, left_x, top_y, left_w, bullets_h, bullets,
                    size=b_size, line_spacing=b_line_mult, space_after=b_space_after)
        if callout:
            gap = 0.18
            callout_y = top_y + bullets_h + gap
            inner_w = left_w - 0.36
            callout_h = stack_height_in([callout], 22, inner_w, line_mult=1.1, space_after_pt=0) + 0.24
            callout_h = max(callout_h, 0.55)
            add_callout(slide, left_x, callout_y, left_w, callout_h, callout, size=22)
        img_path = os.path.join(ASSETS, "fig_" + slide_data["id"].replace(".", "_") + ".png")
        if os.path.exists(img_path):
            add_picture_fit(slide, img_path, right_x, 1.85, right_w, 4.85)
    else:
        top_y = 1.95
        bullets_w = 11.0
        b_size, b_line_mult, b_space_after = 22, 1.15, 14
        bullets_h = stack_height_in(["・" + b for b in bullets], b_size, bullets_w,
                                     line_mult=b_line_mult, space_after_pt=b_space_after)
        add_bullets(slide, 0.85, top_y, bullets_w, bullets_h, bullets,
                    size=b_size, line_spacing=b_line_mult, space_after=b_space_after)
        if callout:
            gap = 0.18
            callout_w = 9.8
            callout_y = top_y + bullets_h + gap
            inner_w = callout_w - 0.36
            callout_h = stack_height_in([callout], 22, inner_w, line_mult=1.1, space_after_pt=0) + 0.24
            callout_h = max(callout_h, 0.75)
            add_callout(slide, 1.5, callout_y, callout_w, callout_h, callout, size=22)

    add_footer(slide, slide_data.get("sources", []), page_num)
    add_notes(slide, slide_data.get("notes", ""))
    return slide


# ======================================================================
# スライド生成: まとめ
# ======================================================================
def build_summary(prs, blank_layout, scaffold, page_num):
    slide = prs.slides.add_slide(blank_layout)
    add_nav(slide, None)
    add_title_rule(slide, scaffold["summary"]["title"])
    add_bullets(slide, 0.85, 2.0, 11.2, 4.6, scaffold["summary"]["bullets"],
                size=22, line_spacing=1.25, space_after=18)
    add_footer(slide, [], page_num)
    add_notes(slide, scaffold["summary"].get("notes", ""))
    return slide


# ======================================================================
# スライド生成: 出典・参考
# ======================================================================
def build_references(prs, blank_layout, scaffold, chapters_data, page_num):
    slide = prs.slides.add_slide(blank_layout)
    add_nav(slide, None)
    add_title_rule(slide, "出典・参考")

    all_sources = []
    for n in range(1, 9):
        for s in chapters_data[n]["slides"]:
            all_sources.extend(s.get("sources", []))
    uniq_sorted = sorted(set(all_sources), key=lambda s: s.lower())

    box, tf = new_textbox(slide, 0.6, 1.80, 11.8, 0.45, wrap=True, anchor=MSO_ANCHOR.TOP)
    p = tf.paragraphs[0]
    r = p.add_run()
    set_run(r, text=scaffold["references_note"], size=10, color=MUTED, italic=True)

    # 54件規模のため3カラムに分けて配置し、実測フォントで欠落なく収める。
    n_cols = 3
    col_gap = 0.25
    col_w = (11.8 - col_gap * (n_cols - 1)) / n_cols
    k, m = divmod(len(uniq_sorted), n_cols)
    chunks, start = [], 0
    for i in range(n_cols):
        cnt = k + (1 if i < m else 0)
        chunks.append(uniq_sorted[start:start + cnt])
        start += cnt

    col_top, col_h = 2.35, 4.55
    for i, chunk in enumerate(chunks):
        col_x = 0.6 + i * (col_w + col_gap)
        _, tfc = new_textbox(slide, col_x, col_top, col_w, col_h, wrap=True, anchor=MSO_ANCHOR.TOP)
        first = True
        for src in chunk:
            p = tfc.paragraphs[0] if first else tfc.add_paragraph()
            first = False
            p.line_spacing = 1.05
            p.space_after = Pt(3)
            r = p.add_run()
            set_run(r, text="・" + src, size=10, color=INK)

    add_footer(slide, [], page_num)
    return slide


# ======================================================================
# メイン組み立て
# ======================================================================
def main():
    scaffold, chapters_data = load_data()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank_layout = prs.slide_layouts[6]

    page = 1
    build_cover(prs, blank_layout, scaffold)  # page 1

    page = 2
    build_toc(prs, blank_layout, scaffold, page)  # page 2

    for chapter_num in range(1, 9):
        page += 1
        build_chapter_opener(prs, blank_layout, chapter_num, scaffold, chapters_data, page)
        for s in chapters_data[chapter_num]["slides"]:
            page += 1
            build_content_slide(prs, blank_layout, chapter_num, s, page)

    page += 1
    build_summary(prs, blank_layout, scaffold, page)  # 46

    page += 1
    build_references(prs, blank_layout, scaffold, chapters_data, page)  # 47

    assert len(prs.slides) == TOTAL_SLIDES, f"スライド数不一致: {len(prs.slides)} != {TOTAL_SLIDES}"
    assert page == TOTAL_SLIDES, f"ページ番号不一致: {page} != {TOTAL_SLIDES}"

    prs.save(OUT_PATH)
    print(f"saved: {OUT_PATH} ({len(prs.slides)} slides)")


if __name__ == "__main__":
    main()
