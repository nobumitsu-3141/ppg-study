#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""川副式スライド レイアウト・ルール機械検査（read-only）.

pptx を開いて、スライド作成ルールの機械検出可能な違反を列挙する:
  - 本文フォント < 22pt（出典10.5/ページ番号12/パンくず/図キャプション16 は除外）
  - 図形がスライド枠外にはみ出す（右端/下端）
  - タイトルがスライド左端に密着（左マージン不足）
  - 図・枠・テキストの重なり（bbox 交差。picture/塗り箱が絡むもの）
使い方: python3 slide_lint.py <deck.pptx> [<deck2.pptx> ...]
"""
import sys
from pptx import Presentation
from pptx.util import Cm, Emu
from pptx.enum.shapes import MSO_SHAPE_TYPE

EMU_CM = 360000
SW = int(33.867 * EMU_CM)
SH = int(19.05 * EMU_CM)
TOL = int(0.20 * EMU_CM)          # 枠外はみ出し許容
TITLE_MIN_LEFT = int(1.3 * EMU_CM)  # タイトル左マージン最小
UNDERLINE_Y = int(3.85 * EMU_CM)  # 金下線のおよそのy
BODY_MIN_PT = 22.0
SMALL_OK = {10.5, 12.0, 16.0}     # 出典/ページ番号/パンくず数字/キャプション等の既知小サイズ
OVERLAP_MIN = int((0.6 * EMU_CM) ** 2)  # 重なり最小面積（誤検出抑制）


def _emu(v):
    return int(v) if v is not None else None


def shape_box(sp):
    try:
        l, t, w, h = _emu(sp.left), _emu(sp.top), _emu(sp.width), _emu(sp.height)
    except Exception:
        return None
    if None in (l, t, w, h):
        return None
    return (l, t, w, h)


def has_text(sp):
    return sp.has_text_frame and sp.text_frame.text.strip() != ""


def is_picture(sp):
    return sp.shape_type == MSO_SHAPE_TYPE.PICTURE


def is_filled_box(sp):
    # 角丸/矩形の塗り箱（枠）。テキスト有無問わず視覚要素。
    try:
        return sp.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE
    except Exception:
        return False


def overlap_area(a, b):
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    ix = max(0, min(ax + aw, bx + bw) - max(ax, bx))
    iy = max(0, min(ay + ah, by + bh) - max(ay, by))
    return ix * iy


def contains(a, b):
    """a が b をほぼ内包（箱の中に文字＝重なりでない、を除外するため）。"""
    ax, ay, aw, ah = a
    bx, by, bw, bh = b
    return (ax - TOL <= bx and ay - TOL <= by and
            ax + aw + TOL >= bx + bw and ay + ah + TOL >= by + bh)


def runs_sizes(sp):
    out = []
    if not sp.has_text_frame:
        return out
    for para in sp.text_frame.paragraphs:
        for r in para.runs:
            sz = r.font.size
            if sz is not None:
                out.append((round(sz.pt, 1), r.text[:18]))
    return out


# --- Meiryo安全係数(本番PowerPointのMeiryoはLibreOffice置換のHiraginoより大きい) ---
JP_W = 1.12   # 横幅の割り増し
JP_H = 1.18   # 行高の割り増し
BREADCRUMB_X = int(19.0 * EMU_CM)  # パンくず開始x(概算)


def _units(text):
    u = 0.0
    for ch in text:
        o = ord(ch)
        if ch == " ": u += 0.30
        elif o <= 0x24F: u += 0.56
        elif 0x2080 <= o <= 0x208E: u += 0.56
        else: u += 1.0
    return u


def _text_w_emu(text, pt):
    return _units(text) * (pt / 72 * 2.54) * EMU_CM * JP_W


def find_breadcrumb_x0(slide):
    """右上のパンくずペナント群の左端x(あれば)。無ければNone。"""
    xs = []
    for sp in slide.shapes:
        b = shape_box(sp)
        if not b:
            continue
        l, t, w, h = b
        if is_filled_box(sp) and t < int(1.9 * EMU_CM) and h < int(1.9 * EMU_CM) \
                and w < int(2.2 * EMU_CM) and l > int(17.0 * EMU_CM):
            xs.append(l)
    return min(xs) if xs else None


def est_overflow(sp, box):
    """テキストボックスの推定描画高さ - 枠高さ (cmで正なら超過)。"""
    if not sp.has_text_frame:
        return None
    l, t, w, h = box
    total = 0.0
    for para in sp.text_frame.paragraphs:
        runs = [r for r in para.runs if r.text.strip()]
        if not runs:
            total += 0.2 * EMU_CM
            continue
        pt = max((r.font.size.pt if r.font.size else 22.0) for r in runs)
        txt = "".join(r.text for r in runs)
        line_w = _text_w_emu(txt, pt)
        lines = max(1, -(-int(line_w) // max(1, w)))   # ceil(line_w / box_w)
        ls = para.line_spacing if isinstance(para.line_spacing, (int, float)) else 1.2
        line_h = (pt / 72 * 2.54) * EMU_CM * ls * JP_H
        sa = para.space_after.pt if para.space_after is not None else 6.0
        total += lines * line_h + (sa / 72 * 2.54) * EMU_CM
    return (total - h) / EMU_CM


def lint(path):
    prs = Presentation(path)
    slides = list(prs.slides)
    issues = []
    for si, slide in enumerate(slides):
        if si == 0:
            continue  # 表紙は別扱い
        page = si  # フッター番号相当（表紙除く1始まり）
        bc_x0 = find_breadcrumb_x0(slide)   # パンくず左端x(あれば)
        boxes = []  # (kind, box, shape)
        for sp in slide.shapes:
            box = shape_box(sp)
            is_title = getattr(sp, "is_placeholder", False) and sp.placeholder_format is not None and sp.placeholder_format.idx == 0
            # フォント検査
            for pt, txt in runs_sizes(sp):
                if is_title:
                    continue
                if pt < BODY_MIN_PT and pt not in SMALL_OK and pt > 12.5:
                    issues.append((page, "FONT<22", f"{pt}pt「{txt}」"))
            # 枠外検査
            if box:
                l, t, w, h = box
                if l + w > SW + TOL:
                    issues.append((page, "OOB-右", f"右端 {(l+w)/EMU_CM:.1f}cm>33.9"))
                if t + h > SH + TOL:
                    issues.append((page, "OOB-下", f"下端 {(t+h)/EMU_CM:.1f}cm>19.1"))
            # タイトル: 左マージン / パンくず侵入 / 2行化(Meiryo想定の推定幅で)
            if is_title and box and sp.has_text_frame:
                ttxt = sp.text_frame.text.replace("\n", "")
                tright = box[0] + _text_w_emu(ttxt, 44.0)
                if box[0] < TITLE_MIN_LEFT:
                    issues.append((page, "TITLE左詰め", f"左 {box[0]/EMU_CM:.2f}cm<1.3"))
                if bc_x0 is not None and tright > bc_x0 - int(0.2 * EMU_CM):
                    issues.append((page, "TITLEパンくず侵入", f"推定右端 {tright/EMU_CM:.1f}cm>パンくず{bc_x0/EMU_CM:.1f}cm : {ttxt[:16]}"))
                if tright > box[0] + box[3 - 1]:  # box[2]=width
                    issues.append((page, "TITLE2行化", f"推定幅超過 : {ttxt[:16]}"))
            # 本文あふれ(推定描画高さ vs 枠高さ) + 各行の折返し(1行に収まらない)
            if (not is_title) and box and sp.has_text_frame and sp.text_frame.text.strip():
                cy = box[1] + box[3] / 2
                if int(4.0 * EMU_CM) <= cy <= int(18.2 * EMU_CM):
                    ov = est_overflow(sp, box)
                    if ov is not None and ov > 0.3:
                        issues.append((page, "本文あふれ", f"推定+{ov:.1f}cm超過「{sp.text_frame.text.strip()[:14]}」"))
                    # 各段落がMeiryo想定で枠幅を超える=2行以上に折り返す(単語途中改行の温床)
                    # 本文(≥20pt)のみ対象。参考文献/注釈(10.5pt等)は折返して当然なので除外。
                    for para in sp.text_frame.paragraphs:
                        runs = [r for r in para.runs if r.text.strip()]
                        if not runs:
                            continue
                        pt = max((r.font.size.pt if r.font.size else 22.0) for r in runs)
                        if pt < 20:
                            continue
                        ptxt = "".join(r.text for r in runs)
                        if _text_w_emu(ptxt, pt) > box[2] + int(0.15 * EMU_CM):
                            issues.append((page, "本文折返し", f"「{ptxt[:18]}」"))
            # 重なり候補として収集（本文ゾーンの picture/塗り箱/テキスト箱のみ。
            # ヘッダ帯=タイトル/パンくず/下線、フッタ帯=出典/ページ番号 は除外）
            if box and not is_title and (is_picture(sp) or is_filled_box(sp) or has_text(sp)):
                cy = box[1] + box[3] / 2
                if int(4.2 * EMU_CM) <= cy <= int(17.8 * EMU_CM):
                    kind = "pic" if is_picture(sp) else ("box" if is_filled_box(sp) else "txt")
                    boxes.append((kind, box, sp))
        # 重なり検査：picture か 塗り箱 が絡む対だけ。txt-txt同士・入れ子(箱の中の文字)は除外
        for i in range(len(boxes)):
            for j in range(i + 1, len(boxes)):
                k1, b1, _ = boxes[i]
                k2, b2, _ = boxes[j]
                if k1 == "txt" and k2 == "txt":
                    continue
                if contains(b1, b2) or contains(b2, b1):
                    continue
                a = overlap_area(b1, b2)
                if a > OVERLAP_MIN:
                    issues.append((page, "重なり", f"{k1}×{k2} {a/(EMU_CM**2):.1f}cm²"))
    return len(slides), issues


def main():
    for path in sys.argv[1:]:
        try:
            n, issues = lint(path)
        except Exception as e:
            print(f"\n### {path}\n  ERROR: {e}")
            continue
        cats = {}
        for _, cat, _ in issues:
            cats[cat] = cats.get(cat, 0) + 1
        print(f"\n### {path.split('/')[-1]}  ({n}枚)")
        print(f"  違反 合計 {len(issues)} 件: " + (", ".join(f"{k}={v}" for k, v in sorted(cats.items())) or "なし"))
        bypage = {}
        for pg, cat, det in issues:
            bypage.setdefault(pg, []).append(f"{cat}:{det}")
        for pg in sorted(bypage):
            print(f"    p{pg}: " + " / ".join(bypage[pg][:6]) + (" …" if len(bypage[pg]) > 6 else ""))


if __name__ == "__main__":
    main()
