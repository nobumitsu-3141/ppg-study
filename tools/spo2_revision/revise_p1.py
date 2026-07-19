from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import PP_PLACEHOLDER
path='/tmp/claude-0/-home-user-ppg-study/0e62960e-385d-58f8-942d-ed4e5a333b2d/scratchpad/revised.pptx'
p=Presentation(path)
TITLE=PP_PLACEHOLDER.TITLE  # (1)
changed=0
for s in p.slides:
    for sh in s.shapes:
        if sh.is_placeholder and sh.placeholder_format.type==TITLE and sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                for r in para.runs:
                    if r.font.size is None or abs(r.font.size.pt-44)>0.01:
                        r.font.size=Pt(44); changed+=1
# p36 タイトル短縮（パンくず侵入解消）
for sh in p.slides[35].shapes:
    if sh.is_placeholder and sh.placeholder_format.type==TITLE and sh.has_text_frame:
        runs=sh.text_frame.paragraphs[0].runs
        if runs:
            runs[0].text='6.2  限界への対応'
            for r in runs[1:]:
                r._r.getparent().remove(r._r)
p.save(path)
print('タイトル44pt化した run 数:', changed)
