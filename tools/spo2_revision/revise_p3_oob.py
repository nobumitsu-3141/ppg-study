import sys
sys.path.insert(0,'/root/.claude/skills/slide-format/scripts')
import slide_lint as L
from pptx import Presentation
path=sys.argv[1]
p=Presentation(path)
EMU=L.EMU_CM; SW=L.SW; SH=L.SH; TOL=L.TOL
M=int(0.3*EMU)
fixed=[]
for si,slide in enumerate(p.slides):
    if si==0: continue
    for sp in slide.shapes:
        b=L.shape_box(sp)
        if not b: continue
        l,t,w,h=b; ch=False
        if l+w>SW+TOL:
            newl=SW-w-M
            if newl>=M: sp.left=int(newl); ch=True
            else:  # 広すぎる→幅縮小
                sp.width=int(SW-l-M); ch=True
        if t+h>SH+TOL:
            newt=SH-h-M
            if newt>=M: sp.top=int(newt); ch=True
            else:
                sp.height=int(SH-t-M); ch=True
        if ch: fixed.append((si+1, (sp.text_frame.text.strip()[:16] if sp.has_text_frame else '<図>')))
p.save(path)
print('OOB是正:', len(fixed))
for f in fixed: print('   p%d [%s]'%f)
