import sys
sys.path.insert(0,'/root/.claude/skills/slide-format/scripts')
import slide_lint as L
from pptx import Presentation
from pptx.util import Pt
path=f'{sys.argv[1]}'
p=Presentation(path)
EMU=L.EMU_CM
def is_title(sp):
    return getattr(sp,'is_placeholder',False) and sp.placeholder_format is not None and sp.placeholder_format.idx==0
def flagged_pt(pt):
    pt=round(pt,1); return (12.5 < pt < 22.0) and pt!=16.0
fixed=0; deferred=[]
for si,slide in enumerate(p.slides):
    if si==0: continue
    for sp in slide.shapes:
        if not sp.has_text_frame or is_title(sp): continue
        box=L.shape_box(sp)
        if not box: continue
        runs_flagged=[(para,r,r.font.size) for para in sp.text_frame.paragraphs for r in para.runs
                      if r.font.size and flagged_pt(r.font.size.pt)]
        if not runs_flagged: continue
        # 試しに22ptへ
        for _,r,_ in runs_flagged: r.font.size=Pt(22)
        # lintと同じ判定: あふれ / 折返し
        cy=box[1]+box[3]/2
        in_zone= int(4.0*EMU)<=cy<=int(18.2*EMU)
        ov=L.est_overflow(sp,box)
        overflow_bad= in_zone and (ov is not None and ov>0.3)
        wrap_bad=False
        for para in sp.text_frame.paragraphs:
            rs=[r for r in para.runs if r.text.strip()]
            if not rs: continue
            pt=max((r.font.size.pt if r.font.size else 22.0) for r in rs)
            if pt<20: continue
            ptxt=''.join(r.text for r in rs)
            if L._text_w_emu(ptxt,pt) > box[2]+int(0.15*EMU): wrap_bad=True
        if overflow_bad or wrap_bad:
            for _,r,sz in runs_flagged: r.font.size=sz   # revert
            deferred.append((si+1, sp.text_frame.text.strip()[:18], 'あふれ' if overflow_bad else '折返し'))
        else:
            fixed+=1
p.save(path)
print(f'安全に22pt化した枠: {fixed}')
print(f'保留(枠拡大/分割が必要): {len(deferred)}')
for d in deferred[:40]: print('   保留 p%d [%s] %s'%d)
