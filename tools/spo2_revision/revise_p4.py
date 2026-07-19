import sys
sys.path.insert(0,'/root/.claude/skills/slide-format/scripts')
import slide_lint as L
from pptx import Presentation
from pptx.util import Pt, Emu
EMU=L.EMU_CM; SW=L.SW; SH=L.SH
GAP=int(0.3*EMU); TOLW=int(0.15*EMU)
path=sys.argv[1]
p=Presentation(path)
def is_title(sp):
    return getattr(sp,'is_placeholder',False) and sp.placeholder_format is not None and sp.placeholder_format.idx==0
def flagged(pt):
    pt=round(pt,1); return 12.5<pt<22.0 and pt!=16.0
def box_viol(sp,box):
    v=0
    for para in sp.text_frame.paragraphs:
        for r in para.runs:
            if r.font.size:
                pt=round(r.font.size.pt,1)
                if pt<22 and pt not in L.SMALL_OK and pt>12.5: v+=1
    cy=box[1]+box[3]/2
    if int(4.0*EMU)<=cy<=int(18.2*EMU):
        ov=L.est_overflow(sp,box)
        if ov is not None and ov>0.3: v+=1
        for para in sp.text_frame.paragraphs:
            rs=[r for r in para.runs if r.text.strip()]
            if not rs: continue
            pt=max((r.font.size.pt if r.font.size else 22.0) for r in rs)
            if pt<20: continue
            ptxt=''.join(r.text for r in rs)
            if L._text_w_emu(ptxt,pt)>box[2]+TOLW: v+=1
    return v
def overlaps(box, others):
    n=0
    for b in others:
        if L.contains(box,b) or L.contains(b,box): continue
        if L.overlap_area(box,b)>L.OVERLAP_MIN: n+=1
    return n
kept=0
for si,slide in enumerate(p.slides):
    if si==0: continue
    # 各スライドの他図形(pic/box/txt)bbox
    allshapes=[(sp,L.shape_box(sp)) for sp in slide.shapes]
    for sp in slide.shapes:
        if not sp.has_text_frame or is_title(sp): continue
        box=L.shape_box(sp)
        if not box or not sp.text_frame.text.strip(): continue
        others=[b for (o,b) in allshapes if o is not sp and b and (L.is_picture(o) or L.is_filled_box(o) or (o.has_text_frame and o.text_frame.text.strip()))]
        before=box_viol(sp,box); before_ol=overlaps(box,others)
        if before==0: continue
        # snapshot
        snap=[(r,r.font.size) for para in sp.text_frame.paragraphs for r in para.runs]
        gl,gt,gw,gh=sp.left,sp.top,sp.width,sp.height
        # 1) flagged fonts -> 22
        for para in sp.text_frame.paragraphs:
            for r in para.runs:
                if r.font.size and flagged(r.font.size.pt): r.font.size=Pt(22)
        # 2) 必要幅/高さを算出
        l,t,w,h=box
        need_w=w
        for para in sp.text_frame.paragraphs:
            rs=[r for r in para.runs if r.text.strip()]
            if not rs: continue
            pt=max((r.font.size.pt if r.font.size else 22.0) for r in rs)
            if pt<20: continue
            need_w=max(need_w, int(L._text_w_emu(''.join(r.text for r in rs),pt))+int(0.2*EMU))
        ovcm=L.est_overflow(sp,box)
        need_h=h + (int(ovcm*EMU)+int(0.2*EMU) if (ovcm and ovcm>0) else 0)
        # 3) 空きへ拡張
        if need_w>w:
            cand=[o[0] for o in others if o[0]>=l+w-TOLW and not (o[1]>=t+h or o[1]+o[3]<=t)]
            maxright=min([SW-GAP]+([min(cand)-GAP] if cand else []))
            neww=min(need_w, maxright-l)
            if neww>w+10: sp.width=int(neww)
        if need_h>h:
            cand=[o[1] for o in others if o[1]>=t+h-TOLW and not (o[0]>=l+w or o[0]+o[2]<=l)]
            maxbot=min([SH-GAP]+([min(cand)-GAP] if cand else []))
            newh=min(need_h, maxbot-t)
            if newh>h+10: sp.height=int(newh)
        # 4) 評価
        nb=L.shape_box(sp)
        after=box_viol(sp,nb); after_ol=overlaps(nb,others)
        if after<before and after_ol<=before_ol:
            kept+=1
        else:
            for r,sz in snap: r.font.size=sz
            sp.left,sp.top,sp.width,sp.height=gl,gt,gw,gh
p.save(path)
print('採用した枠(改善):', kept)
