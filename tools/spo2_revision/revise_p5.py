import sys, copy
sys.path.insert(0,'/root/.claude/skills/slide-format/scripts')
import slide_lint as L
from pptx import Presentation
from pptx.util import Pt
from pptx.oxml.ns import qn
EMU=L.EMU_CM; SW=L.SW; SH=L.SH; GAP=int(0.3*EMU); TOLW=int(0.15*EMU)
path=sys.argv[1]; p=Presentation(path)
BEFORE={'　',' ','、','／','。','：','）',')'}   # この文字の後ろで改行可
AFTER={'＋','→'}                                  # この文字の前で改行可
def is_title(sp):
    return getattr(sp,'is_placeholder',False) and sp.placeholder_format is not None and sp.placeholder_format.idx==0
def wrap_paras(sp,box):
    out=[]
    for para in sp.text_frame.paragraphs:
        rs=[r for r in para.runs if r.text.strip()]
        if not rs: continue
        pt=max((r.font.size.pt if r.font.size else 22.0) for r in rs)
        if pt<20: continue
        txt=''.join(r.text for r in rs)
        if L._text_w_emu(txt,pt)>box[2]+TOLW: out.append((para,pt,txt))
    return out
def box_viol(sp,box):
    v=0
    for para in sp.text_frame.paragraphs:
        for r in para.runs:
            if r.font.size:
                pt=round(r.font.size.pt,1)
                if pt<22 and pt not in L.SMALL_OK and pt>12.5: v+=1
    cy=box[1]+box[3]/2
    if int(4.0*EMU)<=cy<=int(18.2*EMU):
        ov=L.est_overflow(sp,box)
        if ov is not None and ov>0.3: v+=1
        for para in sp.text_frame.paragraphs:
            rs=[r for r in para.runs if r.text.strip()]
            if not rs: continue
            pt=max((r.font.size.pt if r.font.size else 22.0) for r in rs)
            if pt<20: continue
            if L._text_w_emu(''.join(r.text for r in rs),pt)>box[2]+TOLW: v+=1
    return v
def overlaps(box,others):
    n=0
    for b in others:
        if L.contains(box,b) or L.contains(b,box): continue
        if L.overlap_area(box,b)>L.OVERLAP_MIN: n+=1
    return n
def best_split(txt,pt,boxw):
    cands=[i for i in range(1,len(txt)) if (txt[i-1] in BEFORE) or (txt[i] in AFTER)]
    best=None
    for i in cands:
        a,b=txt[:i].strip(),txt[i:].strip()
        if not a or not b: continue
        wa,wb=L._text_w_emu(a,pt),L._text_w_emu(b,pt)
        if wa<=boxw+TOLW and wb<=boxw+TOLW:
            score=abs(wa-wb)
            if best is None or score<best[0]: best=(score,i)
    return best[1] if best else None
def split_para(para,k):
    # para のテキストを k 文字目で2段落に（run境界/run内どちらも対応, 書式維持）
    runs=para.runs; acc=0; new_p=copy.deepcopy(para._p)
    # new_p の runs を全削除、pPr は残す
    for r in new_p.findall(qn('a:r')): new_p.remove(r)
    # 元 para: k以降を new_p へ
    moved=[]; cut_done=False
    idx=0
    for r in list(runs):
        tlen=len(r.text)
        if acc+tlen<=k:
            acc+=tlen; idx+=1; continue
        if acc>=k:
            moved.append(r); 
        else:
            off=k-acc
            # run内で分割
            r2=copy.deepcopy(r._r)
            # r: 前半, r2: 後半
            r._r.find(qn('a:t')).text=r.text[:off]
            r2.find(qn('a:t')).text=r.text[off:]
            new_p.append(r2)
        acc+=tlen
    for r in moved:
        new_p.append(copy.deepcopy(r._r)); r._r.getparent().remove(r._r)
    para._p.addnext(new_p)
def grow_h(sp,slide,others):
    b=L.shape_box(sp); l,t,w,h=b
    ov=L.est_overflow(sp,b)
    if not ov or ov<=0.3: return
    need=h+int(ov*EMU)+int(0.2*EMU)
    cand=[o[1] for o in others if o[1]>=t+h-TOLW and not (o[0]>=l+w or o[0]+o[2]<=l)]
    maxbot=min([SH-GAP]+([min(cand)-GAP] if cand else []))
    newh=min(need,maxbot-t)
    if newh>h+10: sp.height=int(newh)
kept=0
for si,slide in enumerate(p.slides):
    if si==0: continue
    allsh=[(o,L.shape_box(o)) for o in slide.shapes]
    for sp in slide.shapes:
        if not sp.has_text_frame or is_title(sp): continue
        box=L.shape_box(sp)
        if not box or not sp.text_frame.text.strip(): continue
        others=[b for (o,b) in allsh if o is not sp and b and (L.is_picture(o) or L.is_filled_box(o) or (o.has_text_frame and o.text_frame.text.strip()))]
        wps=wrap_paras(sp,box)
        if not wps: continue
        before=box_viol(sp,box); before_ol=overlaps(box,others)
        xml_snap=copy.deepcopy(sp.text_frame._txBody)
        gl,gt,gw,gh=sp.left,sp.top,sp.width,sp.height
        did=False
        for para,pt,txt in wps:
            k=best_split(txt,pt,box[2])
            if k: split_para(para,k); did=True
        if did: grow_h(sp,slide,others)
        nb=L.shape_box(sp)
        after=box_viol(sp,nb); after_ol=overlaps(nb,others)
        if did and after<before and after_ol<=before_ol:
            kept+=1
        else:
            sp.text_frame._txBody.getparent().replace(sp.text_frame._txBody,xml_snap)
            sp.left,sp.top,sp.width,sp.height=gl,gt,gw,gh
p.save(path)
print('折返し改善した枠:',kept)
