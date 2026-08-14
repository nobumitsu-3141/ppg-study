# -*- coding: utf-8 -*-
"""R2 stage4: 参考文献を「スライドでの登場順」に振り直し、参考文献ページを作り直す"""
import re
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
from deckkit import (add_text, set_title, set_source, set_pageno, clone_slide,
                     slide_index, drop_shapes, INK, GREY)
from r2_common import rewrite_nav, title_text
from refs import REFS, SHORT

prs = Presentation('r2_s3.pptx')
SL = lambda n: list(prs.slides)[n - 1]

NUM_TOKEN = re.compile(r'(\d{1,2})\.\s*(?=[A-Za-zÀ-ÿ])')
BARE_LIST = re.compile(r'^\s*\d{1,2}(\s*,\s*\d{1,2})*\s*$')


def para_text(para):
    return ''.join(r.text for r in para.runs)


def set_para(para, text):
    runs = para.runs
    if not runs:
        return
    runs[0].text = text
    for r in runs[1:]:
        r.text = ''


# ═══════════ 0) 表記ゆれの正規化（番号なしの Ahmed 2005 など）
for i, s in enumerate(prs.slides, 1):
    for sh in s.shapes:
        if sh.name != 'Source' or not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            t = orig = para_text(para)
            if not t.strip():
                continue
            if title_text(prs.slides[i - 1]) == '6.3  DN-less 信号対策' \
                    and '48.' not in t and '47.' in t:
                t = t.replace('47. Attivissimo 2023',
                              '47. Attivissimo 2023　／　48. Domínguez-Hernández 2026')
            t2 = (t.replace('、　Ahmed 2005', '　／　53. Ahmed 2005')
                   .replace('、Ahmed 2005', '　／　53. Ahmed 2005')
                   .replace('36. Wu 2010　38. Pal 2024', '36. Wu 2010　／　38. Pal 2024'))
            if t2 != orig:
                set_para(para, t2)

# ═══════════ 1) 引用箇所の収集（登場順）
targets = []          # (slide_idx, shape, para, kind)
order = []


def note_num(n):
    if n in REFS and n not in order:
        order.append(n)


REF_PAGES = {i for i, s in enumerate(prs.slides, 1) if title_text(s).startswith('参考文献')}
# 章扉とメニューは目次テキストが「7. PDA…」の形なので走査から外す
DIVIDERS = {i for i, s in enumerate(prs.slides, 1)
            if title_text(s) == 'メニュー' or re.match(r'^\d+\.\s', title_text(s))}
REF_PAGES |= DIVIDERS

for i, s in enumerate(prs.slides, 1):
    if i in REF_PAGES:
        continue
    # Source を先に
    for sh in s.shapes:
        if sh.name == 'Source' and sh.has_text_frame:
            for para in sh.text_frame.paragraphs:
                t = para_text(para)
                if not t.strip():
                    continue
                nums = [int(m.group(1)) for m in NUM_TOKEN.finditer(t)]
                if nums:
                    targets.append((i, para, 'source'))
                    for n in nums:
                        note_num(n)
    # 本文中の番号（表の出典列・文献チップ）
    for sh in s.shapes:
        if sh.name in ('Source', 'PageNo', 'Title 1', 'タイトル 1') or not sh.has_text_frame:
            continue
        if 'Off-page Connector' in sh.name:      # 章ナビ帯の番号は文献番号ではない
            continue
        for para in sh.text_frame.paragraphs:
            t = para_text(para)
            if not t.strip():
                continue
            if BARE_LIST.match(t):
                nums = [int(x) for x in re.findall(r'\d{1,2}', t)]
                # 出典列だけを対象にする（右端に置いた列、または複数番号の並び）
                is_cite_col = (sh.left is not None and sh.left > 10.4 * 914400)
                if not (is_cite_col or ',' in t):
                    continue
                if not all(n in REFS for n in nums):
                    continue
                targets.append((i, para, 'bare'))
                for n in nums:
                    note_num(n)
            elif NUM_TOKEN.match(t.lstrip()):
                targets.append((i, para, 'chip'))
                m = NUM_TOKEN.match(t.lstrip())
                note_num(int(m.group(1)))

NEW = {old: k for k, old in enumerate(order, start=1)}
print('使用文献:', len(order), '件 ／ 未使用で削除:',
      sorted(SHORT[n] for n in set(REFS) - set(order)))
FIRST = {}
for i, para, kind in targets:
    for m in NUM_TOKEN.finditer(para_text(para)):
        FIRST.setdefault(int(m.group(1)), i)
    if kind == 'bare':
        for x in re.findall(r'\d{1,2}', para_text(para)):
            FIRST.setdefault(int(x), i)

# ═══════════ 2) 番号の書き換え
for i, para, kind in targets:
    t = para_text(para)
    if kind == 'bare':
        nums = [int(x) for x in re.findall(r'\d{1,2}', t)]
        set_para(para, ', '.join(str(NEW[n]) for n in nums if n in NEW))
    else:
        def rep(m):
            n = int(m.group(1))
            return f'{NEW[n]}. ' if n in NEW else m.group(0)
        set_para(para, NUM_TOKEN.sub(rep, t))

# 出典行を新しい番号の昇順に並べ替える
TOKEN_SPLIT = re.compile(r'\s*(?:／|/)\s*')
for i, s_ in enumerate(prs.slides, 1):
    for sh in s_.shapes:
        if sh.name != 'Source' or not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            t = para_text(para)
            if not t.strip() or '／' not in t:
                continue
            note = ''
            m = re.search(r'（[^）]*）\s*$', t)
            if m:
                note, t = m.group(0), t[:m.start()]
            parts = [x.strip() for x in TOKEN_SPLIT.split(t) if x.strip()]
            keyed = []
            for x in parts:
                mm = re.match(r'^(\d{1,2})\.', x)
                keyed.append((int(mm.group(1)) if mm else 999, x))
            keyed.sort(key=lambda kv: kv[0])
            set_para(para, '　／　'.join(x for _, x in keyed) + note)

# ノート内の「N. 著者」表記も同じ対応表で置き換える
def renum_notes(txt):
    def rep(m):
        n = int(m.group(1))
        return f'{NEW[n]}. ' if n in NEW else ''
    return NUM_TOKEN.sub(rep, txt)


for s in prs.slides:
    if s.has_notes_slide:
        tf = s.notes_slide.notes_text_frame
        nt = renum_notes(tf.text)
        if nt != tf.text:
            tf.text = nt

# ═══════════ 3) 参考文献ページの作り直し
ref_idx = [i for i, s in enumerate(prs.slides, 1) if title_text(s).startswith('参考文献')]
base = SL(ref_idx[0])
PER = 14
pages = [order[k:k + PER] for k in range(0, len(order), PER)]

while len(ref_idx) < len(pages):
    ns = clone_slide(prs, base, insert_after_idx=ref_idx[-1] - 1 + 1)
    ref_idx.append(slide_index(prs, ns) + 1)
while len(ref_idx) > len(pages):
    i = ref_idx.pop()
    el = list(prs.slides._sldIdLst)[i - 1]
    rid = el.get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
    prs.slides._sldIdLst.remove(el)
    prs.part.drop_rel(rid)

for pi, (i, group) in enumerate(zip(ref_idx, pages)):
    slide = SL(i)
    drop_shapes(slide, lambda sh: sh.has_text_frame and sh.top is not None
                and 1.6 * 914400 < sh.top < 6.9 * 914400)
    set_title(slide, '参考文献' if pi == 0 else f'参考文献（{pi + 1}）')
    rewrite_nav(slide, None)
    set_source(slide, None)
    half = (len(group) + 1) // 2
    for ci, chunk in enumerate((group[:half], group[half:])):
        if not chunk:
            continue
        add_text(slide, 0.60 + ci * 6.20, 1.90, 5.95, 5.10,
                 [f'{NEW[n]}. {REFS[n]}' for n in chunk],
                 size=10.5, color=INK, space_after=7, line_spacing=1.15)

# ═══════════ 4) ページ番号の再採番（同じタイトルの連続＝1論理ページ）
vis = [(i, s) for i, s in enumerate(prs.slides, 1) if s.element.get('show') != '0']
logical, prev = [], None
for i, s in vis:
    key = title_text(s) or f'(fig{i})'
    if prev is not None and key == prev and not key.startswith('(fig'):
        logical[-1].append(s)
    else:
        logical.append([s])
    prev = key
total = len(logical)
for k, group in enumerate(logical, start=1):
    for s in group:
        set_pageno(s, None if k == 1 else f'{k}/{total}')
for i, s in enumerate(prs.slides, 1):
    if s.element.get('show') == '0':
        set_pageno(s, None)

prs.save('r2_s4.pptx')
print('r2 stage4 ok:', len(prs.slides._sldIdLst), 'slides ／ 論理ページ', total)
print('新旧対応:', {NEW[o]: SHORT[o] for o in order})
