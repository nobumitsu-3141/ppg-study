# -*- coding: utf-8 -*-
"""R2 stage3: 5・66 の図解版追加、60 の要約補完と改善案スライド"""
from pptx import Presentation
from pptx.util import Emu
from pptx.enum.text import PP_ALIGN
from deckkit import (add_text, add_box, add_arrow, add_line, add_freeform,
                     set_notes, set_title, set_source, set_pageno, clone_slide,
                     strip_content, slide_index, INK, GOLD, BLUE, VERM, TEAL,
                     RED, GREY, LGREY)
from r2_common import rewrite_nav, title_text
from refs import cite
from newpages import ppg_curve, WAVE_YOUNG, WAVE_MID, WAVE_OLD, WAVE_FLAT

prs = Presentation('r2_s2a.pptx')
PALE = "F2F2F2"


def idx_of_title(t, nth=0):
    return [i for i, s in enumerate(prs.slides, 1) if title_text(s) == t][nth]


def S(n):
    return list(prs.slides)[n - 1]


tpl = S(idx_of_title('5.3  SDPPG の主な知見'))


def new_after(anchor, title, ch, source=None, notes=None):
    s = clone_slide(prs, tpl, insert_after_idx=slide_index(prs, anchor))
    strip_content(s)
    set_title(s, title)
    rewrite_nav(s, ch)
    set_source(s, source)
    set_pageno(s, None)
    if notes:
        set_notes(s, notes)
    return s


def replace_text(slide, pairs):
    for sh in slide.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            runs = para.runs
            if not runs:
                continue
            whole = ''.join(r.text for r in runs)
            new = whole
            for a, b in pairs:
                new = new.replace(a, b)
            if new != whole:
                runs[0].text = new
                for r in runs[1:]:
                    r.text = ''


# ══════════════════════════ 5 の直後：SaO₂ と SpO₂ の関係（図解版）
anchor = S(5)
s = new_after(anchor, '1.1  SpO₂ と SaO₂', 1, cite(25, 26, 53))

add_text(s, 0.55, 1.82, 5.85, 0.45, '動脈血のヘモグロビン', size=24, bold=True, color=INK)
# O2Hb / HHb のブロック図
for i in range(10):
    col = VERM if i < 9 else BLUE
    add_box(s, 0.60 + i * 0.56, 2.32, 0.48, 0.95, fill=col, line=None, radius=0.12)
add_text(s, 0.60, 3.40, 5.7, 0.95,
         [[('■', {'color': VERM}), (' 酸素と結合した  O₂Hb', {})],
          [('■', {'color': BLUE}), (' 結合していない  HHb', {})]],
         size=22, color=INK, space_after=3)

add_box(s, 0.55, 4.42, 5.85, 1.20, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.70, 4.64, 5.55, 0.80,
         [[('SaO₂', {'bold': True, 'color': GOLD, 'size': 26}),
           (' ＝ O₂Hb の割合', {'bold': True, 'size': 25})],
          [('この例なら 90%', {'size': 23})]],
         size=25, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=3)

# 右：測り方の2ルート
add_text(s, 6.75, 1.82, 6.10, 0.45, '測り方は 2 通り', size=24, bold=True, color=INK)

add_box(s, 6.75, 2.28, 6.10, 1.62, fill=PALE, line=BLUE, line_w=2.25)
add_text(s, 6.95, 2.44, 5.70, 1.30,
         [[('採血 → CO オキシメータ', {'size': 23})],
          [('SaO₂', {'bold': True, 'color': BLUE, 'size': 27})],
          [('酸素化評価のゴールドスタンダード', {'size': 22})]],
         size=23, color=INK, align=PP_ALIGN.CENTER, space_after=4)

add_box(s, 6.75, 4.06, 6.10, 1.62, fill=PALE, line=VERM, line_w=2.25)
add_text(s, 6.95, 4.22, 5.70, 1.30,
         [[('指に光を当てる（非侵襲）', {'size': 23})],
          [('SpO₂', {'bold': True, 'color': VERM, 'size': 27})],
          [('SaO₂ を経皮的に推定した値', {'size': 22})]],
         size=23, color=INK, align=PP_ALIGN.CENTER, space_after=4)

add_box(s, 0.55, 5.88, 12.30, 1.00, fill=PALE, line=RED, line_w=2.25)
add_text(s, 0.75, 6.06, 11.90, 0.62,
         'パルスオキシメーターは、どうやって SaO₂ を推定しているのか？',
         size=26, bold=True, color=INK, align=PP_ALIGN.CENTER)

set_notes(s, """1.1 SpO₂ と SaO₂（図解版）

・SaO₂（動脈血酸素飽和度）＝ 動脈血中で酸素と結合しているヘモグロビンの割合。
　SO₂ は「酸化ヘモグロビン O₂Hb と、O₂Hb ＋ 脱酸素化ヘモグロビン RHb の和との比」として
　100% × O₂Hb /（O₂Hb + RHb）と定義される（Ahmed 2005）。
　いわゆる functional saturation で、COHb・MetHb は分母に含めない。
・SaO₂ は動脈採血＋CO オキシメータで実測する。酸素化評価のゴールドスタンダード。
・SpO₂ は指に光を当てて非侵襲に SaO₂ を推定した値。実測値ではなく推定値である。

Ahmed A, et al. Hemoglobin oxygen saturation discrepancy using various methods in
patients with sickle cell vaso-occlusive painful crisis. Eur J Haematol 2005;74:309-14.
（鎌状赤血球症の血管閉塞発作という条件下で、パルスオキシメトリ SpO₂、CO オキシメトリの
　SO₂・FO₂Hb、酸素解離曲線からの計算値 SaO₂ の一致度を比較した研究。
　各指標の定義の違いが乖離を生むことを示している。）

→ 次ページ以降で「どうやって推定しているのか」を追う。

出典：25. Aoyagi 2003 ／ 26. Chan 2013 ／ Ahmed 2005""")

# ══════════════════════════ 60：Hashimoto より下の要約を補完
s60 = S(idx_of_title('5.3  SDPPG の主な知見'))
replace_text(s60, [('Hashimoto 2006', 'Hashimoto 2002'),
                   ('J - SHIPP', '（J-SHIPP 研究）')])
for sh in s60.shapes:
    if not sh.has_text_frame:
        continue
    t = sh.text_frame.text.strip()
    if t == 'PWVとSDPPGは別物':
        from deckkit import add_text as _at
        sh.text_frame.paragraphs[0].runs[0].text = 'PWV と SDPPG の相関は弱く、中心部と末梢部で別の情報'
        p2 = sh.text_frame.add_paragraph()
        r = p2.add_run()
        r.text = '心拍数は d/a・AGI と負、b/a と正に相関（治療中高血圧 294 名）'
        from pptx.util import Pt as _Pt
        from pptx.dml.color import RGBColor as _RGB
        from deckkit import _force_ea, FONT
        r.font.size = _Pt(22)
        r.font.name = FONT
        r.font.color.rgb = _RGB.from_string(INK)
        _force_ea(r, FONT)
    if t.startswith('臓器障害との関連は弱い'):
        runs = sh.text_frame.paragraphs[0].runs
        runs[0].text = '頸動脈 IMT など臓器障害との関連は弱いか有意でない'
        for r in runs[1:]:
            r.text = ''
        if len(sh.text_frame.paragraphs) > 1:
            rr = sh.text_frame.paragraphs[1].runs
            if rr:
                rr[0].text = 'b/a・d/a・AGI の強い規定因子は 年齢・性・収縮期血圧・心拍数'
                for r in rr[1:]:
                    r.text = ''

# ══════════════════════════ 60 の直後：改善案スライド
s = new_after(s60, '5.3  SDPPG は何を測れたか', 5, cite(16, 18, 19, 30, 32))

cols = [
    ('急性のトーヌス', BLUE,
     [('16. Takazawa 1998', '薬で d/a が動く'),
      ('', ['AT II −0.62', 'NTG −0.25'])], '○'),
    ('慢性の加齢', TEAL,
     [('16. Takazawa 1998', 'AGI と年齢 r=0.80'),
      ('30. Otsuka 2006', ['Framingham risk', 'score と相関'])], '○'),
    ('臓器障害・動脈硬化', VERM,
     [('18. Bortolotto 2000', ['規定因子は', 'PWV が優位']),
      ('32. Tabara 2016', '関連は弱い')], '×'),
]
for i, (q, col, items, mark) in enumerate(cols):
    x = 0.50 + i * 4.20
    add_box(s, x, 1.88, 3.90, 4.05, fill=PALE, line=col, line_w=2.25)
    add_text(s, x + 0.12, 2.06, 3.66, 0.5, q, size=24, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.12, 2.66, 3.66, 0.5, mark, size=30, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    y = 3.34
    for name, what in items:
        if name:
            add_text(s, x + 0.15, y, 3.60, 0.42, name, size=22, bold=True, color=INK)
            y += 0.44
        lines = what if isinstance(what, list) else [what]
        add_text(s, x + 0.15, y, 3.60, 0.42 * len(lines), lines, size=22,
                 color=GREY, space_after=1)
        y += 0.16 + 0.40 * len(lines)

add_box(s, 0.50, 6.10, 12.30, 0.85, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.70, 6.28, 11.90, 0.5,
         '19. Hashimoto 2002 ： SDPPG は PWV の代用ではなく別の情報',
         size=23, bold=True, color=INK, align=PP_ALIGN.CENTER)

set_notes(s, """5.3 SDPPG は何を測れたか（前ページの改善案）

前ページは「文献名 → 知見」を 5 行並べる形だった。何が言えて何が言えないのかを
つかみやすくするため、**問いの軸で 3 列に組み替えた**のがこのページである。

【急性のトーヌス軸：○】
・16. Takazawa 1998：上行大動脈圧を同時記録した 39 名。d/a はアンジオテンシン II で
　−0.40±0.13 → −0.62±0.19、ニトログリセリンで −0.25±0.12（いずれも P<0.001）。
　血管収縮で d/a はより負に、血管拡張で 0 に近づく。＝薬理学的に鋭敏に動く。

【慢性の加齢軸：○】
・16. Takazawa 1998：健診受診者 600 名。加齢指数 AGI は y = 0.023x − 1.515（r=0.80, P<0.001）。
・30. Otsuka 2006：一般集団 211 名。SDPTG 指標は Framingham risk score と相関
　（b/a r=0.43(男)/0.54(女)、d/a r=−0.38/−0.58）。高リスク者の判別は感度・特異度とも 0.85。

【臓器障害・動脈硬化軸：×】
・18. Bortolotto 2000：高血圧 524 例。動脈硬化性病変あり群で AGI −0.093 vs なし −0.271（P<.001）
　だが、ロジスティック回帰では PWV が有意な独立規定因子で SDPTG-AI の寄与は弱かった。
・32. Tabara 2016（J-SHIPP）：b/a・d/a・AGI の強い規定因子は年齢・性・収縮期血圧・心拍数で、
　頸動脈 IMT など臓器障害との関連は弱いか有意でない（b/a で β=0.069, P=0.002）。

【位置づけ】
・19. Hashimoto 2002：治療中高血圧 294 名。PWV と SDPTG 指標は年齢・血圧という共通因子を
　持つが相互の相関は弱く、中心部と末梢部について異なる情報を与える。
　＝ SDPPG は PWV の代用ではない。別の軸として扱う。

【まとめ】
SDPPG は「いま血管トーヌスがどう動いたか」と「年齢相応か」には答えられるが、
「臓器障害があるか」には単独では答えられない。""")

# ══════════════════════════ 66 の直後：切痕が消えるとどうなるか（図解版）
s66 = S(idx_of_title('6.2  極端な動脈硬化'))
s = new_after(s66, '6.2  切痕が消えると', 6, cite(6, 22, 36, 38))

waves = [('若年・柔らかい', WAVE_YOUNG, '切痕が見える', BLUE, '○'),
         ('中間', WAVE_MID, '切痕が浅い', TEAL, '△'),
         ('高齢・硬い', WAVE_OLD, '肩だけになる', VERM, '△'),
         ('極端な動脈硬化', WAVE_FLAT, '目印が消える', RED, '×')]
for i, (label, comps, note, col, mark) in enumerate(waves):
    x = 0.45 + i * 3.15
    add_text(s, x, 1.85, 2.95, 0.42, label, size=23, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_freeform(s, ppg_curve(x + 0.10, 2.32, 2.75, 1.30, comps), color=col, width=2.5)
    add_line(s, x + 0.10, 3.64, x + 2.85, 3.64, color=LGREY, width=1.25)
    add_text(s, x, 3.74, 2.95, 0.42, note, size=22, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, x, 4.20, 2.95, 0.45, mark, size=27, bold=True, color=col,
             align=PP_ALIGN.CENTER)

add_text(s, 0.45, 4.82, 12.43, 0.4, 'SI・RI・SDPPG が必要とする「点」を拾えるか',
         size=22, bold=True, color=GREY, align=PP_ALIGN.CENTER)

add_box(s, 0.45, 5.28, 6.05, 1.32, fill=PALE, line=RED, line_w=2.25)
add_text(s, 0.62, 5.48, 5.70, 0.95,
         [[('DN-less signal', {'bold': True, 'color': RED, 'size': 25})],
          [('切痕が見えない PPG 波形のこと', {'size': 22})]],
         size=24, bold=True, color=INK, space_after=3)

add_box(s, 6.83, 5.28, 6.05, 1.32, fill=PALE, line=BLUE, line_w=2.25)
add_text(s, 7.00, 5.48, 5.70, 0.95,
         [[('全体の ', {}), ('14%', {'bold': True, 'color': BLUE, 'size': 27}),
           ('（25,286 / 169,787）', {'size': 22})],
          [('残り 86% は切痕が見える', {'size': 22})]],
         size=24, bold=True, color=INK, space_after=3)

add_text(s, 0.45, 6.72, 12.43, 0.4,
         '→ 問題は「全員見えない」ことではなく、見えない人だけが解析から落ちること',
         size=23, bold=True, color=INK)

set_notes(s, """6.2 切痕が消えるとどうなるか（前ページの図解版）

・動脈硬化が進むと反射波が早く戻り、切痕は「明瞭 → 浅い → 下降脚の肩 → 消失」と変化する。
・SI（＝身長 ÷ ΔT）も RI（＝P2/P1）も「P2 という点」の時刻と高さを必要とするため、
　点が無い波形では原理的に計算できない。SDPPG（二次微分）も、下降脚が平滑になるほど
　極値が不安定になる。
・切痕が見えない PPG 信号を DN-less signal と呼ぶ（38. Pal 2024）。
・Cunningham 2023（UK Biobank, n=169,787）では切痕が欠如していたのは 25,286 名（14%）。
　裏返せば 86% では切痕が見えている。「高齢では何も見えない」は誇張である。
・36. Wu 2010（Atherosclerosis 2010;213:173-7、428 名）は、加齢・動脈硬化例で波形が歪み
　従来パラメータの決定が妨げられることを問題にし、拡張期ピークが判別できない例でも
　使える指標として NCT（正規化 crest time）と CTR（crest time ratio）を提案した。

【論点】
本当の問題は「全員が見えない」ことではなく、見えない 14% だけが解析から系統的に
脱落し、選択バイアスを生むことである。次章の PDA はこの 14% を救えるかという話になる。

出典：6. Cunningham 2023 ／ 22. Elgendi 2012 ／ 36. Wu 2010 ／ 38. Pal 2024""")

prs.save('r2_s3.pptx')
print('r2 stage3 ok:', len(prs.slides._sldIdLst), 'slides')
