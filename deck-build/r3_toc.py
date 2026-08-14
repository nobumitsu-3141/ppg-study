# -*- coding: utf-8 -*-
"""目次ページの共通レンダラ（3階層対応・行数に応じて文字と行間を自動調整）"""
import re
from lxml import etree
from pptx.util import Emu, Pt
from pptx.enum.text import MSO_ANCHOR
from deckkit import NS, INK, GOLD, LGREY, FONT

# 版面：金色下線のすぐ下からページ番号の手前まで目一杯使う
BOX_X, BOX_Y = 0.94, 1.68
BOX_W, BOX_H = 11.57, 5.42
JP_H = 1.18                      # メイリオの行高安全係数

# 行数に応じた（見出しpt, 小見出しpt, 行間）の候補。上から順に入るものを採る
BANDS = [
    (34, 30, 1.45),
    (32, 29, 1.42),
    (30, 27, 1.40),
    (30, 27, 1.30),
    (28, 26, 1.32),
    (28, 25, 1.24),
    (26, 24, 1.26),
    (26, 24, 1.16),
    (25, 23, 1.16),
    (24, 22, 1.14),
    (24, 22, 1.06),
    (23, 22, 1.02),
]


def _fits(lines, main, sub, ls):
    total = 0.0
    for lv, _, _ in lines:
        pt = main if lv == 0 else sub
        total += pt / 72.0 * ls * JP_H
    return total <= BOX_H


def pick_band(lines):
    for main, sub, ls in BANDS:
        if _fits(lines, main, sub, ls):
            return main, sub, ls
    return BANDS[-1]


def build_lines(chapters, active_ch, active_sub=None, sub2=None):
    """[(level, 番号, 見出し), ...] を組み立てる"""
    out = []
    for i in sorted(chapters):
        out.append((0, f'{i}.', chapters[i][0]))
        if i != active_ch:
            continue
        for entry in chapters[i][1]:
            num, _, label = entry.partition('　')
            out.append((1, num, label))
            if sub2 and active_sub and num == active_sub:
                for n2, l2 in sub2.get(num, []):
                    out.append((2, n2, l2))
    return out


def render_agenda(shape, lines, main=None, sub=None, ls=None, active_ch=None):
    if main is None:
        main, sub, ls = pick_band(lines)
    shape.left, shape.top = Emu(int(BOX_X * 914400)), Emu(int(BOX_Y * 914400))
    shape.width, shape.height = Emu(int(BOX_W * 914400)), Emu(int(BOX_H * 914400))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    txBody = tf._txBody
    for p in txBody.findall('a:p', NS):
        txBody.remove(p)
    for lv, num, label in lines:
        on = (lv > 0) or (active_ch is not None and num.rstrip('.') == str(active_ch))
        pt = main if lv == 0 else sub
        indent = '' if lv == 0 else ('　' if lv == 1 else '　　　')
        num_col = GOLD if on else LGREY
        txt_col = INK if on else LGREY
        bold_lbl = (lv == 0 and on)
        p = etree.SubElement(txBody, '{%s}p' % NS['a'])
        pPr = etree.SubElement(p, '{%s}pPr' % NS['a'])
        lnSpc = etree.SubElement(pPr, '{%s}lnSpc' % NS['a'])
        etree.SubElement(lnSpc, '{%s}spcPct' % NS['a']).set(
            'val', str(int(round(ls * 100000))))
        for txt, col, bold in ((indent + num + ' ', num_col, True),
                               (label, txt_col, bold_lbl)):
            r = etree.SubElement(p, '{%s}r' % NS['a'])
            rPr = etree.SubElement(r, '{%s}rPr' % NS['a'])
            rPr.set('sz', str(int(pt * 100)))
            rPr.set('b', '1' if bold else '0')
            sf = etree.SubElement(rPr, '{%s}solidFill' % NS['a'])
            etree.SubElement(sf, '{%s}srgbClr' % NS['a']).set('val', col)
            for tag in ('latin', 'ea', 'cs'):
                etree.SubElement(rPr, '{%s}%s' % (NS['a'], tag)).set('typeface', FONT)
            etree.SubElement(r, '{%s}t' % NS['a']).text = txt


def collect_sub2(prs, title_text):
    """デック内の 3 階層見出し（2.1.1 など）を 2.1 ごとにまとめる"""
    sub2 = {}
    for s in prs.slides:
        m = re.match(r'^(\d+\.\d+)\.(\d+)\s+(.*)$', title_text(s).strip())
        if not m:
            continue
        parent, k, label = m.group(1), m.group(2), m.group(3).strip()
        num = f'{parent}.{k}'
        lst = sub2.setdefault(parent, [])
        if not any(n == num for n, _ in lst):
            lst.append((num, label))
    for v in sub2.values():
        v.sort(key=lambda t: [int(x) for x in t[0].split('.')])
    return sub2


# ─────────────────────────────────────────── 節目次（2欄）
LCOL_X, LCOL_W = 0.90, 5.50
RCOL_X, RCOL_W = 6.90, 5.90

BANDS_COL = [
    (32, 29, 1.45), (30, 27, 1.42), (30, 27, 1.32), (28, 26, 1.38),
    (28, 25, 1.28), (26, 24, 1.32), (26, 24, 1.20), (24, 22, 1.22),
    (24, 22, 1.10), (23, 22, 1.04),
]


def _fits_w(lines, main, sub, ls, h):
    total = sum((main if it[0] == 0 else sub) / 72.0 * ls * JP_H for it in lines)
    return total <= h


def pick_band_col(lines, h=BOX_H):
    for main, sub, ls in BANDS_COL:
        if _fits_w(lines, main, sub, ls, h):
            return main, sub, ls
    return BANDS_COL[-1]


def render_column(shape, lines, x, w, main, sub, ls, active_ch=None,
                  y=BOX_Y, h=BOX_H, anchor_middle=True):
    from pptx.enum.text import MSO_ANCHOR as _A
    shape.left, shape.top = Emu(int(x * 914400)), Emu(int(y * 914400))
    shape.width, shape.height = Emu(int(w * 914400)), Emu(int(h * 914400))
    tf = shape.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = _A.MIDDLE if anchor_middle else _A.TOP
    txBody = tf._txBody
    for p in txBody.findall('a:p', NS):
        txBody.remove(p)
    for lv, num, label, *rest in lines:
        opt = rest[0] if rest else {}
        on = opt.get('on')
        if on is None:
            on = (lv > 0) or (active_ch is not None and num.rstrip('.') == str(active_ch))
        pt = opt.get('size', main if lv == 0 else sub)
        indent = opt.get('indent', '' if lv == 0 else '　')
        num_col = opt.get('num_color', GOLD if on else LGREY)
        txt_col = opt.get('color', INK if on else LGREY)
        bold_lbl = opt.get('bold', lv == 0 and on)
        p = etree.SubElement(txBody, '{%s}p' % NS['a'])
        pPr = etree.SubElement(p, '{%s}pPr' % NS['a'])
        lnSpc = etree.SubElement(pPr, '{%s}lnSpc' % NS['a'])
        etree.SubElement(lnSpc, '{%s}spcPct' % NS['a']).set(
            'val', str(int(round(ls * 100000))))
        segs = [(indent + num + (' ' if num else ''), num_col, True),
                (label, txt_col, bold_lbl)]
        for txt, col, bold in segs:
            if not txt:
                continue
            r = etree.SubElement(p, '{%s}r' % NS['a'])
            rPr = etree.SubElement(r, '{%s}rPr' % NS['a'])
            rPr.set('sz', str(int(pt * 100)))
            rPr.set('b', '1' if bold else '0')
            sf = etree.SubElement(rPr, '{%s}solidFill' % NS['a'])
            etree.SubElement(sf, '{%s}srgbClr' % NS['a']).set('val', col)
            for tag in ('latin', 'ea', 'cs'):
                etree.SubElement(rPr, '{%s}%s' % (NS['a'], tag)).set('typeface', FONT)
            etree.SubElement(r, '{%s}t' % NS['a']).text = txt
