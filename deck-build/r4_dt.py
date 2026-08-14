# -*- coding: utf-8 -*-
"""v1.1 → v1.2：4章に「ΔT の構造的な強み」3 枚を図解で追加する

挿入位置は「4.3 なぜ SI・RI か」（現 74 枚目）の直後。
ページ番号は挿入点より後ろだけを +3 し、総数を 55 → 58 に書き換える
（既存の論理ページ区切りはユーザーの手編集を尊重してそのまま使う）。
"""
import re
import sys
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from deckkit import (add_text, add_box, add_line, add_freeform, add_arrow,
                     set_notes, slide_index, set_pageno, INK, GREY, LGREY, FONT)
from newpages import Builder, ppg_curve

SRC, OUT = sys.argv[1], sys.argv[2]
prs = Presentation(SRC)
S = lambda n: list(prs.slides)[n - 1]

# ── この資料 4 章の配色語彙（既存スライドから採取）────────────────
NAVY, NAVY_BG = '2E5395', 'EAF0F8'      # 青パネル
ORNG, ORNG_BG = 'C55A11', 'FDF0E6'      # 橙パネル
TEAL, TEAL_BG = '00A8AA', 'E4F5F5'      # 青緑パネル
GOLD, GOLD_BG = 'BF9000', 'FFF3D6'      # 金＝要点
REDD, REDD_BG = 'C00000', 'FCEAEA'      # 赤＝警告
WHITE = 'FFFFFF'

PANEL_L, PANEL_R, PANEL_W = 0.55, 6.79, 5.94
BAN_X, BAN_W = 0.55, 12.18

b = Builder(prs, S(74))
anchor = S(74)


def panel(s, x, y, w, h, line, fill, name):
    return add_box(s, x, y, w, h, fill=fill, line=line, line_w=1.5,
                   radius=0.06, name=name)


def banner(s, y, h, line, fill, text, size=26, color=INK):
    add_box(s, BAN_X, y, BAN_W, h, fill=fill, line=line, line_w=2.25)
    add_text(s, BAN_X + 0.2, y, BAN_W - 0.4, h, text, size=size, bold=True,
             color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def head(s, x, y, w, text, color):
    add_text(s, x + 0.22, y, w - 0.44, 0.52, text, size=24, bold=True,
             color=color, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def seg(s, x, y, w, h, fill, label, size=22, color=WHITE):
    """塗りつぶしの帯＋その中央に白文字"""
    add_box(s, x, y, w, h, fill=fill, line=None, radius=0.10)
    if label:
        add_text(s, x, y, w, h, label, size=size, bold=True, color=color,
                 align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def peaks(pts):
    """PPT 座標（y は下向き）の点列から、山＝y の極小を左から拾う"""
    out = []
    for i in range(2, len(pts) - 2):
        y = pts[i][1]
        if y < pts[i - 2][1] and y < pts[i + 2][1] and y <= pts[i - 1][1] and y <= pts[i + 1][1]:
            if not out or pts[i][0] - out[-1][0] > 0.3:
                out.append(pts[i])
    return out


# ══════════════════════════════════════════ ① ΔT は PEP を含まない
s1 = b.new(anchor, '4.3  ΔT は PEP を含まない', 4,
           source='34. Sugo 2012　／　43. Beutel 2021')

PY, PH = 1.95, 3.10
# --- 左：PWTT は心電図が起点 --------------------------------------
panel(s1, PANEL_L, PY, PANEL_W, PH, NAVY, NAVY_BG, 'l-bg')
head(s1, PANEL_L, PY + 0.16, PANEL_W, 'PWTT ＝ 心電図が起点', NAVY)

BX0, BXM, BX1 = 0.85, 2.85, 6.20
add_text(s1, BX0, 2.58, 1.40, 0.42, 'R 波', size=16, color=GREY,
         align=PP_ALIGN.LEFT)
add_text(s1, BXM - 0.75, 2.58, 1.50, 0.42, '弁が開く', size=16, color=GREY,
         align=PP_ALIGN.CENTER)
add_text(s1, BX1 - 1.55, 2.58, 1.55, 0.42, '脈波が届く', size=16, color=GREY,
         align=PP_ALIGN.RIGHT)
for xt in (BX0, BXM, BX1):
    add_line(s1, xt, 3.00, xt, 3.63, color=GREY, width=1.0)
seg(s1, BX0, 3.06, BXM - BX0, 0.52, ORNG, 'PEP')
seg(s1, BXM, 3.06, BX1 - BXM, 0.52, NAVY, 'VTT')
add_text(s1, BX0, 3.68, BXM - BX0, 0.56, '心臓側', size=22, bold=True,
         color=ORNG, align=PP_ALIGN.CENTER)
add_text(s1, BXM, 3.68, BX1 - BXM, 0.56, '血管側', size=22, bold=True,
         color=NAVY, align=PP_ALIGN.CENTER)
add_text(s1, BX0, 4.30, 5.35, 0.60, '起点が心臓 → PEP が必ず混ざる',
         size=22, bold=True, color=INK, align=PP_ALIGN.CENTER)

# --- 右：ΔT は 1 拍の中だけ ---------------------------------------
panel(s1, PANEL_R, PY, PANEL_W, PH, TEAL, TEAL_BG, 'r-bg')
head(s1, PANEL_R, PY + 0.16, PANEL_W, 'ΔT ＝ 脈波 1 拍の中', TEAL)

WX, WY, WW, WH = 7.35, 2.72, 5.00, 1.15
pts = ppg_curve(WX, WY, WW, WH, [(0.24, 0.085, 1.00), (0.56, 0.115, 0.52)])
add_line(s1, WX, WY + WH, WX + WW, WY + WH, color=LGREY, width=1.25)
add_freeform(s1, pts, color=INK, width=2.75)
pk = peaks(pts)
p1, p2 = pk[0], pk[1]
for px, py, lab, side in ((p1[0], p1[1], 'P1', -1), (p2[0], p2[1], 'P2', +1)):
    add_line(s1, px, py, px, 4.06, color=TEAL, width=1.25, dash='dash')
    lx = px - 0.95 if side < 0 else px + 0.10
    add_text(s1, lx, py - 0.14, 0.85, 0.42, lab, size=16, bold=True, color=GREY,
             align=PP_ALIGN.RIGHT if side < 0 else PP_ALIGN.LEFT)
add_arrow(s1, p1[0], 4.02, p2[0] - p1[0], 0.18, color=TEAL,
          shape=MSO_SHAPE.LEFT_RIGHT_ARROW)
add_text(s1, (p1[0] + p2[0]) / 2 - 0.70, 4.22, 1.40, 0.60, 'ΔT', size=24,
         bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s1, PANEL_R + 0.22, 4.30, PANEL_W - 0.44, 0.60,
         '外の基準点を使わない', size=22, bold=True, color=INK,
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.BOTTOM)

banner(s1, 5.25, 0.90, GOLD, GOLD_BG,
       'ΔT は近位の基準点を持たない ＝ PEP をゼロ含有', size=26)
banner(s1, 6.22, 0.85, TEAL, TEAL_BG,
       '同じ SpO₂ 波形から「血管だけの時間」が取れる', size=24)

set_notes(s1, '''PWTT（モニタ表示）は心電図の R 波を起点にしているので、
定義上そこに PEP が入り込む。大動脈弁が開いた瞬間を非侵襲・連続に
捉える手段がない以上、PWTT から PEP を除くことはできない。

一方 ΔT は、同じ 1 拍の PPG の中で P1 と P2 の時間差を取るだけなので、
心臓側の基準点をまったく使わない。つまり ΔT は構造的に PEP をゼロ含有する。
これは PWTT・PAT には原理的に不可能な性質であり、しかも追加センサを要さない。

Sugo 2012 は、PEP の変化が PWTT の変化の約半分を占めることを実測で示している。
Beutel 2021 は PAT を心臓区間と血管区間に分ける手法を提案しているが、
そこでは別の信号を必要とする。''')

# ══════════════════════════════════════════ ② ΔT で PWTT を分解する
s2 = b.new(s1, '4.3  ΔT で PWTT を分解する', 4,
           source='43. Beutel 2021　／　47. Kortekaas 2018')

for x, col, bg, txt in ((PANEL_L, NAVY, NAVY_BG, 'VTT ≈ L ÷ PWV'),
                        (PANEL_R, TEAL, TEAL_BG, 'ΔT ≈ 2d ÷ PWV')):
    add_box(s2, x, 1.92, PANEL_W, 0.72, fill=bg, line=col, line_w=1.5, radius=0.10)
    add_text(s2, x, 1.92, PANEL_W, 0.72, txt, size=26, bold=True, color=col,
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
add_box(s2, BAN_X, 2.74, BAN_W, 0.68, fill=GOLD_BG, line=GOLD, line_w=2.25)
add_text(s2, BAN_X + 0.2, 2.74, BAN_W - 0.4, 0.68,
         'どちらも 1/PWV に比例　→　PWTT ＝ PEP ＋ k・ΔT',
         size=24, bold=True, color=INK, align=PP_ALIGN.CENTER,
         anchor=MSO_ANCHOR.MIDDLE)

LBL_X, LBL_W = 0.55, 1.85
RX0, RXM, RX1 = 2.55, 4.95, 12.55
ROWS = [
    (3.62, 'PWTT', INK, None),
    (4.26, 'k・ΔT', TEAL, None),
    (4.90, '差', ORNG, None),
]
for y, lab, col, _ in ROWS:
    add_text(s2, LBL_X, y, LBL_W, 0.52, lab, size=22, bold=True, color=col,
             align=PP_ALIGN.RIGHT, anchor=MSO_ANCHOR.MIDDLE)
# 1 行目：PWTT ＝ PEP ＋ VTT
seg(s2, RX0, 3.62, RXM - RX0, 0.52, ORNG, 'PEP')
seg(s2, RXM, 3.62, RX1 - RXM, 0.52, NAVY, 'VTT')
# 2 行目：k・ΔT を VTT の真下に置く
seg(s2, RXM, 4.26, RX1 - RXM, 0.52, TEAL, 'ΔT で置き換える')
# 3 行目：引き算の残り＝PEP
seg(s2, RX0, 4.90, RXM - RX0, 0.52, ORNG, 'PEP')
add_text(s2, RXM + 0.25, 4.90, 7.00, 0.52, '＝ 心臓側の時間だけが残る',
         size=22, bold=True, color=INK, align=PP_ALIGN.LEFT,
         anchor=MSO_ANCHOR.MIDDLE)

banner(s2, 5.60, 0.70, GOLD, GOLD_BG, 'ΔPWTT − k・ΔΔT ≈ ΔPEP', size=28)
banner(s2, 6.38, 0.70, REDD, REDD_BG,
       '成立条件 ― k ＝ L ÷ 2d が術中に一定であること', size=24, color=REDD)

set_notes(s2, '''VTT も ΔT も、どちらも 1/PWV に比例する時間である。
であれば ΔT を VTT の代理として使い、PWTT から差し引けば、
残るのは PEP の変化分ということになる。

esCCO の限界は「大動脈弁開放を捉えていない」一点に集約されるが、
この方法は弁開放を捉えずに、PEP が動いた分だけを差として取り出す。
Beutel 2021 が別の信号で解こうとした問題を、同じ SpO₂ 波形だけで
扱えることになる。

ただし成立条件がある。上肢（大動脈→指、VTT）と下半身（大動脈→反射点→戻り、
ΔT）は別の血管床であり、k ＝ L ÷ 2d が術中も一定という保証はない。
昇圧薬が両者に同じ比率で効くとは限らない。ここは実測で確かめる必要がある。

なお Kortekaas 2018 は、安静時なら PEP の個人内変動は小さく
（SD 1.6〜12.0 ms）、PWTT の変化はほぼ VTT の変化とみなせると報告している。
とすれば ΔT に期待すべきは「常に PEP を剥がす」ことではなく、
「PEP が実際に動いた区間を見つけて、そこだけ別扱いにする」ことかもしれない。''')

# ══════════════════════════════════════════ ③ 個人内なら身長は要らない
s3 = b.new(s2, '4.3  身長は要らない', 4,
           source='6. Millasseau 2002　／　28. von Wowern 2015')

PY3, PH3 = 1.95, 3.30
panel(s3, PANEL_L, PY3, PANEL_W, PH3, ORNG, ORNG_BG, 'l-bg')
head(s3, PANEL_L, PY3 + 0.16, PANEL_W, '個人間で並べる', ORNG)
add_text(s3, PANEL_L + 0.30, 2.74, PANEL_W - 0.60, 2.32,
         ['・SI ＝ 身長 ÷ ΔT', '・速度に換算するため身長が要る',
          '・cfPWV と r ＝ 0.65', '・順位づけには足りない'],
         size=22, bold=True, color=INK, line_spacing=1.30)

panel(s3, PANEL_R, PY3, PANEL_W, PH3, TEAL, TEAL_BG, 'r-bg')
head(s3, PANEL_R, PY3 + 0.16, PANEL_W, '同じ人を追う', TEAL)
add_text(s3, PANEL_R + 0.30, 2.74, PANEL_W - 0.60, 2.32,
         ['・身長は定数 → 約分で消える', '・ΔSI ÷ SI ＝ − ΔΔT ÷ ΔT',
          '・被験者内 CV ＝ 9.6 %', '・追跡には十分'],
         size=22, bold=True, color=INK, line_spacing=1.30)

banner(s3, 5.42, 0.82, GOLD, GOLD_BG, '見るのは ΔT の個人内変化だけでよい', size=26)
banner(s3, 6.34, 0.78, TEAL, TEAL_BG,
       '血管の情報を得るのに患者背景の入力は要らない', size=24)

set_notes(s3, '''SI ＝ 身長 ÷ ΔT の「身長」は、ΔT を速度の次元に直して
個人間で並べるためだけに入っている。同じ人を追うなら身長は定数なので、
比を取った瞬間に約分で消える。身長を捨てるのは妥協ではなく、恒等式である。

その結果、効いてくる統計量が入れ替わる。Millasseau 2002 が報告した
cfPWV との相関 r ＝ 0.65 は「集団の中で個人を序列づけるには不足」を意味するが、
同じ論文の被験者内変動係数 9.6 %（週間隔）は「同じ人を追うには十分」を意味する。
同じ論文の中で結論が逆を向くのは、両者が別の問いに答えているからである。

したがって「SI は cfPWV の代替にならない」という批判は、
PWV を主張しないこの使い方には当たらない。''')

# ══════════════════ 参考文献 47 を追加（末尾に足すので既存番号は動かない）
from lxml import etree
from deckkit import NS
from pptx.dml.color import RGBColor

ref = ('47. Kortekaas MC, et al. Small intra-individual variability of the '
       'pre-ejection period justifies the use of pulse transit time as '
       'approximation of the vascular transit. PLoS One 2018;13(10):e0204105.')
import copy as _copy
box = [sh for sh in S(87).shapes if sh.name == 'TextBox 16'][0]
last = box.text_frame.paragraphs[-1]
last._p.addnext(_copy.deepcopy(last._p))
para = box.text_frame.paragraphs[-1]
runs = para.runs
runs[0].text = ref
for r in runs[1:]:
    r.text = ''
print('参考文献 47 を追加')

# ══════════════════ ページ番号：挿入点より後ろを +3、総数 55 → 58
NEW_TOTAL = 58
CUT = 45                      # 「4.3 なぜ SI・RI か」＝ 45/55
for i, s in enumerate(prs.slides, 1):
    if s in (s1, s2, s3):
        continue
    for sh in list(s.shapes):
        if sh.name != 'PageNo':
            continue
        m = re.match(r'^(\d+)\s*/\s*(\d+)$', sh.text_frame.text.strip())
        if not m:
            continue
        n = int(m.group(1))
        set_pageno(s, f'{n if n <= CUT else n + 3}/{NEW_TOTAL}')
for k, s in enumerate((s1, s2, s3), start=CUT + 1):
    set_pageno(s, f'{k}/{NEW_TOTAL}')

prs.save(OUT)
print(f'{OUT}: {len(prs.slides._sldIdLst)} slides '
      f'（挿入 {slide_index(prs, s1)+1}〜{slide_index(prs, s3)+1} 枚目）')
