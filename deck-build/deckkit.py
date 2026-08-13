"""川副式デック編集ヘルパー（python-pptx ベース）"""
import copy
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from lxml import etree

EMU = 914400
NS = {
    'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
    'p': 'http://schemas.openxmlformats.org/presentationml/2006/main',
    'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
}

FONT = "Meiryo"
GOLD = "BF9000"
BLUE = "0072B2"
VERM = "D55E00"
TEAL = "00A8AA"
RED = "C00000"
GREY = "808080"
LGREY = "D9D9D9"
INK = "262626"

SRC_X, SRC_Y = 720000, 6606000        # 出典テキストボックスの定位置
PAGE_X, PAGE_Y = 11049000, 6606000    # ページ番号の定位置

CHAPTERS = ["基礎", "加齢", "循環", "麻酔", "定量化", "まとめ", "限界"]


def inches(v):
    return int(round(v * EMU))


# ---------------------------------------------------------------- slide clone
def clone_slide(prs, src_slide, insert_after_idx=None):
    """src_slide を複製した新スライドを返す。insert_after_idx で位置指定。"""
    layout = src_slide.slide_layout
    new = prs.slides.add_slide(layout)
    # レイアウト由来のプレースホルダを一旦すべて削除
    for shp in list(new.shapes):
        shp._element.getparent().remove(shp._element)
    # 図形をディープコピー
    for shp in src_slide.shapes:
        new.shapes._spTree.append(copy.deepcopy(shp._element))
    # 画像などの関係を移送（レイアウト・ノートは複製先で新規に持たせる）
    for rid, rel in src_slide.part.rels.items():
        if rel.reltype.endswith('/slideLayout') or rel.reltype.endswith('/notesSlide'):
            continue
        if rid in new.part.rels:
            continue
        if rel.is_external:
            new.part.rels.get_or_add_ext_rel(rel.reltype, rel.target_ref)
        else:
            new.part.rels._add_relationship(rel.reltype, rel._target, rid)
    # 位置調整
    if insert_after_idx is not None:
        move_slide(prs, len(prs.slides) - 1, insert_after_idx + 1)
    return new


def move_slide(prs, from_idx, to_idx):
    sldIdLst = prs.slides._sldIdLst
    ids = list(sldIdLst)
    el = ids[from_idx]
    sldIdLst.remove(el)
    ids2 = list(sldIdLst)
    if to_idx >= len(ids2):
        sldIdLst.append(el)
    else:
        ids2[to_idx].addprevious(el)


def slide_index(prs, slide):
    for i, s in enumerate(prs.slides):
        if s.slide_id == slide.slide_id:
            return i
    return -1


# ------------------------------------------------------------------- cleaning
def keep_only(slide, names):
    """指定名の図形だけ残して他を削除"""
    for shp in list(slide.shapes):
        if shp.name not in names:
            shp._element.getparent().remove(shp._element)


def drop_shapes(slide, pred):
    for shp in list(slide.shapes):
        if pred(shp):
            shp._element.getparent().remove(shp._element)


def strip_content(slide):
    """タイトル・章ナビ・ページ番号・出典以外を削除"""
    def is_keeper(shp):
        n = shp.name
        if 'Off-page Connector' in n:
            return True
        if n in ('Title 1', 'タイトル 1', 'Source'):
            return True
        if shp.has_text_frame and shp.top is not None and shp.top > 6400000 \
           and shp.left is not None and shp.left > 10500000:
            return True   # ページ番号
        return False
    drop_shapes(slide, lambda s: not is_keeper(s))


# ---------------------------------------------------------------------- title
def get_title(slide):
    for shp in slide.shapes:
        if shp.name in ('Title 1', 'タイトル 1'):
            return shp
    return None


def set_title(slide, text, size=46):
    shp = get_title(slide)
    if shp is None:
        return None
    tf = shp.text_frame
    txBody = tf._txBody
    for p in txBody.findall('a:p', NS):
        txBody.remove(p)
    p = etree.SubElement(txBody, '{%s}p' % NS['a'])
    pPr = etree.SubElement(p, '{%s}pPr' % NS['a'])
    pPr.set('algn', 'l')
    r = etree.SubElement(p, '{%s}r' % NS['a'])
    rPr = etree.SubElement(r, '{%s}rPr' % NS['a'])
    rPr.set('sz', str(int(size * 100)))
    rPr.set('dirty', '0')
    for tag in ('latin', 'ea', 'cs'):
        e = etree.SubElement(rPr, '{%s}%s' % (NS['a'], tag))
        e.set('typeface', FONT)
    t = etree.SubElement(r, '{%s}t' % NS['a'])
    t.text = text
    return shp


# ------------------------------------------------------------------- nav chips
def set_nav(slide, active_ch):
    """active_ch: 1..7 or None（全て灰）"""
    chips = [s for s in slide.shapes if 'Off-page Connector' in s.name]
    chips.sort(key=lambda s: s.left)
    for i, chip in enumerate(chips, start=1):
        on = (active_ch is not None and i == active_ch)
        fill = TEAL if on else LGREY
        col = "FFFFFF" if on else GREY
        spPr = chip._element.find('p:spPr', NS)
        for sf in spPr.findall('a:solidFill', NS):
            spPr.remove(sf)
        sf = etree.Element('{%s}solidFill' % NS['a'])
        clr = etree.SubElement(sf, '{%s}srgbClr' % NS['a'])
        clr.set('val', fill)
        geom = spPr.find('a:prstGeom', NS)
        geom.addnext(sf)
        for rPr in chip._element.iter('{%s}rPr' % NS['a']):
            for old in rPr.findall('a:solidFill', NS):
                rPr.remove(old)
            sf2 = etree.Element('{%s}solidFill' % NS['a'])
            c2 = etree.SubElement(sf2, '{%s}srgbClr' % NS['a'])
            c2.set('val', col)
            rPr.insert(0, sf2)


# ---------------------------------------------------------------------- source
def set_source(slide, text):
    """左下の出典を統一書式で設定（既存は置換）"""
    drop_shapes(slide, lambda s: s.has_text_frame and s.top is not None
                and s.top > 6300000 and (s.left is None or s.left < 10500000))
    if not text:
        return None
    box = slide.shapes.add_textbox(Emu(SRC_X), Emu(SRC_Y), Emu(int(11.9 * EMU)), Emu(200000))
    box.name = 'Source'
    tf = box.text_frame
    tf.word_wrap = True
    bp = tf._txBody.find('a:bodyPr', NS)
    for k, v in (('lIns', '0'), ('tIns', '0'), ('rIns', '0'), ('bIns', '0'), ('anchor', 't')):
        bp.set(k, v)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(10.5)
    f.bold = False
    f.name = FONT
    f.color.rgb = RGBColor.from_string(GREY)
    _force_ea(run, FONT)
    return box


def set_pageno(slide, text):
    drop_shapes(slide, lambda s: s.has_text_frame and s.top is not None
                and s.top > 6300000 and s.left is not None and s.left > 10500000)
    if not text:
        return None
    box = slide.shapes.add_textbox(Emu(PAGE_X), Emu(PAGE_Y), Emu(inches(1.10)), Emu(inches(0.24)))
    box.name = 'PageNo'
    tf = box.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = text
    f = run.font
    f.size = Pt(12)
    f.name = FONT
    f.color.rgb = RGBColor.from_string(GREY)
    _force_ea(run, FONT)
    return box


def _force_ea(run, face):
    rPr = run._r.get_or_add_rPr()
    for tag in ('ea', 'cs'):
        for old in rPr.findall('a:%s' % tag, NS):
            rPr.remove(old)
        e = etree.SubElement(rPr, '{%s}%s' % (NS['a'], tag))
        e.set('typeface', face)


# ------------------------------------------------------------------ text boxes
def add_text(slide, x, y, w, h, lines, size=24, bold=False, color=INK,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space_after=0,
             wrap=True, line_spacing=None, name=None):
    """lines: str | list[str] | list[list[(text, {opts})]]"""
    box = slide.shapes.add_textbox(Emu(inches(x)), Emu(inches(y)),
                                   Emu(inches(w)), Emu(inches(h)))
    if name:
        box.name = name
    tf = box.text_frame
    tf.word_wrap = wrap
    tf.vertical_anchor = anchor
    bp = tf._txBody.find('a:bodyPr', NS)
    for k, v in (('lIns', '0'), ('tIns', '0'), ('rIns', '0'), ('bIns', '0')):
        bp.set(k, v)
    if isinstance(lines, str):
        lines = [lines]
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        if space_after:
            p.space_after = Pt(space_after)
        if line_spacing:
            p.line_spacing = line_spacing
        segs = line if isinstance(line, list) else [(line, {})]
        for seg in segs:
            txt, opt = (seg if isinstance(seg, tuple) else (seg, {}))
            run = p.add_run()
            run.text = txt
            f = run.font
            f.size = Pt(opt.get('size', size))
            f.bold = opt.get('bold', bold)
            f.name = FONT
            f.color.rgb = RGBColor.from_string(opt.get('color', color))
            _force_ea(run, FONT)
    return box


def add_box(slide, x, y, w, h, fill=None, line=None, line_w=1.5, radius=None,
            shape=MSO_SHAPE.ROUNDED_RECTANGLE, name=None):
    shp = slide.shapes.add_shape(shape, Emu(inches(x)), Emu(inches(y)),
                                 Emu(inches(w)), Emu(inches(h)))
    if name:
        shp.name = name
    shp.shadow.inherit = False
    if fill:
        shp.fill.solid()
        shp.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = RGBColor.from_string(line)
        shp.line.width = Pt(line_w)
    else:
        shp.line.fill.background()
    if radius is not None and shape == MSO_SHAPE.ROUNDED_RECTANGLE:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    shp.text_frame.text = ''
    return shp


def add_arrow(slide, x, y, w, h, color=GREY, shape=MSO_SHAPE.RIGHT_ARROW):
    shp = slide.shapes.add_shape(shape, Emu(inches(x)), Emu(inches(y)),
                                 Emu(inches(w)), Emu(inches(h)))
    shp.shadow.inherit = False
    shp.fill.solid()
    shp.fill.fore_color.rgb = RGBColor.from_string(color)
    shp.line.fill.background()
    shp.text_frame.text = ''
    return shp


def add_line(slide, x1, y1, x2, y2, color=GREY, width=1.5, dash=None):
    from pptx.enum.shapes import MSO_CONNECTOR
    cn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Emu(inches(x1)),
                                    Emu(inches(y1)), Emu(inches(x2)), Emu(inches(y2)))
    cn.line.color.rgb = RGBColor.from_string(color)
    cn.line.width = Pt(width)
    if dash:
        ln = cn.line._get_or_add_ln()
        d = etree.SubElement(ln, '{%s}prstDash' % NS['a'])
        d.set('val', dash)
    return cn


def add_freeform(slide, pts, color=INK, width=2.5, dash=None):
    """pts: [(x_in, y_in), ...] を折れ線で描く"""
    b = slide.shapes.build_freeform(Emu(inches(pts[0][0])), Emu(inches(pts[0][1])))
    b.add_line_segments([(Emu(inches(x)), Emu(inches(y))) for x, y in pts[1:]],
                        close=False)
    shp = b.convert_to_shape()
    shp.fill.background()
    shp.line.color.rgb = RGBColor.from_string(color)
    shp.line.width = Pt(width)
    shp.shadow.inherit = False
    if dash:
        ln = shp.line._get_or_add_ln()
        d = etree.SubElement(ln, '{%s}prstDash' % NS['a'])
        d.set('val', dash)
    return shp


# -------------------------------------------------------------------- notes
def set_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text


def append_notes(slide, text):
    tf = slide.notes_slide.notes_text_frame
    cur = tf.text
    tf.text = (cur + "\n\n" + text) if cur.strip() else text


# --------------------------------------------------------------------- tables
def add_table(slide, x, y, w, h, rows, cols, col_widths=None):
    gf = slide.shapes.add_table(rows, cols, Emu(inches(x)), Emu(inches(y)),
                                Emu(inches(w)), Emu(inches(h)))
    tbl = gf.table
    if col_widths:
        for i, cw in enumerate(col_widths):
            tbl.columns[i].width = Emu(inches(cw))
    return gf, tbl


def cell_text(cell, lines, size=22, bold=False, color=INK, align=PP_ALIGN.LEFT,
              fill=None, anchor=MSO_ANCHOR.MIDDLE):
    tf = cell.text_frame
    tf.word_wrap = True
    cell.vertical_anchor = anchor
    cell.margin_left = Emu(inches(0.08))
    cell.margin_right = Emu(inches(0.08))
    cell.margin_top = Emu(inches(0.04))
    cell.margin_bottom = Emu(inches(0.04))
    if fill:
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor.from_string(fill)
    else:
        cell.fill.background()
    if isinstance(lines, str):
        lines = [lines]
    for p in list(tf.paragraphs)[1:]:
        p._p.getparent().remove(p._p)
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        segs = line if isinstance(line, list) else [(line, {})]
        for seg in segs:
            txt, opt = (seg if isinstance(seg, tuple) else (seg, {}))
            run = p.add_run()
            run.text = txt
            f = run.font
            f.size = Pt(opt.get('size', size))
            f.bold = opt.get('bold', bold)
            f.name = FONT
            f.color.rgb = RGBColor.from_string(opt.get('color', color))
            _force_ea(run, FONT)


def table_style_plain(gf):
    """バンド塗り・見出し強調をオフにして自前で塗る"""
    tblPr = gf.table._tbl.find('.//{%s}tblPr' % NS['a'])
    if tblPr is not None:
        tblPr.set('firstRow', '0')
        tblPr.set('bandRow', '0')
