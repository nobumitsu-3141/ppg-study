# -*- coding: utf-8 -*-
"""Stage 3: 参考文献ページの再構成と、ページ番号の一括再採番"""
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from deckkit import (add_text, set_title, set_nav, set_source, set_pageno,
                     clone_slide, strip_content, slide_index, drop_shapes,
                     INK, GOLD, GREY)
from refs import REFS

prs = Presentation('deck_s2d.pptx')


def titles():
    out = []
    for s in prs.slides:
        t = ''
        for sh in s.shapes:
            if sh.name in ('Title 1', 'タイトル 1') and sh.has_text_frame:
                t = sh.text_frame.text.strip()
        out.append((s, t))
    return out


# ------------------------------------------------------ 1) 参考文献ページ再構成
ref_slides = [s for s, t in titles() if t.startswith('参考文献')]
base = ref_slides[0]

PER_PAGE = 14           # 1 ページあたり（2 段組 × 7）
nums = sorted(REFS)
pages = [nums[i:i + PER_PAGE] for i in range(0, len(nums), PER_PAGE)]

# 既存 2 枚を使い、足りない分を複製で足す
while len(ref_slides) < len(pages):
    idx = slide_index(prs, ref_slides[-1])
    ns = clone_slide(prs, base, insert_after_idx=idx)
    ref_slides.append(ns)

for pi, (slide, group) in enumerate(zip(ref_slides, pages)):
    # 本文テキストボックスだけ落とす（タイトル・章ナビ・ページ番号は残す）
    drop_shapes(slide, lambda sh: sh.has_text_frame and sh.top is not None
                and 1.6 * 914400 < sh.top < 6.9 * 914400)
    set_title(slide, '参考文献' if pi == 0 else f'参考文献（{pi + 1}）')
    set_nav(slide, None)
    set_source(slide, None)
    half = (len(group) + 1) // 2
    for ci, chunk in enumerate((group[:half], group[half:])):
        if not chunk:
            continue
        x = 0.60 + ci * 6.20
        add_text(slide, x, 1.90, 5.95, 5.10,
                 [f'{n}. {REFS[n]}' for n in chunk],
                 size=10.5, color=INK, space_after=7, line_spacing=1.15)

print('reference pages:', len(pages))

# ------------------------------------------------------------ 2) ページ番号
# 同一タイトルが連続するビルドアップ群は 1 つの論理ページとして数える
tl = titles()
visible = [(s, t) for s, t in tl if s.element.get('show') != '0']

logical, prev = [], None
for s, t in visible:
    key = t if t else '(figure)'
    if prev is not None and key == prev and key != '(figure)':
        logical[-1][1].append(s)
    else:
        logical.append((key, [s]))
    prev = key

total = len(logical)
for i, (key, group) in enumerate(logical, start=1):
    for s in group:
        # 表紙・メニューは番号なし
        set_pageno(s, None if i == 1 else f'{i}/{total}')

# 非表示スライドには番号を置かない
for s, t in tl:
    if s.element.get('show') == '0':
        set_pageno(s, None)

print('logical pages:', total, '/ physical slides:', len(prs.slides._sldIdLst))

prs.save('deck_v67.pptx')
print('saved deck_v67.pptx')
