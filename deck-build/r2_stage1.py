# -*- coding: utf-8 -*-
"""R2 stage1: 章立ての再編（6章廃止・限界を6章・PDA を7章に独立）＋ 26/27 の復元"""
import copy
from pptx import Presentation
from pptx.util import Emu
from deckkit import set_title, get_title, set_source, drop_shapes, slide_index, NS
from r2_common import CH, rewrite_nav, fill_divider, fill_menu, title_text

prs = Presentation('deck_v68.pptx')
S = list(prs.slides)


def sl(n):
    return S[n - 1]


# ═════════════════════ 1) 26・27 を v6.6 の見た目に戻す
# stage1 で入れたタイトルが画像の背面に回り「2.1」だけが見える壊れ方をしていた。
# 図版1枚ページなので、タイトルと章ナビを外して元の体裁へ戻す。
for n in (26, 27):
    t = sl(n)
    drop_shapes(t, lambda sh: sh.name in ('Title 1', 'タイトル 1')
                or 'Off-page Connector' in sh.name)

# ═════════════════════ 2) 章番号の振り直し（7章 → 6章／PDA を 7章へ）
RETITLE = {
    64: '6. 限界',
    65: '6.1  自動解析の壁',
    66: '6.2  極端な動脈硬化',
    67: '6.3  DN-less 信号対策',
    68: '6.3  DN-less 信号対策',
    69: '7.1  PDA とは',
    70: '7.3  分解の手順',
    71: '7.2  基底関数の種類',
    72: '7.5  Rubins 2008',
    73: '7.5  Goswami 2010',
    74: '7.5  Couceiro 2012',
    75: '7.5  Couceiro 2015',
    76: '7.5  Tigges 2017',
    77: '7.5  Fleischhauer',
    78: '7.5  Basso 2024',
    79: '7.5  Baruch 2014',
    80: '7.5  Pal 2024',
    81: '7.6  PDA の限界',
    82: '7.7  残された問い',
}
for n, t in RETITLE.items():
    set_title(sl(n), t)

# ═════════════════════ 3) 章ナビ帯の文言と強調を全ページ更新
CHAPTER_OF = {}
for i, s in enumerate(prs.slides, 1):
    t = title_text(s)
    if not t:
        CHAPTER_OF[i] = None
        continue
    head = t.split()[0].rstrip('.').split('.')[0]
    CHAPTER_OF[i] = int(head) if head.isdigit() else None
# 表紙・メニュー・参考文献は全灰
for i in (1, 3, 83, 84, 85, 86):
    CHAPTER_OF[i] = None
# 非表示の下書きページ
for i in (87, 88, 89, 90, 91, 92):
    CHAPTER_OF[i] = None

for i, s in enumerate(prs.slides, 1):
    rewrite_nav(s, CHAPTER_OF.get(i))

# ═════════════════════ 4) 章扉・メニューの本文を差し替え
fill_menu(sl(3))
for n, ch in ((4, 1), (21, 2), (34, 3), (44, 4), (51, 5), (64, 6)):
    fill_divider(sl(n), ch)
    set_title(sl(n), f'{ch}. {CH[ch][0]}')
    rewrite_nav(sl(n), ch)

# ═════════════════════ 5) PDA 章の扉を 69 の直前に新設（4章扉を型に使う）
from deckkit import clone_slide, set_pageno
tpl = sl(44)                      # 「4. 麻酔中の波形変化」の扉
idx = slide_index(prs, sl(68))    # 68（DN-less の2枚目）の直後
d = clone_slide(prs, tpl, insert_after_idx=idx)
set_title(d, f'7. {CH[7][0]}')
rewrite_nav(d, 7)
fill_divider(d, 7)
set_source(d, None)
set_pageno(d, None)

prs.save('r2_s1.pptx')
print('r2 stage1 ok:', len(prs.slides._sldIdLst), 'slides')
