# -*- coding: utf-8 -*-
"""新規追加ページの版面を lint と同じ推定式で自動調整する"""
import sys
from pptx import Presentation
from pptx.util import Emu

EMU_CM = 360000
SW = int(33.867 * EMU_CM)
JP_W, JP_H = 1.12, 1.18

NEW_TITLES = {
    '用語の整理', '1.5  PPG基礎まとめ', '4.2  6分類と血行動態',
    '5.3  SDPPG 文献の位置づけ', '5.3  SDPPG 文献の取捨',
    '5.3  b/a と SI', '5.3  d/a と RI', '6. 解析まとめ', '6.1  二軸で読む',
    '7.1  切痕が消えると', '7.1  DN-less 信号とは', '7.1  切痕が無いときの対策',
    '7.1  対策の全体像',
    '7.2  PDA という考え方', '7.2  PDA の手順', '7.2  どんな山を使うか',
    '7.2  PDA の限界', '7.2  残されている問い',
}


def units(text):
    u = 0.0
    for ch in text:
        o = ord(ch)
        if ch == " ":
            u += 0.30
        elif o <= 0x24F or 0x2080 <= o <= 0x208E:
            u += 0.56
        else:
            u += 1.0
    return u


def tw(text, pt):
    return units(text) * (pt / 72 * 2.54) * EMU_CM * JP_W


def is_new(slide):
    for sh in slide.shapes:
        if sh.name in ('Title 1', 'タイトル 1') and sh.has_text_frame:
            t = sh.text_frame.text.strip()
            return t in NEW_TITLES or t.startswith('7.2  文献') or t.startswith('7.2  ①') \
                or t.startswith('7.2  ②') or t.startswith('7.2  ③') or t.startswith('7.2  ④') \
                or t.startswith('7.2  ⑤') or t.startswith('7.2  ⑥') or t.startswith('7.2  ⑦') \
                or t.startswith('7.2  ⑧') or t.startswith('7.2  ⑨') or t.startswith('参考文献')
    return False


def fix(path, out):
    prs = Presentation(path)
    wraps = []
    for si, slide in enumerate(prs.slides, 1):
        if not is_new(slide):
            continue
        for sh in slide.shapes:
            if not sh.has_text_frame or sh.name in ('Title 1', 'タイトル 1', 'Source', 'PageNo'):
                continue
            if sh.top is None or sh.width is None:
                continue
            if not sh.text_frame.text.strip():
                continue
            w = sh.width
            total = 0.0
            need_w = 0
            for para in sh.text_frame.paragraphs:
                runs = [r for r in para.runs if r.text.strip()]
                if not runs:
                    total += 0.2 * EMU_CM
                    continue
                pt = max((r.font.size.pt if r.font.size else 22.0) for r in runs)
                txt = "".join(r.text for r in runs)
                lw = tw(txt, pt)
                need_w = max(need_w, int(lw))
                lines = max(1, -(-int(lw) // max(1, w)))
                ls = para.line_spacing if isinstance(para.line_spacing, (int, float)) else 1.2
                line_h = (pt / 72 * 2.54) * EMU_CM * ls * JP_H
                sa = para.space_after.pt if para.space_after is not None else 6.0
                total += lines * line_h + (sa / 72 * 2.54) * EMU_CM
                if pt >= 20 and lw > w + int(0.15 * EMU_CM):
                    wraps.append((si, round(pt), txt[:44]))
            # 高さを推定描画高さに合わせる（縮めない）
            need_h = int(total * 1.03)
            if need_h > (sh.height or 0):
                sh.height = Emu(need_h)
        # 図形が枠外に出ないよう最終調整
        for sh in slide.shapes:
            if sh.left is None or sh.width is None:
                continue
            if sh.left + sh.width > SW:
                sh.width = Emu(max(EMU_CM, SW - sh.left - int(0.1 * EMU_CM)))
    prs.save(out)
    print('== 幅が足りず自動折返しになる段落 ==')
    for si, pt, txt in wraps:
        print(f'  p{si} {pt}pt  {txt}')
    print('件数:', len(wraps))


if __name__ == '__main__':
    fix(sys.argv[1], sys.argv[2])
