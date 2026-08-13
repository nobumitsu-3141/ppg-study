# -*- coding: utf-8 -*-
"""Stage 2a: 2 / 20 / 49 / 60 / 61 / 62 / 66 / 67 の改訂ページを直後に挿入"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from deckkit import (add_text, add_box, add_arrow, add_line, add_freeform,
                     set_notes, INK, GOLD, BLUE, VERM, TEAL, RED, GREY, LGREY)
from newpages import Builder, ppg_curve, comp_curve, WAVE_YOUNG, WAVE_MID, WAVE_OLD, WAVE_FLAT
from refs import cite

prs = Presentation('deck_s1.pptx')
S = list(prs.slides)
B = Builder(prs, S[59])          # 5.3 SDPPG の主な知見＝タイトル＋章ナビ＋出典を持つ素直な版面

PALE = "F2F2F2"

# ══════════════════════════════════════════════ 2 の直後：用語の整理（図解版）
s = B.new(S[1], '用語の整理', 1, cite(1, 25, 26))

add_box(s, 0.55, 1.95, 3.85, 1.95, fill=PALE, line=BLUE, line_w=2.25)
add_text(s, 0.75, 2.20, 3.45, 0.5, 'PPG', size=30, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s, 0.75, 2.80, 3.45, 1.0,
         ['光を当てて', '血液量の変化を測る'], size=23, color=INK, align=PP_ALIGN.CENTER,
         space_after=4)

add_arrow(s, 4.55, 2.66, 0.62, 0.55, color=GREY)

add_box(s, 5.32, 1.95, 3.85, 1.95, fill=PALE, line=VERM, line_w=2.25)
add_text(s, 5.52, 2.20, 3.45, 0.5, 'PPG 波形', size=30, bold=True, color=VERM, align=PP_ALIGN.CENTER)
add_text(s, 5.52, 2.80, 3.45, 1.0,
         ['得られた波', '＝ 容積脈波'], size=23, color=INK, align=PP_ALIGN.CENTER, space_after=4)

add_arrow(s, 9.32, 2.66, 0.62, 0.55, color=GREY)

add_box(s, 10.09, 1.95, 2.72, 1.95, fill=PALE, line=TEAL, line_w=2.25)
add_text(s, 10.24, 2.20, 2.42, 0.5, 'SpO₂', size=30, bold=True, color=TEAL, align=PP_ALIGN.CENTER)
add_text(s, 10.09, 2.80, 2.72, 1.0,
         ['波から計算した', '酸素飽和度'], size=22, color=INK, align=PP_ALIGN.CENTER, space_after=4)

# 下段：測定のイメージ
add_text(s, 0.55, 4.30, 4.2, 0.42, '① 指に光を当てる', size=24, bold=True, color=INK)
add_box(s, 0.75, 4.90, 1.15, 1.05, fill='FDE9D9', line=GOLD, line_w=1.5)
add_text(s, 0.78, 5.20, 1.1, 0.4, '発光', size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
add_arrow(s, 2.05, 5.24, 0.95, 0.38, color=GOLD)
add_box(s, 3.15, 4.90, 1.15, 1.05, fill=PALE, line=GREY, line_w=1.5)
add_text(s, 3.18, 5.20, 1.1, 0.4, '受光', size=22, bold=True, color=GREY, align=PP_ALIGN.CENTER)

add_text(s, 5.32, 4.30, 4.2, 0.42, '② 波が得られる', size=24, bold=True, color=INK)
pts = ppg_curve(5.45, 4.85, 3.55, 1.15, WAVE_YOUNG)
add_freeform(s, pts, color=VERM, width=2.5)
add_line(s, 5.45, 6.02, 9.00, 6.02, color=LGREY, width=1.25)

add_text(s, 10.09, 4.30, 3.0, 0.42, '③ 数値になる', size=24, bold=True, color=INK)
add_box(s, 10.09, 4.85, 2.72, 1.10, fill=PALE, line=TEAL, line_w=2.0)
add_text(s, 10.14, 5.10, 2.62, 0.5, '98 %', size=30, bold=True, color=TEAL, align=PP_ALIGN.CENTER)

add_text(s, 0.55, 6.35, 12.3, 0.5,
         [[('本セッションは ', {}), ('②の「波そのもの」', {'bold': True, 'color': VERM}),
           (' を読み直す話', {})]], size=26, bold=True, color=INK)

set_notes(s, """用語の整理（図解版）

・PPG（photoplethysmography, 光電式容積脈波記録法）＝ 皮膚に LED などの光を当て、
　血液量の変化による光の透過・反射量の変化を検出して脈波を測定する非侵襲的な光学的測定法。
・得られた脈波を PPG 波形、容積脈波（digital volume pulse, DVP）などと称する。
・パルスオキシメトリーは、この PPG 信号を利用して SpO₂（経皮的動脈血酸素飽和度）を
　非侵襲的に推定する測定原理。
・つまりモニタは SpO₂ という数値を出すためにすでに PPG を取り込んでいる。
　本セッションはその副産物である「波形」をもう一度読み直す作業にあたる。
出典：1. Allen 2007 ／ 25. Aoyagi 2003 ／ 26. Chan 2013""")

# ══════════════════════════════════════════ 20 の直後：PPG 基礎まとめ（図解版）
s = B.new(S[19], '1.5  PPG基礎まとめ', 1, cite(1, 2, 25, 26))

cards = [
    ('①', 'SpO₂ の正体', ['赤色光と赤外光の', '吸収の比 R を', '健常者の較正表に', '当てはめた推定値'], BLUE),
    ('②', 'PPG が測るもの', ['圧ではなく', '動脈血の', '容積の変化', '（局所の信号）'], VERM),
    ('③', '波形の成り立ち', ['前進波 ＋ 反射波', 'の合成', '反射点の場所は', '同定できない'], TEAL),
]
for i, (num, head, body, col) in enumerate(cards):
    x = 0.55 + i * 4.16
    add_box(s, x, 1.90, 3.86, 3.30, fill=PALE, line=col, line_w=2.25)
    add_text(s, x + 0.20, 2.06, 0.6, 0.45, num, size=28, bold=True, color=col)
    add_text(s, x + 0.80, 2.10, 2.90, 0.45, head, size=25, bold=True, color=col)
    add_text(s, x + 0.25, 2.75, 3.40, 2.30, body, size=23, color=INK, space_after=6)

add_box(s, 0.55, 5.50, 12.23, 1.20, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.75, 5.72, 11.8, 0.8,
         [[('では、この波形は ', {}), ('どのような因子', {'bold': True, 'color': GOLD}),
           (' に影響を受けるのか', {})],
          [('→　次章から 加齢・循環・麻酔 の順に見ていく', {'size': 24})]],
         size=27, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=4)

set_notes(s, """1.5 PPG 基礎まとめ（図解版）

① SpO₂ は赤色光と赤外光の（AC/DC）比 R を、健常人の低酸素負荷試験で作った較正表に
　 当てはめて SaO₂ を推定した値である。較正データは SaO₂ 70〜100% の範囲でしか取得できず、
　 70% 未満は外挿である。
② PPG 波形は圧ではなく、測定部位の動脈血液量（容積）の拍動性変化を見ている。
　 したがって振幅は中心大動脈圧の絶対値ではなく、局所の血管コンプライアンス・灌流・
　 局所脈圧に依存する。
③ PPG 波形は前進波と反射波の合成によって形成されるが、反射点の解剖学的な場所を
　 同定することは原理的にできない（成立するのは「経路」としての解釈だけ）。

→ 次章以降で、この波形がどのような因子に影響を受けるかを見ていく。
出典：1. Allen 2007 ／ 2. Politi 2016 ／ 25. Aoyagi 2003 ／ 26. Chan 2013""")

# ══════════════════════════════════ 49 の直後：6分類と血行動態（要点に絞った版）
s = B.new(S[48], '4.2  6分類と血行動態', 4, cite(14, note='心臓手術 15 名・190 エピソード'))

trio = [('Class I・II', '血管収縮', BLUE), ('Class III', '正常', '4E8F45'),
        ('Class IV〜VI', '血管拡張', VERM)]
for i, (cls, state, col) in enumerate(trio):
    x = 0.55 + i * 4.16
    add_box(s, x, 1.95, 3.86, 1.35, fill=col, line=None)
    add_text(s, x + 0.10, 2.14, 3.66, 0.45, cls, size=27, bold=True,
             color='FFFFFF', align=PP_ALIGN.CENTER)
    add_text(s, x + 0.10, 2.68, 3.66, 0.45, state, size=25, bold=True,
             color='FFFFFF', align=PP_ALIGN.CENTER)

add_box(s, 0.55, 3.60, 12.23, 1.45, fill=PALE, line=None)
add_text(s, 0.75, 3.82, 11.8, 0.55, '収縮期血圧との相関', size=25, bold=True,
         color=INK, align=PP_ALIGN.CENTER)
add_text(s, 0.75, 4.35, 11.8, 0.6, 'r ＝ −0.90　（P<0.0001）', size=32, bold=True,
         color=BLUE, align=PP_ALIGN.CENTER)

add_box(s, 0.55, 5.30, 12.23, 1.35, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.75, 5.52, 11.8, 0.95,
         ['低血圧の検出　感度 100%　・　特異度 97.9%',
          '高血圧の検出　感度 94.9%　・　特異度 99.2%'],
         size=26, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=4)

set_notes(s, """4.2 6分類と血行動態（要点版）

Tusman 2019：心臓手術（CABG）15 名・190 エピソード。
PPG 波形を「切痕の位置 × 振幅」で 6 クラスに分類し、観血的血行動態と対比した。

・Class I・II ＝ 血管収縮（SAP・SVR 高）
・Class III  ＝ 正常トーヌス（SAP・SVR 基準）
・Class IV〜VI ＝ 血管拡張（SAP・SVR 低）

PPG 6 分類との相関（いずれも P<0.0001）
・収縮期血圧 SAP        r = −0.90
・全身血管抵抗 SVR      r = −0.72
・血管コンプライアンス   r = +0.82

検出能
・低血圧の検出 感度 100%・特異度 97.9%
・高血圧の検出 感度 94.9%・特異度 99.2%

※前ページに SVR・コンプライアンスを含む全数値を残してあるので、
　口頭ではそちらを補って説明する。
出典：14. Tusman 2019""")

# ═══════════════════════════════ 60 の直後：SDPPG の文献をどう取捨するか（改訂）
s = B.new(S[59], '5.3  文献の取捨', 5, cite(16, 18, 19, 30, 32))

add_text(s, 0.45, 1.80, 3.30, 0.35, '文献', size=22, bold=True, color=GREY)
add_text(s, 4.35, 1.80, 6.00, 0.35, '主な知見', size=22, bold=True, color=GREY)
add_text(s, 10.60, 1.80, 2.30, 0.35, '判定', size=22, bold=True, color=GREY,
         align=PP_ALIGN.CENTER)

rows = [
    ('16. Takazawa 1998', ['薬で d/a が動く', '（AT II −0.62 ／ NTG −0.25）'], '残す', GOLD),
    ('30. Otsuka 2006', ['Framingham risk score と相関', '（b/a 0.43–0.54）'], '残す', GOLD),
    ('18. Bortolotto 2000', ['規定因子としては PWV が優位'], '残す', GOLD),
    ('32. Tabara 2016', ['臓器障害との関連は弱い', '（J-SHIPP）'], '足す', VERM),
    ('19. Hashimoto 2002', ['PWV と SDPPG は別情報'], '足す', VERM),
]
y = 2.20
for name, finding, judge, col in rows:
    add_box(s, 0.45, y, 12.43, 0.78, fill=PALE, line=None)
    add_text(s, 0.62, y + 0.20, 3.55, 0.45, name, size=22, bold=True, color=INK)
    add_text(s, 4.35, y + (0.09 if len(finding) > 1 else 0.20), 5.95, 0.62,
             finding, size=22, color=INK, space_after=1)
    add_text(s, 10.60, y + 0.16, 2.30, 0.5, judge, size=26, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    y += 0.85

add_text(s, 0.45, 6.60, 12.43, 0.35,
         '※ 引くべき文献はない（作用機序・外的妥当性・限界と役割が分かれている）',
         size=22, bold=True, color=GREY)

set_notes(s, """5.3 SDPPG 文献の位置づけ ― 足すべき文献・引くべき文献

【結論】3 件（Takazawa 1998／Bortolotto 2000／Otsuka 2006）だけでは
「SDPPG は有望」という方向に偏る。否定側の 2 件を足すのが望ましい。引くべき文献はない。

■ 残す
・16. Takazawa 1998：血管作動薬で d/a が用量依存的に変化（AT II −0.62／NTG −0.25）、
  AGI と年齢 r=0.80。d/a を「機能的な末梢血管緊張」と結びつけた原典であり外せない。
・30. Otsuka 2006（Circ J 2006;70:304-10、211 名）：SDPTG 指標が Framingham risk score と
  相関（b/a r=0.43(男)/0.54(女)、d/a r=−0.38/−0.58）。一般集団での外的妥当性を与える。
・18. Bortolotto 2000：動脈硬化あり AGI −0.093／なし −0.271（P<0.001）。ただし
  「独立規定因子としては PWV が優位」という否定的側面が主眼なので 1 行に圧縮してよい。

■ 足す（今回追加を推奨）
・32. Tabara 2016（J-SHIPP、Hypertens Res 2016;39:552-6）：b/a・d/a・AGI の強い規定因子は
  年齢・性・収縮期血圧・心拍数であり、頸動脈 IMT など臓器障害との関連は弱いか有意でない
  （b/a では β=0.069, P=0.002）。SDPPG を過大評価しないための歯止めとして必須。
・19. Hashimoto 2002：治療中高血圧 294 名。PWV と SDPTG 指標は年齢・血圧という共通因子を
  持つが相互の相関は弱く、中心部と末梢部について別の情報を与える。既に本デックの文献表
  （#19）に入っているのに、このページでは使われていなかった。

■ 引く候補
・なし。3 件はそれぞれ役割が違う（作用機序／外的妥当性／限界）ため、いずれも残す。""")

prs.save('deck_s2a.pptx')
print('stage2a ok. slides =', len(prs.slides._sldIdLst))
