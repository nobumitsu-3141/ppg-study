# -*- coding: utf-8 -*-
"""Stage 1: 全ページの出典書式統一・目次修正・ページ番号再採番・誤記修正"""
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from deckkit import (set_source, set_pageno, set_title, set_nav, get_title,
                     add_text, NS, INK, GOLD, LGREY, FONT, _force_ea)
from refs import cite
from pptx.enum.text import PP_ALIGN

prs = Presentation('deck_v66.pptx')
S = list(prs.slides)   # 元の 1..76 を index-1 でアクセス


def sl(n):
    return S[n - 1]


# ---------------------------------------------------------------- 1) 出典統一
SOURCES = {
    2:  cite(25, 26),
    5:  cite(25, 26),
    6:  cite(26),
    7:  cite(1), 8: cite(1), 9: cite(1),
    10: cite(26), 11: cite(26), 12: cite(25, 26), 13: cite(25, 26),
    14: cite(25, 26),
    15: cite(1),
    16: cite(1, 2), 17: cite(1, 2), 18: cite(1, 2),
    19: cite(1, 2, 4),
    20: cite(1, 2, 25, 26),
    22: cite(5), 23: cite(5), 24: cite(5), 25: cite(5), 26: cite(5), 27: cite(5),
    28: cite(6, 37), 29: cite(6, 37), 30: cite(6), 31: cite(6),
    32: cite(4, 5, 6), 33: cite(4, 5, 6),
    35: cite(8), 36: cite(8),
    37: cite(15, 27), 38: cite(15, 27),
    39: cite(4, 6, 27),
    40: cite(10, 12), 41: cite(10, 12),          # ← 「111. Lin 2020」を修正
    42: cite(7, 19, 23, 24),
    43: cite(4, 8, 10, 15, 23),
    45: cite(17),
    46: cite(14), 47: cite(14), 48: cite(14),
    49: cite(14, note="心臓手術 15 名・190 エピソード"),
    50: cite(29, note="全身麻酔導入 61 名"),
    52: cite(4), 53: cite(4),
    54: cite(15), 55: cite(15),
    56: cite(35, note="ICU 48 名"),               # ← 「Lee 2011」を番号付きに
    57: cite(7, 19, 23, 24),
    58: cite(22),                                  # ← 「elgendi2012」を修正
    59: cite(16, 19),                              # ← 「Hashimoto 2002」を番号付きに
    60: cite(16, 18, 30),                          # ← Otsuka 2006 に番号付与
    61: cite(4, 27, 31, 32, 33),                   # ← カタカナ表記を廃し番号付きへ
    62: cite(16, 24, 27, 29, 34),                  # ← Mdlazim2023 → 24. Md Lazim 2020
    63: cite(4, 16, 20, 28),
    65: cite(22),
    66: cite(6, 22, 36),                           # ← 「22. Wu 2010」を分離
    67: cite(38),                                  # ← DN-less の出典を明示
    68: cite(11, 41),                              # ← 【改善案…】メモを削除
}
for n, txt in SOURCES.items():
    set_source(sl(n), txt)

# 出典を持たない扉・目次・文献ページは出典欄を空に
for n in (1, 3, 4, 21, 34, 44, 51, 64, 69, 70):
    set_source(sl(n), None)

# ------------------------------------------------- 2) タイトル・章ナビの補完
# 26・27（Dawber の図版ページ）はタイトルと章ナビが無い → 21 番（2章扉）から複写
src_nav = [s for s in sl(24).shapes if 'Off-page Connector' in s.name]
import copy
for n, title in ((26, '2.1  フラミンガム研究'), (27, '2.1  フラミンガム研究')):
    tgt = sl(n)
    if get_title(tgt) is None:
        tpl = None
        for shp in sl(24).shapes:
            if shp.name in ('Title 1', 'タイトル 1'):
                tpl = shp
                break
        tgt.shapes._spTree.insert(2, copy.deepcopy(tpl._element))
    if not [s for s in tgt.shapes if 'Off-page Connector' in s.name]:
        for chip in src_nav:
            tgt.shapes._spTree.append(copy.deepcopy(chip._element))
    set_title(tgt, title)
    set_nav(tgt, 2)

# 2.3 のタイトルゆれを統一（S32「意味」／S33「理由」→「理由」に統一）
set_title(sl(32), '2.3  切痕平滑化の理由')

# 3.7 → 3.6（3.6 が実在しないための番号ずれを解消）
set_title(sl(43), '3.6  波形変化まとめ')

# --------------------------------------------------------------- 3) 目次の修正
from stage1helpers import CH, fill_divider, fill_menu

DIVIDERS = {4: 1, 21: 2, 34: 3, 44: 4, 51: 5, 64: 7}
for n, ch in DIVIDERS.items():
    fill_divider(sl(n), ch)
    set_title(sl(n), f'{ch}. {CH[ch][0]}')
    set_nav(sl(n), ch)

fill_menu(sl(3))
set_nav(sl(3), None)

# ------------------------------------------------------- 4) 個別の誤記修正
# S56：Lee QY 2011 の対象人数と結論の向きを修正（段落単位で置換）
def replace_in_slide(slide, pairs):
    hit = 0
    for shp in slide.shapes:
        if not shp.has_text_frame:
            continue
        for para in shp.text_frame.paragraphs:
            runs = para.runs
            if not runs:
                continue
            whole = ''.join(r.text for r in runs)
            new = whole
            for old, rep in pairs:
                new = new.replace(old, rep)
            if new != whole:
                runs[0].text = new
                for r in runs[1:]:
                    r.text = ''
                hit += 1
    return hit


n = replace_in_slide(sl(56), [
    ('ICU 患者 64人', 'ICU 患者 48人'),
    ('ICU 患者 64 人', 'ICU 患者 48人'),
    ('単独の特徴量では識別能力が高い', '単独の特徴量では識別能力が低い'),
    ('関連性はまずまず（κ=0.33）', 'RI 単独 κ=0.33／最良の特徴量セット κ=0.57'),
])
print('S56 replacements:', n)

prs.save('deck_s1.pptx')
print('stage1 saved:', len(prs.slides.__iter__.__self__._sldIdLst), 'slides')
