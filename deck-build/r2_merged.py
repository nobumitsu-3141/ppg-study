# -*- coding: utf-8 -*-
"""2・3・4 章を「波形への影響」1章にまとめた別バージョンを作る（3階層番号）"""
import re
from pptx import Presentation
from deckkit import set_title, set_pageno, drop_shapes, slide_index, NS
from r2_common import rewrite_nav, fill_divider, fill_menu, title_text, CHIPS
import r2_common

prs = Presentation('deck_v69.pptx')
RNS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'


def SL(n):
    return list(prs.slides)[n - 1]


# ══════════ 新しい章立て（2〜4 → 2章に統合、以降を1つ繰り上げ）
MCH = {
    1: ('PPGの基礎', ['1.1　SpO₂の測定原理', '1.2　PPGが測るもの',
                      '1.3　前進波の成因', '1.4　反射波の成因', '1.5　PPG基礎まとめ']),
    2: ('波形への影響', ['2.1　加齢の影響', '2.2　循環の変化', '2.3　麻酔中の変化']),
    3: ('波形の定量化', ['3.1　Stiffness Index (SI)', '3.2　Reflection Index (RI)',
                        '3.3　加速度脈波 SDPPG', '3.4　定量化指標まとめ']),
    4: ('限界', ['4.1　自動解析の壁', '4.2　極端な動脈硬化', '4.3　DN-less 信号対策']),
    5: ('PDA（波形分解）', ['5.1　PDA とは', '5.2　分解のしくみ', '5.3　当てはめ方',
                           '5.4　得られる指標', '5.5　重要文献',
                           '5.6　限界', '5.7　残された問い']),
}
MCHIPS = ['基礎', '影響', '定量化', '限界', 'PDA']

# ══════════ 見出し番号の対応表（旧 → 新）
#  2章 加齢 → 2.1.x ／ 3章 循環 → 2.2.x ／ 4章 麻酔 → 2.3.x
SUBMAP = {
    '2.1': '2.1.1', '2.2': '2.1.2', '2.3': '2.1.3',
    '3.1': '2.2.1', '3.2': '2.2.2', '3.3': '2.2.3',
    '3.4': '2.2.4', '3.5': '2.2.5', '3.6': '2.2.6',
    '4.1': '2.3.1', '4.2': '2.3.2', '4.3': '2.3.3',
    '5.1': '3.1', '5.2': '3.2', '5.3': '3.3', '5.4': '3.4',
    '6.1': '4.1', '6.2': '4.2', '6.3': '4.3',
    '7.1': '5.1', '7.2': '5.2', '7.3': '5.3', '7.4': '5.4',
    '7.5': '5.5', '7.6': '5.6', '7.7': '5.7',
}
CHMAP = {'5.': '3.', '6.': '4.', '7.': '5.'}

# ══════════ 3・4 章の扉を削除（2章に吸収）
kill = [i for i, s in enumerate(prs.slides, 1)
        if title_text(s) in ('3. 波形への影響：循環', '4. 麻酔中の波形変化')]
for i in sorted(kill, reverse=True):
    el = list(prs.slides._sldIdLst)[i - 1]
    prs.part.drop_rel(el.get('{%s}id' % RNS))
    prs.slides._sldIdLst.remove(el)
print('統合により削除した章扉:', kill)

# ══════════ タイトルの番号を振り直し
for s in prs.slides:
    t = title_text(s)
    if not t:
        continue
    m = re.match(r'^(\d+\.\d+)(\.?\s+.*)$', t)
    if m and m.group(1) in SUBMAP:
        set_title(s, SUBMAP[m.group(1)] + m.group(2))
        continue
    m = re.match(r'^(\d+\.)(\s.*)$', t)
    if m and m.group(1) in CHMAP:
        set_title(s, CHMAP[m.group(1)] + m.group(2))

set_title([s for s in prs.slides if title_text(s) == '2. 波形への影響：加齢'][0],
          '2. 波形への影響')

# ══════════ 章ナビ帯を 5 チップに作り直す
r2_common.CHIPS = MCHIPS


def chapter_of_title(t):
    if not t:
        return None
    head = t.split()[0].rstrip('.').split('.')[0]
    return int(head) if head.isdigit() else None


for s in prs.slides:
    t = title_text(s)
    ch = None if (t in ('メニュー', '') or t.startswith('参考文献')) else chapter_of_title(t)
    chips = sorted([sh for sh in s.shapes if 'Off-page Connector' in sh.name],
                   key=lambda sh: sh.left or 0)
    # 右詰めのまま 5 個に減らす（余った 2 個は削除し、残りを右へ寄せる）
    if len(chips) > len(MCHIPS):
        for sh in chips[:len(chips) - len(MCHIPS)]:
            sh._element.getparent().remove(sh._element)
        chips = chips[len(chips) - len(MCHIPS):]
    rewrite_nav(s, ch)

# ══════════ 目次・章扉
fill_menu([s for s in prs.slides if title_text(s) == 'メニュー'][0], MCH)
for s in prs.slides:
    t = title_text(s)
    m = re.match(r'^(\d+)\.\s', t)
    if m and t.count('.') == 1:
        ch = int(m.group(1))
        if ch in MCH:
            fill_divider(s, ch, MCH, sub_size=2200, main_size=2400)
            set_title(s, f'{ch}. {MCH[ch][0]}')

# ══════════ ページ番号の再採番
vis = [s for s in prs.slides if s.element.get('show') != '0']
logical, prev = [], None
for s in vis:
    key = title_text(s) or f'(fig{id(s)})'
    if prev is not None and key == prev and not key.startswith('(fig'):
        logical[-1].append(s)
    else:
        logical.append([s])
    prev = key
total = len(logical)
for k, group in enumerate(logical, start=1):
    for s in group:
        set_pageno(s, None if k == 1 else f'{k}/{total}')
for s in prs.slides:
    if s.element.get('show') == '0':
        set_pageno(s, None)

prs.save('deck_v69_merged.pptx')
print('統合版 ok:', len(prs.slides._sldIdLst), 'slides ／ 論理ページ', total)
