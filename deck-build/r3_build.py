# -*- coding: utf-8 -*-
"""R3: 2.x の節目次ページを全ての節の頭に挿入し、全目次ページの版面を広げる"""
import re
import sys
from pptx import Presentation
from pptx.util import Emu
from deckkit import (set_title, set_source, set_pageno, clone_slide,
                     slide_index, drop_shapes, INK, GOLD, GREY, LGREY)
from r2_common import rewrite_nav, title_text, find_agenda, CH
import r2_common
from r3_toc import (build_lines, render_agenda, pick_band, collect_sub2,
                    render_column, pick_band_col, LCOL_X, LCOL_W, RCOL_X, RCOL_W,
                    BOX_Y, BOX_H)

SRC, OUT, MODE = sys.argv[1], sys.argv[2], sys.argv[3]   # MODE: 'merged' | 'plain'
prs = Presentation(SRC)

MCH = {
    1: ('PPGの基礎', ['1.1　SpO₂の測定原理', '1.2　PPGが測るもの',
                      '1.3　前進波の成因', '1.4　反射波の成因', '1.5　PPG基礎まとめ']),
    2: ('波形への影響', ['2.1　加齢の影響', '2.2　循環の変化', '2.3　麻酔中の変化']),
    3: ('波形の定量化', ['3.1　Stiffness Index (SI)', '3.2　Reflection Index (RI)',
                        '3.3　加速度脈波 SDPPG', '3.4　定量化指標まとめ']),
    4: ('限界', ['4.1　自動解析の壁', '4.2　極端な動脈硬化', '4.3　DN-less 信号対策']),
    5: ('PDA（波形分解）', ['5.1　PDA とは', '5.2　分解のしくみ', '5.3　当てはめ方',
                           '5.4　得られる指標', '5.5　重要文献',
                           '5.6　限界', '5.7　残された問い']),
}
if MODE == 'merged':
    CHAPTERS, CHIPS = MCH, ['基礎', '影響', '定量化', '限界', 'PDA']
else:
    CHAPTERS, CHIPS = CH, ['基礎', '加齢', '循環', '麻酔', '定量化', '限界', 'PDA']
r2_common.CHIPS = CHIPS

SUB2 = collect_sub2(prs, title_text)
print('3階層見出し:', {k: [n for n, _ in v] for k, v in SUB2.items()})


def is_divider(s):
    t = title_text(s)
    return bool(re.match(r'^\d+\.\s', t)) and t.count('.') == 1


def ch_of(s):
    m = re.match(r'^(\d+)\.', title_text(s))
    return int(m.group(1)) if m else None


# ══════════════ 1) 既存の「試作」節目次ページを見分ける
#   章扉と同じタイトルで、本文に 3 階層番号を含むページ
def has_lv2(s):
    ag = find_agenda(s)
    return ag is not None and re.search(r'\d+\.\d+\.\d+', ag.text_frame.text) is not None


trial = [i for i, s in enumerate(prs.slides, 1) if is_divider(s) and has_lv2(s)]
print('既存の節目次ページ:', trial)

# ══════════════ 2) 節目次ページを各節の先頭に用意する
SEC_PAGES = set()
if MODE == 'merged':
    tpl_idx = trial[0] if trial else [i for i, s in enumerate(prs.slides, 1)
                                      if is_divider(s) and ch_of(s) == 2][0]
    tpl = list(prs.slides)[tpl_idx - 1]

    for parent in sorted(SUB2, key=lambda p: [int(x) for x in p.split('.')]):
        ch = int(parent.split('.')[0])
        # その節の最初のスライド
        first = None
        for i, s in enumerate(prs.slides, 1):
            if title_text(s).startswith(parent + '.'):
                first = i
                break
        if first is None:
            continue
        prev = list(prs.slides)[first - 2]
        if is_divider(prev) and has_lv2(prev):
            target = prev                       # すでにある（23ページの試作など）
        else:
            target = clone_slide(prs, tpl, insert_after_idx=first - 2)
        set_title(target, f'{ch}. {CHAPTERS[ch][0]}')
        set_source(target, None)
        set_pageno(target, None)
        rewrite_nav(target, ch)
        # 左欄：章と節の一覧（該当節を強調）
        left = []
        for i in sorted(CHAPTERS):
            left.append((0, f'{i}.', CHAPTERS[i][0]))
            if i != ch:
                continue
            for entry in CHAPTERS[i][1]:
                num, _, label = entry.partition('　')
                cur = (num == parent)
                left.append((1, num, label,
                             {'on': True, 'bold': cur,
                              'color': INK if cur else GREY,
                              'num_color': GOLD if cur else LGREY}))
        # 右欄：この節の中身
        sec_label = dict((e.partition('　')[0], e.partition('　')[2])
                         for e in CHAPTERS[ch][1])[parent]
        right = [(0, parent, sec_label, {'on': True, 'bold': True})]
        right += [(1, n2, l2, {'on': True, 'indent': ''}) for n2, l2 in SUB2[parent]]
        m1, s1, l1 = pick_band_col(left)
        m2, s2, l2_ = pick_band_col(right)
        m, sb, lsp = min(m1, m2), min(s1, s2), min(l1, l2_)
        render_column(find_agenda(target), left, LCOL_X, LCOL_W, m, sb, lsp,
                      active_ch=ch, anchor_middle=False)
        # 右欄は新しいテキストボックスで
        drop_shapes(target, lambda sh: sh.name == 'SubToc')
        box = target.shapes.add_textbox(
            Emu(int(RCOL_X * 914400)), Emu(int(BOX_Y * 914400)),
            Emu(int(RCOL_W * 914400)), Emu(int(BOX_H * 914400)))
        box.name = 'SubToc'
        render_column(box, right, RCOL_X, RCOL_W, m, sb, lsp, anchor_middle=False)
        SEC_PAGES.add(target.slide_id)
        print(f'  節目次 {parent}: 左 {len(left)} 行 / 右 {len(right)} 行 '
              f'/ {m}pt・{sb}pt・行間{lsp}')

# ══════════════ 3) 章扉とメニューの版面を作り直す
from lxml import etree
from deckkit import NS, GOLD, INK

for s in prs.slides:
    if is_divider(s) and s.slide_id not in SEC_PAGES:
        ch = ch_of(s)
        if ch not in CHAPTERS:
            continue
        # 左欄＝章の一覧、右欄＝その章の節。2欄にして文字を大きく取る
        left = [(0, f'{i}.', CHAPTERS[i][0]) for i in sorted(CHAPTERS)]
        right = []
        for entry in CHAPTERS[ch][1]:
            num, _, label = entry.partition('　')
            right.append((1, num, label, {'on': True, 'indent': ''}))
        m1, s1, l1 = pick_band_col(left)
        m2, s2, l2_ = pick_band_col(right)
        m, sb, lsp = min(m1, m2), min(s1, s2), min(l1, l2_)
        render_column(find_agenda(s), left, LCOL_X, LCOL_W, m, sb, lsp,
                      active_ch=ch, anchor_middle=False)
        drop_shapes(s, lambda sh: sh.name == 'SubToc')
        box = s.shapes.add_textbox(
            Emu(int(RCOL_X * 914400)), Emu(int(BOX_Y * 914400)),
            Emu(int(RCOL_W * 914400)), Emu(int(BOX_H * 914400)))
        box.name = 'SubToc'
        render_column(box, right, RCOL_X, RCOL_W, m, sb, lsp, anchor_middle=False)

# メニューは全章を強調表示（番号=金／見出し=黒・太字）
menu = [s for s in prs.slides if title_text(s) == 'メニュー']
if menu:
    ag = find_agenda(menu[0])
    lines = [(0, f'{i}.', CHAPTERS[i][0]) for i in sorted(CHAPTERS)]
    render_agenda(ag, lines, active_ch=None)
    for p in ag.text_frame._txBody.findall('a:p', NS):
        for k, r in enumerate(p.findall('a:r', NS)):
            clr = r.find('a:rPr/a:solidFill/a:srgbClr', NS)
            if clr is not None:
                clr.set('val', GOLD if k == 0 else INK)
            r.find('a:rPr', NS).set('b', '1')

# ══════════════ 4) ページ番号の再採番
vis = [s for s in prs.slides if s.element.get('show') != '0']
logical, prev = [], None
for k, s in enumerate(vis):
    key = title_text(s) or f'(fig{k})'
    is_sec = s.slide_id in SEC_PAGES
    if prev is not None and key == prev and not key.startswith('(fig') and not is_sec:
        logical[-1].append(s)
    else:
        logical.append([s])
    prev = f'(sec{k})' if is_sec else key
total = len(logical)
for k, group in enumerate(logical, start=1):
    for s in group:
        set_pageno(s, None if k == 1 else f'{k}/{total}')
for s in prs.slides:
    if s.element.get('show') == '0':
        set_pageno(s, None)

prs.save(OUT)
print(f'{OUT}: {len(prs.slides._sldIdLst)} slides / 論理ページ {total}')
