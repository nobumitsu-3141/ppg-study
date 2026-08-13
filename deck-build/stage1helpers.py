# -*- coding: utf-8 -*-
"""章立て（目次）の定義と、章扉／メニューの本文差し替え"""
from lxml import etree
from pptx.util import Emu
from deckkit import NS, INK, GOLD, LGREY, FONT

CH = {
    1: ('PPGの基礎', ['1.1　SpO₂の測定原理', '1.2　PPGが測るもの',
                      '1.3　前進波の成因', '1.4　反射波の成因', '1.5　PPG基礎まとめ']),
    2: ('波形への影響：加齢', ['2.1　フラミンガム研究', '2.2　UK Biobank',
                              '2.3　切痕平滑化の理由']),
    3: ('波形への影響：循環', ['3.1　前負荷', '3.2　後負荷（末梢トーヌス）',
                              '3.3　後負荷の違い', '3.4　心収縮力', '3.5　心拍数',
                              '3.6　波形変化まとめ']),
    4: ('麻酔中の波形変化', ['4.1　トーヌスと切痕', '4.2　波形の6分類',
                            '4.3　導入時の低血圧']),
    5: ('波形の定量化', ['5.1　Stiffness Index (SI)', '5.2　Reflection Index (RI)',
                        '5.3　加速度脈波 SDPPG', '5.4　定量化指標まとめ']),
    6: ('解析まとめ', ['6.1　二軸で読む']),
    7: ('限界', ['7.1　自動解析の壁', '7.2　波形分解 PDA']),
}


def find_agenda(slide):
    cands = []
    for shp in slide.shapes:
        if not shp.has_text_frame or shp.top is None:
            continue
        if 'Off-page Connector' in shp.name or shp.name in ('Title 1', 'タイトル 1'):
            continue
        if shp.top < 1.3 * 914400 or shp.top > 2.4 * 914400:
            continue
        n = len([p for p in shp.text_frame.paragraphs if p.text.strip()])
        cands.append((n, shp))
    if not cands:
        return None
    cands.sort(key=lambda t: -t[0])
    return cands[0][1]


def _para(txBody, parts, size):
    p = etree.SubElement(txBody, '{%s}p' % NS['a'])
    for txt, col, bold in parts:
        r = etree.SubElement(p, '{%s}r' % NS['a'])
        rPr = etree.SubElement(r, '{%s}rPr' % NS['a'])
        rPr.set('sz', str(size))
        rPr.set('b', '1' if bold else '0')
        sf = etree.SubElement(rPr, '{%s}solidFill' % NS['a'])
        c = etree.SubElement(sf, '{%s}srgbClr' % NS['a'])
        c.set('val', col)
        for tag in ('latin', 'ea', 'cs'):
            e = etree.SubElement(rPr, '{%s}%s' % (NS['a'], tag))
            e.set('typeface', FONT)
        t = etree.SubElement(r, '{%s}t' % NS['a'])
        t.text = txt


def fill_divider(slide, active):
    """章扉の目次テキストを、実在するスライド構成に合わせて再構成"""
    target = find_agenda(slide)
    if target is None:
        print('  ! agenda box not found')
        return
    target.left, target.top = Emu(int(0.94 * 914400)), Emu(int(1.63 * 914400))
    target.width, target.height = Emu(int(11.57 * 914400)), Emu(int(5.47 * 914400))
    txBody = target.text_frame._txBody
    for p in txBody.findall('a:p', NS):
        txBody.remove(p)
    for i in range(1, 8):
        on = (i == active)
        _para(txBody, [(f'{i}. ', GOLD if on else LGREY, True),
                       (CH[i][0], INK if on else LGREY, on)], 2400)
        if on:
            for sub in CH[i][1]:
                num, _, label = sub.partition('　')
                _para(txBody, [('　' + num + ' ', GOLD, True), (label, INK, False)], 2200)


def fill_menu(slide):
    target = find_agenda(slide)
    if target is None:
        return
    txBody = target.text_frame._txBody
    for p in txBody.findall('a:p', NS):
        txBody.remove(p)
    for i in range(1, 8):
        _para(txBody, [(f'{i}.　', GOLD, True), (CH[i][0], INK, True)], 2600)
