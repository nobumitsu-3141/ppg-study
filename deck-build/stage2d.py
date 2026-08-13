# -*- coding: utf-8 -*-
"""Stage 2d: 68（7.2 PDA）の直後に PDA 解説スライド群＋重要文献まとめを挿入"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from deckkit import (add_text, add_box, add_arrow, add_line, add_freeform,
                     set_notes, INK, GOLD, BLUE, VERM, TEAL, RED, GREY, LGREY)
from newpages import (Builder, ppg_curve, comp_curve, WAVE_YOUNG, WAVE_OLD, WAVE_FLAT)
from refs import cite

prs = Presentation('deck_s2c.pptx')
T = {}
for s in prs.slides:
    for sh in s.shapes:
        if sh.name in ('Title 1', 'タイトル 1') and sh.has_text_frame:
            T.setdefault(sh.text_frame.text.strip(), []).append(s)

B = Builder(prs, T['5.3  SDPPG の主な知見'][0])
PALE = "F2F2F2"
cur = T['7.2  PDA'][0]


def nxt(title, source, notes=None):
    global cur
    cur = B.new(cur, title, 7, source, notes)
    return cur


# ══════════════════════════════════════════════ 7.2-1 PDA とは
s = nxt('7.2  PDA という考え方', cite(39, 40, 41))
COMPS = [(0.22, 0.085, 1.00), (0.55, 0.125, 0.46)]
add_text(s, 0.55, 1.85, 6.0, 0.42, '見えている波（実測）', size=24, bold=True, color=INK)
add_freeform(s, ppg_curve(0.60, 2.30, 5.30, 1.60, COMPS), color=INK, width=3.0)
add_line(s, 0.60, 3.92, 5.90, 3.92, color=LGREY, width=1.25)

add_text(s, 0.55, 4.20, 6.0, 0.42, '＝ 2つの山の足し算とみなす', size=24, bold=True, color=INK)
add_freeform(s, comp_curve(0.60, 4.68, 5.30, 1.55, COMPS, COMPS[0]), color=BLUE,
             width=2.5, dash='dash')
add_freeform(s, comp_curve(0.60, 4.68, 5.30, 1.55, COMPS, COMPS[1]), color=VERM,
             width=2.5, dash='dash')
add_line(s, 0.60, 6.25, 5.90, 6.25, color=LGREY, width=1.25)
add_text(s, 0.60, 6.34, 2.6, 0.4, '前進波', size=22, bold=True, color=BLUE)
add_text(s, 3.30, 6.34, 2.6, 0.4, '反射波', size=22, bold=True, color=VERM)

add_box(s, 6.75, 1.85, 6.10, 2.10, fill=PALE, line=RED, line_w=2.25)
add_text(s, 6.92, 2.05, 5.86, 1.75,
         [[('これまでの壁', {'bold': True, 'color': RED, 'size': 25})],
          [('P2 が「点」として見えない波形では', {'size': 23})],
          [('SI も RI も計算できない', {'size': 23, 'bold': True})]],
         size=23, color=INK, space_after=7)

add_box(s, 6.75, 4.20, 6.10, 2.30, fill=PALE, line=GOLD, line_w=2.25)
add_text(s, 6.92, 4.42, 5.86, 1.95,
         [[('PDA', {'bold': True, 'color': GOLD, 'size': 25}),
           ('（脈波分解解析）', {'size': 21})],
          [('波形を山の重ね合わせに当てはめ', {'size': 23})],
          [('反射波の 時刻と高さ を推定する', {'size': 23, 'bold': True})],
          [('→ 切痕が消えていても値が出せる', {'size': 23, 'color': GOLD, 'bold': True})]],
         size=23, color=INK, space_after=6)

set_notes(s, """7.2 PDA という考え方

PDA（pulse decomposition analysis, 脈波分解解析）は、PPG の一拍を
「複数の山（基底関数）の線形和」としてモデル化し、そのパラメータを最小二乗などで
当てはめる手法である。ちょうど複数の音源からの音波が重なって 1 つの波形になるのと同じで、
心臓から来た前進波と、末梢から戻ってきた反射波の重ね合わせとして PPG を説明する。

■ なぜ必要か
従来の SI（＝身長 ÷ ΔT）も RI（＝P2/P1）も、「P2 という点」の時刻と高さを必要とする。
下降脚が平滑化して P2 が局所最大として存在しない波形（DN-less）では原理的に計算できず、
検出できた症例だけを解析する選択バイアスが生じる。
PDA は点を探さずモデルを当てはめるので、切痕が消えていても
反射成分の時刻 τ と振幅 a を数値として返せる。

■ 得られる量
・T1_2（前進波と反射波の時間差）→ SI に相当
・R1_2（両者の振幅比）→ RI に相当
・LVET（左室駆出時間）も導出できる

出典：39. Rubins 2008 ／ 40. Goswami 2010 ／ 41. Couceiro 2015""")

# ══════════════════════════════════════════════ 7.2-2 手順
s = nxt('7.2  PDA の手順', cite(41, 42, 43))
steps = [('①', '波を切り出す', ['一拍ごとに', '取り出して', '正規化する'], BLUE),
         ('②', '山の形を決める', ['ガウス関数など', '基底関数の', '種類と個数'], TEAL),
         ('③', '当てはめる', ['最小二乗などで', '山の位置・高さ', '幅を推定する'], GOLD),
         ('④', '指標を作る', ['山どうしの', '時間差と振幅比が', 'SI・RI になる'], VERM)]
for i, (num, head, body, col) in enumerate(steps):
    x = 0.45 + i * 3.16
    add_box(s, x, 1.95, 3.02, 3.55, fill=PALE, line=col, line_w=2.25)
    add_text(s, x + 0.15, 2.12, 2.72, 0.48, num, size=28, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 2.70, 2.72, 0.48, head, size=24, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.10, 3.35, 2.82, 1.95, body, size=22, color=INK,
             align=PP_ALIGN.CENTER, space_after=6)
    if i < 3:
        add_arrow(s, x + 3.06, 3.35, 0.10, 0.32, color=GREY)

add_box(s, 0.45, 5.75, 12.42, 1.10, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.65, 5.95, 12.02, 0.75,
         ['むずかしいのは ② ： 「どの形の山を何個使うか」が決まっていない',
          '同じ波形でも、選び方しだいで答えが変わる'],
         size=25, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=4)

set_notes(s, """7.2 PDA の手順

① 一拍ごとに切り出し、幅と振幅を正規化する。
② 基底関数（山の形）の種類と個数を決める。ここが本手法の急所。
   Gaussian／Gamma／skewed-Gaussian／Rayleigh／log-normal、カーネル数は 2〜5 と
   研究ごとに割れている。
③ 重み付き最小二乗などでパラメータ（位置・高さ・幅）を推定する。
   初期値の選び方で解が動くことが知られている。
④ 得られた山どうしの時間差 T1_2 と振幅比 R1_2 を、それぞれ SI・RI に対応させる。

■ ② が決着していないことの意味
Tigges 2017 は 7,805 拍を対象に 4 種の基底関数とモデル次数を総当たりし、AICc で選択した。
最頻の最良モデル（Gamma 3 個）でも全体の 28.1%、次点（Rayleigh 2 個）が 14.4% で、
両者を足しても半数に満たない。単一の最良モデルは存在しない。
これは「分解」という逆問題が一意に解けていないことの直接の表れである。

出典：41. Couceiro 2015 ／ 42. Tigges 2017 ／ 43. Fleischhauer 2020""")

# ══════════════════════════════════════════════ 7.2-3 基底関数
s = nxt('7.2  どんな山を使うか', cite(42, 43, 44))
import math


def shape_curve(x0, y0, w, h, f, n=140):
    vals = [f(i / n) for i in range(n + 1)]
    vmax = max(vals) or 1.0
    return [(x0 + w * (i / n), y0 + h - h * (v / vmax)) for i, v in enumerate(vals)]


def g(t):
    return math.exp(-((t - 0.5) ** 2) / (2 * 0.15 ** 2))


def gam(t):
    return (t ** 2.2) * math.exp(-6.0 * t) if t > 0 else 0.0


def skew(t):
    z = (t - 0.30) / 0.22
    return math.exp(-z * z / 2) * (1 + math.erf(3.5 * z / math.sqrt(2)))


kinds = [('Gaussian', g, ['左右対称の釣鐘', '扱いやすいが', '実波形は非対称'], BLUE),
         ('Gamma', gam, ['立ち上がりが急で', '下りが緩やか', '実波形に近い'], TEAL),
         ('skewed-Gaussian', skew, ['対称の山に歪みを', '1つ足しただけ', '意味が保たれる'], VERM)]
for i, (name, f, body, col) in enumerate(kinds):
    x = 0.60 + i * 4.20
    add_text(s, x, 1.90, 3.70, 0.42, name, size=25, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_freeform(s, shape_curve(x + 0.25, 2.42, 3.20, 1.45, f), color=col, width=2.75)
    add_line(s, x + 0.25, 3.89, x + 3.45, 3.89, color=LGREY, width=1.25)
    add_text(s, x + 0.10, 4.05, 3.50, 1.45, body, size=22, color=INK,
             align=PP_ALIGN.CENTER, space_after=5)

add_box(s, 0.45, 5.55, 12.42, 1.30, fill=PALE, line=RED, line_w=2.25)
add_text(s, 0.65, 5.75, 12.02, 0.95,
         ['7,805 拍で総当たりしても正解の形は決まらない',
          '最頻の Gamma 3 個でも 28.1% どまり（Tigges 2017）'],
         size=25, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=4)

set_notes(s, """7.2 どんな山（基底関数）を使うか

・Gaussian（ガウス関数）＝左右対称の釣鐘型。数学的に扱いやすいが、
　実際の脈波は「立ち上がりが急で、下りが緩やか」で非対称なので、当てはまりに限界がある。
・Gamma 分布＝非対称な減衰を表現できる。おもりをつけたバネが摩擦で徐々に振動を落とす
　ように、片側だけゆっくり戻る形。Fleischhauer 2020 では Gamma と Gaussian の組合せが
　最良で、カーネル 2 個が最もノイズ・体動に頑健（評価指標が 14.09% 改善）だった。
・skewed-Gaussian＝通常のガウス関数に歪みパラメータを 1 つ足したもの。
　対称の山を土台にしているので頂点の位置や幅といった物理的意味が保たれやすい。
　Basso 2024 は MIMIC-III の 8,000 拍で Gamma-Gaussian より有意に高精度、かつ
　ランダム初期値に対する感度が低く一貫して頑健だと報告した
　（裏を返せば従来モデルは初期値で答えが動く）。
・ほかに Rayleigh 関数、log-normal＋Gaussian（Sološenko 2017）なども使われている。

Tigges 2017 の総当たり（7,805 拍・AICc）では、最頻の最良モデル（スケール付き Gamma 3 個）
でも 28.1%、次点の Rayleigh 2 個が 14.4%。単一の最良モデルは存在しない。

出典：42. Tigges 2017 ／ 43. Fleischhauer 2020 ／ 44. Basso 2024""")


# ══════════════════════════════════════════ 重要文献まとめ（共通レイアウト）
def paper(title, num, headline, pico, takeaway, color=GOLD, extra_note=''):
    s = nxt(title, cite(num))
    add_box(s, 0.45, 1.82, 12.43, 1.00, fill=PALE, line=None)
    add_text(s, 0.65, 2.06, 12.03, 0.55, headline, size=25, bold=True, color=color)
    y = 3.08
    labels = ['対象', '方法', '結果']
    for lab, body in zip(labels, pico):
        n = 1 if isinstance(body, str) else len(body)
        add_text(s, 0.50, y, 1.30, 0.45, lab, size=22, bold=True, color=GREY)
        add_text(s, 1.90, y, 11.00, 0.45 * n, body, size=22, color=INK, space_after=2)
        y += 0.52 + 0.42 * (n - 1)
    add_box(s, 0.45, 5.92, 12.43, 1.00, fill='FFF3D6', line=GOLD, line_w=2.25)
    add_text(s, 0.65, 6.18, 12.03, 0.55, takeaway, size=24, bold=True, color=INK,
             align=PP_ALIGN.CENTER)
    from refs import REFS
    set_notes(s, f"{title}\n\n{REFS[num]}\n\n{extra_note}")
    return s


paper('7.2  Rubins 2008', 39,
      '① 「点を探す」から「山を当てはめる」へ',
      ['健常者 40 名。指と耳の PPG を同時記録',
       ['収縮期波と拡張期波を、それぞれ 2 つの',
        'ガウス関数の和で当てはめ、RI と AI を算出'],
       '従来の微分法と同等に波形解析ができた'],
      'ガウス分解から RI を出す発想は 2008 年に既出', BLUE,
      extra_note="""比較対照は同一波形から微分法で求めた同じ指標。
「分解由来の RI」の最初期の報告として位置づけられる。
PMID 18855034""")

paper('7.2  Goswami 2010', 40,
      '② Rayleigh 2 波合成モデルで SI と RI を同時に導出',
      ['健常者と治療中の高血圧者の指尖容積脈波 113 信号',
       ['Rayleigh 関数による 2 波合成（TPS）モデル',
        'RI・SI・脈波伝播速度・立ち上がり遅延を導出'],
       '従来のランドマーク法・微分法とよく一致した'],
      '分解した反射波の時刻から SI を出すのも既出', BLUE,
      extra_note="""新指標 differential pulse spread（DPS）も提案している。
基底関数がガウス関数とは限らないことを早くから示した例でもある。
PMID 20734136""")

paper('7.2  Couceiro 2012', 11,
      '③ 多ガウス分解から左室駆出時間 LVET を取り出す',
      ['心エコーによる参照値が得られた被験者（人数は抄録に記載なし）',
       ['収縮期相・拡張期相をガウス関数でモデル化し、',
        '3 次微分から駆出の開始・終了を決めて LVET を推定'],
       ['推定誤差 15.84 ± 13.56 ms（既存法 23.01 ± 14.60）',
        '心エコーとの相関は同等（0.73 vs 0.75）']],
      '目視では取り出せない生理的イベントが出てくる', TEAL,
      extra_note="""比較対照は Chan らの既存 LVET 推定アルゴリズムと心エコー参照値。
PMID 23366792""")

paper('7.2  Couceiro 2015', 41,
      '④ 5ガウス分解から血圧・血管抵抗まで対比した到達点',
      [['LVET の検証は健常＋心血管疾患 68 名',
        '血圧・血管抵抗の対比は血行動態不安定な 43 名'],
       ['5 つのガウス関数で分解し、SI・RI に加えて',
        'T1_2（時間差）・R1_2（振幅比）など 6 個を算出'],
       ['LVET 絶対誤差 15.41 ± 13.66 ms（ρ=0.78）',
        'R1_2 は全参照値と低相関。最高でも ρ=0.45']],
      '振幅比は当てにならない。時間側でも ρ=0.45', VERM,
      extra_note="""本デックにとって最も重要な 1 本。
「分解由来の SI・RI を血圧・総末梢血管抵抗係数と対比する」という設問は
2015 年に既に実施済みであり、しかも成績が芳しくない。
失神を経験した患者では SI と収縮期血圧・平均血圧の一致が最高で 0.57。
→ 主要評価項目は振幅側ではなく時間側に置くべき、という設計上の含意。
PMID 26235798""")

paper('7.2  Tigges 2017', 42,
      '⑤ モデル選択は決着していない ― 逆問題が一意に解けない',
      ['ほとんどの波形形態を含む指尖容積脈波 7,805 拍',
       ['4 種類の基底関数とモデル次数を総当たりし、',
        '補正赤池情報量規準（AICc）で最良モデルを選択'],
       ['最頻の最良モデル（Gamma 3 個）でも 28.1%',
        '次点の Rayleigh 2 個が 14.4%。足しても半数未満']],
      '「どの山を何個使うか」はデータからは決まらない', RED,
      extra_note="""単一の最良モデルは存在しない。
これは分解という逆問題が一意に解けていないことの直接の根拠であり、
分解成分に「前進波」「反射波」という部位名を割り当てること自体の弱さにもつながる。
PMID 29060777""")

paper('7.2  Fleischhauer', 43,
      '⑥ Gamma ＋ Gaussian、カーネル 2 個が最も頑健',
      ['合成脈波と、PPG イメージングの実測データ',
       ['各種の脈波分解アルゴリズムを実装し、',
        'カーネルの種類と個数を変えて耐性を比較'],
       ['Gamma と Gaussian の組合せが他を上回った',
        'カーネル 2 個が最も頑健で、指標が 14.09% 改善']],
      '複雑にするほど良いわけではない ― 2 個で十分', TEAL,
      extra_note="""初期値設定の影響も検討している。
アルゴリズム間の差と初期値依存性を実証した文献。
PMID 33021236""")

paper('7.2  Basso 2024', 44,
      '⑦ 非対称な山で初期値依存を減らす',
      ['MIMIC-III から抽出した PPG 脈波 8,000 拍',
       ['非対称分解モデル（skewed-Gaussian）を提案',
        '参照モデルは Gamma-Gaussian モデル'],
       ['残差平方和による評価で参照モデルより高精度',
        'ランダム初期値への感度が低く、一貫して頑健']],
      '裏を返せば、従来モデルは初期値で答えが動いていた', VERM,
      extra_note="""Bland-Altman でバイアスも評価している（数値は抄録に記載なし）。
基底関数の選択が決着していないことの実例であると同時に、
非対称基底の採用という対策の根拠でもある。
PMID 39577084""")

paper('7.2  Baruch 2014', 45,
      '⑧ 分解由来の指標は中心動脈圧の推定に使われている',
      ['心臓カテーテル検査を受ける患者 63 名（平均 62.7 歳）',
       ['CareTaker の非侵襲的動脈脈波を PDA で分解し、',
        '振幅比 P2P1 と時間差 T13 を抽出'],
       ['P2P1 と収縮期圧、T13 と脈圧に有意な相関',
        '（収縮期 R²=0.92、拡張期 R²=0.78、p<0.0001）']],
      '応用範囲の拡張という切り口では新規性は乏しい', GREY,
      extra_note="""中心ラインカテーテルによる中心動脈圧の実測値と比較。
Bland-Altman は AAMI SP-10 のトレンド基準内。
PDA を製品化している系譜（CareTaker／Vitalstream）があり、
実装を検討する場合は特許調査が別途必要になる。
PMID 25005686""")

paper('7.2  Pal 2024', 38,
      '⑨ DN-less 信号でも切痕を拾う ― IEM 法（用語の出典）',
      ['周術期データセット MLORD の 17,327 例',
       ['反復包絡平均（IEM）法による重複切痕の自動検出',
        '参照基準は熟練研究者による手動マーキング'],
       ['DN 検出平均誤差 IEM 0.0046 秒（SD 0.0029）',
        '従来の二次微分法は 0.0968 秒（SD 0.0909）']],
      '「DN-less signals」という語はこの論文に由来する', GOLD,
      extra_note="""収縮期相持続時間 SPD の相関は PPG で R²=0.98（n=86,764, p<.001）。
SNR −12 dB まで高い検出率を維持。
分解法（PDA）ではなく「切痕を直接検出する」系統の到達点であり、
PDA の比較対象になる。
PMID 38901273""")

# ══════════════════════════════════════════════ 7.2 限界
s = nxt('7.2  PDA の限界', cite(41, 42, 44, 46))
lims = [('①', '解が一意に決まらない',
         ['基底関数の種類と個数が決着しておらず、初期値でも解が動く'], RED),
        ('②', '「第2の山＝反射波」とは限らない',
         ['数値モデルでは第2ピークの主因はインピーダンス不整合'], RED),
        ('③', '振幅側の成績が弱い',
         ['振幅比はすべての参照値と低相関。時間側でも ρ=0.45'], VERM)]
y = 1.92
for num, head, body, col in lims:
    add_box(s, 0.45, y, 12.43, 1.26, fill=PALE, line=None)
    add_text(s, 0.68, y + 0.16, 0.62, 0.45, num, size=27, bold=True, color=col)
    add_text(s, 1.42, y + 0.18, 11.20, 0.45, head, size=25, bold=True, color=col)
    add_text(s, 1.42, y + 0.72, 11.20, 0.45, body, size=23, color=INK)
    y += 1.38

add_box(s, 0.45, 6.14, 12.43, 0.78, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.65, 6.32, 12.03, 0.45,
         '当てはまりが良くなっても、その山が何かという帰属は保証されない',
         size=25, bold=True, color=INK, align=PP_ALIGN.CENTER)

set_notes(s, """7.2 PDA の限界

① 逆問題が一意に解けていない
　　Gaussian／Gamma／skewed-Gaussian／Rayleigh、カーネル数 2〜5 と研究ごとに最良が異なる
　　（Tigges 2017 の総当たりで最頻モデルでも 28.1%）。Basso 2024 が示した初期値依存性も
　　同じ根を持つ。

② 物理的な帰属が数値モデルから支持されていない
　　Epstein 2014 は手の主要動脈を含む 75 本の動脈網を表現した非線形 1 次元モデルで、
　　SI が大動脈 PWV の直接の代用ではないこと、さらに
　　**第 2 ピークは末梢反射ではなく動脈網内部のインピーダンス不整合が主因**であり、
　　上肢の末梢反射はむしろ第 1 ピークの到達時刻を遅らせることを示した。
　　「2 番目の山＝末梢からの反射波」というラベルは自明ではない。
　　当てはまりがいくら改善しても、この帰属が誤っていれば
　　そこから作った RI・SI の生理学的意味は変わってしまう。

③ 振幅側の成績が一貫して弱い
　　Couceiro 2015 の R1_2（モデル由来の RI）はすべての参照値と低相関。
　　時間側の T1_2 ですら TPRI との ρ=0.45 にとどまる。

→ 検証設計としては、真値が既知の仮想被験者データベース（Charlton 2019）で
　 先に帰属を確かめるのが筋である。

出典：41. Couceiro 2015 ／ 42. Tigges 2017 ／ 44. Basso 2024 ／ 46. Epstein 2014""")

# ══════════════════════════════════════════════ 7.2 残された空白
s = nxt('7.2  残されている問い', cite(23, 29, 35, 41))
add_box(s, 0.45, 1.90, 6.05, 2.05, fill=PALE, line=GREY, line_w=2.0)
add_text(s, 0.65, 2.08, 5.65, 1.70,
         [[('既に済んでいること', {'bold': True, 'color': GREY, 'size': 24})],
          [('・分解から SI・RI を作る', {'size': 22})],
          [('・血圧・血管抵抗と対比する', {'size': 22})],
          [('（2008〜2015 年に完了）', {'size': 22})]],
         size=22, color=INK, space_after=6)

add_box(s, 6.83, 1.90, 6.05, 2.05, fill=PALE, line=GOLD, line_w=2.25)
add_text(s, 7.03, 2.08, 5.65, 1.70,
         [[('まだ誰もやっていないこと', {'bold': True, 'color': GOLD, 'size': 24})],
          [('・全身麻酔下・周術期での検証', {'size': 22})],
          [('・実機モニタ波形での成立性', {'size': 22})],
          [('・退化した波形での「救済」', {'size': 22})]],
         size=22, color=INK, space_after=6)

add_text(s, 0.45, 4.20, 12.43, 0.42, '設問の書き換え', size=24, bold=True, color=GREY)
add_box(s, 0.45, 4.68, 12.43, 1.75, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.70, 4.92, 11.93, 1.35,
         ['既に提案されている分解由来指標を、',
          '①全身麻酔下　②実機モニタ波形　③ランドマーク法が破綻する波形',
          'という3つの未検証条件で、変化の追随性を検証する'],
         size=24, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=5)

add_text(s, 0.45, 6.60, 12.43, 0.4,
         '※ 陰性結果でも情報価値が出る設計にしておく',
         size=22, bold=True, color=GREY)

set_notes(s, """7.2 残されている問い

■ 既出（新規性を否定する側）
・分解由来の RI ＝ Rubins 2008
・分解由来の SI と RI ＝ Goswami 2010
・血圧・総末梢血管抵抗係数との対比 ＝ Couceiro 2015
つまり「推定式で SI・RI を計算し直し、血圧・SVR との関連を調べる」という設問自体に
新規性はない。

■ 空白（検証の条件の側に残っている）
1. 周術期・全身麻酔下という設定
　 先行研究の対象は健常者・傾斜台・カフ圧迫・運動負荷・高血圧外来・血管年齢コホート。
　 麻酔導入時や昇圧薬投与時に分解由来指標が観血血圧と較正済み SVR の変化を追随するかは
　 未検証。周術期に近い先行例（Coutrot 2019、Aguet 2023、Lee 2011、Gratz 2017、
　 Khanna 2024）はいずれも設問がずれている。PubMed 検索でも該当 0 件。
2. 自動ゲイン制御と帯域制限のある実機モニタ波形での成立性
　 先行研究はすべて研究用 PPG か観血動脈圧。臨床パルスオキシメータの表示波形で
　 分解が安定に解けるかは、この構想全体の必要条件でありながら誰も確認していない。
3. 退化した波形における「救済」という仮説そのもの
　 一峰性化・shoulder 化してランドマーク法が破綻する拍で、モデル法が値を返し
　 かつその値が血管抵抗と関連するかを主要評価項目として事前規定した研究は見当たらない。
4. 絶対値の相関ではなく、変化の追随性（concordance rate・polar plot）としての評価。
5. 拍ごとに信頼区間を出し、識別不能な拍を事前規定で除外する運用。

■ 設計上の含意
Couceiro 2015 の結果を踏まえると、主要評価項目は振幅側ではなく時間側に置くべきである。
事前確率は低めに見積もり、陰性結果でも情報価値が出る設計にしておく必要がある。

検証基盤：Charlton 2019（真値が既知の仮想被験者 4,374 名）、VitalDB、pyPPG。""")

prs.save('deck_s2d.pptx')
print('stage2d ok. slides =', len(prs.slides._sldIdLst))
