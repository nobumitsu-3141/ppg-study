# -*- coding: utf-8 -*-
"""v6.8 → v6.9 改訂の共通部品：章立ての再定義と章ナビ帯の書き換え"""
from lxml import etree
from pptx.util import Emu
from deckkit import NS, INK, GOLD, LGREY, GREY, TEAL, FONT, set_title, get_title

# ---- 新しい章立て（6章「解析まとめ」を廃止、PDA を独立章に）----
CHIPS = ['基礎', '加齢', '循環', '麻酔', '定量化', '限界', 'PDA']

CH = {
    1: ('PPGの基礎', ['1.1　SpO₂の測定原理', '1.2　PPGが測るもの',
                      '1.3　前進波の成因', '1.4　反射波の成因', '1.5　PPG基礎まとめ']),
    2: ('波形への影響：加齢', ['2.1　フラミンガム研究', '2.2　UK Biobank',
                              '2.3　切痕平滑化の理由']),
    3: ('波形への影響：循環', ['3.1　前負荷', '3.2　後負荷：末梢',
                              '3.3　後負荷の違い', '3.4　心収縮力', '3.5　心拍数',
                              '3.6　波形変化まとめ']),
    4: ('麻酔中の波形変化', ['4.1　トーヌスと切痕', '4.2　波形の6分類',
                            '4.3　導入時の低血圧']),
    5: ('波形の定量化', ['5.1　Stiffness Index (SI)', '5.2　Reflection Index (RI)',
                        '5.3　加速度脈波 SDPPG', '5.4　定量化指標まとめ']),
    6: ('限界', ['6.1　自動解析の壁', '6.2　極端な動脈硬化', '6.3　DN-less 信号対策']),
    7: ('PDA（波形分解）', ['7.1　PDA とは', '7.2　分解のしくみ', '7.3　当てはめ方',
                           '7.4　得られる指標', '7.5　重要文献',
                           '7.6　限界', '7.7　残された問い']),
}


# ------------------------------------------------------------------ 章ナビ帯
def rewrite_nav(slide, active):
    """チップの文言を新章立てに合わせ、現在章だけ強調する。active=None で全灰"""
    chips = [s for s in slide.shapes if 'Off-page Connector' in s.name]
    if not chips:
        return
    chips.sort(key=lambda s: s.left)
    for i, chip in enumerate(chips, start=1):
        if i > len(CHIPS):
            chip._element.getparent().remove(chip._element)
            continue
        on = (active is not None and i == active)
        fill = TEAL if on else LGREY
        col = 'FFFFFF' if on else GREY
        spPr = chip._element.find('p:spPr', NS)
        for sf in spPr.findall('a:solidFill', NS):
            spPr.remove(sf)
        sf = etree.Element('{%s}solidFill' % NS['a'])
        etree.SubElement(sf, '{%s}srgbClr' % NS['a']).set('val', fill)
        spPr.find('a:prstGeom', NS).addnext(sf)
        # テキスト（1段落目＝番号、2段落目＝ラベル）を書き換える
        paras = chip.text_frame.paragraphs
        texts = [str(i), CHIPS[i - 1]]
        for pi, para in enumerate(paras[:2]):
            runs = para.runs
            if not runs:
                continue
            runs[0].text = texts[pi]
            for r in runs[1:]:
                r.text = ''
        for rPr in chip._element.iter('{%s}rPr' % NS['a']):
            for old in rPr.findall('a:solidFill', NS):
                rPr.remove(old)
            sf2 = etree.Element('{%s}solidFill' % NS['a'])
            etree.SubElement(sf2, '{%s}srgbClr' % NS['a']).set('val', col)
            rPr.insert(0, sf2)


# ------------------------------------------------------------------- 章扉
def find_agenda(slide):
    cands = []
    for shp in slide.shapes:
        if not shp.has_text_frame or shp.top is None:
            continue
        if 'Off-page Connector' in shp.name or shp.name in ('Title 1', 'タイトル 1'):
            continue
        if shp.top < 1.3 * 914400 or shp.top > 2.4 * 914400:
            continue
        n = len([p for p in shp.text_frame.paragraphs if p.text.strip()])
        cands.append((n, shp))
    if not cands:
        return None
    cands.sort(key=lambda t: -t[0])
    return cands[0][1]


def _para(txBody, parts, size):
    p = etree.SubElement(txBody, '{%s}p' % NS['a'])
    for txt, col, bold in parts:
        r = etree.SubElement(p, '{%s}r' % NS['a'])
        rPr = etree.SubElement(r, '{%s}rPr' % NS['a'])
        rPr.set('sz', str(size))
        rPr.set('b', '1' if bold else '0')
        sf = etree.SubElement(rPr, '{%s}solidFill' % NS['a'])
        etree.SubElement(sf, '{%s}srgbClr' % NS['a']).set('val', col)
        for tag in ('latin', 'ea', 'cs'):
            etree.SubElement(rPr, '{%s}%s' % (NS['a'], tag)).set('typeface', FONT)
        etree.SubElement(r, '{%s}t' % NS['a']).text = txt


def fill_divider(slide, active, chapters=None, sub_size=2200, main_size=2400):
    chapters = chapters or CH
    target = find_agenda(slide)
    if target is None:
        return
    target.left, target.top = Emu(int(0.94 * 914400)), Emu(int(1.63 * 914400))
    target.width, target.height = Emu(int(11.57 * 914400)), Emu(int(5.47 * 914400))
    txBody = target.text_frame._txBody
    for p in txBody.findall('a:p', NS):
        txBody.remove(p)
    for i in sorted(chapters):
        on = (i == active)
        _para(txBody, [(f'{i}. ', GOLD if on else LGREY, True),
                       (chapters[i][0], INK if on else LGREY, on)], main_size)
        if on:
            for sub in chapters[i][1]:
                num, _, label = sub.partition('　')
                _para(txBody, [('　' + num + ' ', GOLD, True), (label, INK, False)], sub_size)


def fill_menu(slide, chapters=None):
    chapters = chapters or CH
    target = find_agenda(slide)
    if target is None:
        return
    txBody = target.text_frame._txBody
    for p in txBody.findall('a:p', NS):
        txBody.remove(p)
    for i in sorted(chapters):
        _para(txBody, [(f'{i}.　', GOLD, True), (chapters[i][0], INK, True)], 2600)


# ------------------------------------------------------------ タイトル操作
def title_text(slide):
    sh = get_title(slide)
    return sh.text_frame.text.strip() if sh is not None else ''


def chapter_of(slide):
    """タイトル冒頭の数字から章番号を読む"""
    t = title_text(slide)
    if not t:
        return None
    head = t.split()[0].rstrip('.')
    part = head.split('.')[0]
    return int(part) if part.isdigit() else None
