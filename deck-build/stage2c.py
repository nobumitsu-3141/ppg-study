# -*- coding: utf-8 -*-
"""Stage 2c: 66（極端な動脈硬化・図解）と 67（DN-less）の改訂"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from deckkit import (add_text, add_box, add_arrow, add_line, add_freeform,
                     set_notes, INK, GOLD, BLUE, VERM, TEAL, RED, GREY, LGREY)
from newpages import Builder, ppg_curve, WAVE_YOUNG, WAVE_MID, WAVE_OLD, WAVE_FLAT
from refs import cite

prs = Presentation('deck_s2b.pptx')
T = {}
for s in prs.slides:
    for sh in s.shapes:
        if sh.name in ('Title 1', 'タイトル 1') and sh.has_text_frame:
            T.setdefault(sh.text_frame.text.strip(), []).append(s)

S = list(prs.slides)
B = Builder(prs, T['5.3  SDPPG の主な知見'][0])
PALE = "F2F2F2"

# ═══════════════════════ 66 の直後：切痕が消えるとどうなるか（図解版）
old66 = T['7.1  極端な動脈硬化'][0]
s = B.new(old66, '7.1  切痕が消えると', 7, cite(6, 22, 36))

waves = [('若年・柔らかい', WAVE_YOUNG, '切痕が見える', BLUE, '○'),
         ('中間', WAVE_MID, '切痕が浅い', TEAL, '△'),
         ('高齢・硬い', WAVE_OLD, '肩だけになる', VERM, '△'),
         ('極端な動脈硬化', WAVE_FLAT, '目印が消える', RED, '×')]
for i, (label, comps, note, col, mark) in enumerate(waves):
    x = 0.45 + i * 3.15
    add_text(s, x, 1.88, 2.95, 0.4, label, size=23, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    pts = ppg_curve(x + 0.10, 2.35, 2.75, 1.35, comps)
    add_freeform(s, pts, color=col, width=2.5)
    add_line(s, x + 0.10, 3.72, x + 2.85, 3.72, color=LGREY, width=1.25)
    add_text(s, x, 3.82, 2.95, 0.4, note, size=22, color=INK, align=PP_ALIGN.CENTER)
    add_text(s, x, 4.28, 2.95, 0.45, mark, size=28, bold=True, color=col,
             align=PP_ALIGN.CENTER)

add_text(s, 0.45, 4.92, 12.43, 0.35, 'P2・切痕を「点」として拾えるか',
         size=22, bold=True, color=GREY, align=PP_ALIGN.CENTER)

add_box(s, 0.45, 5.30, 6.05, 1.30, fill=PALE, line=VERM, line_w=2.25)
add_text(s, 0.62, 5.48, 5.70, 0.95,
         [[('切痕が無い人は ', {}), ('14%', {'bold': True, 'color': VERM, 'size': 28})],
          [('169,787 名中 25,286 名', {'size': 22})]],
         size=24, bold=True, color=INK, space_after=3)

add_box(s, 6.83, 5.30, 6.05, 1.30, fill=PALE, line=BLUE, line_w=2.25)
add_text(s, 7.00, 5.48, 5.70, 0.95,
         [[('残り ', {}), ('86%', {'bold': True, 'color': BLUE, 'size': 28}),
           (' は切痕が見える', {})],
          [('（UK Biobank）', {'size': 22})]],
         size=24, bold=True, color=INK, space_after=3)

add_text(s, 0.45, 6.72, 12.43, 0.4,
         '→ 問題は見えない人だけが解析から落ちること',
         size=24, bold=True, color=INK)

set_notes(s, """7.1 切痕が消えるとどうなるか（図解版・数値を修正）

【前ページの記述の訂正】
前ページは「切痕不明瞭なのは PPG を受けた 169,787 名の参加者のうち、
重複切痕は 25,286 名（14%）」と書かれていたが、文がねじれており誤読を招く。
Cunningham 2023（UK Biobank, n=169,787）で 25,286 名（14%）というのは
**切痕が「欠如」していた人数**であり、切痕が見えた人数ではない。
したがって正しくは「切痕が消えているのは 14%、残り 86% は切痕が見える」。
本リポジトリの PPG_reflection_wave_localisation でも
「『高齢者では何も見えない』は誇張。切痕欠如は 14% にとどまる」と整理している。

【このページの論点】
・動脈硬化が進むと反射波が早く戻り、切痕は「浅い切痕 → 下降脚の肩 → 消失」と変化する。
・SI（＝身長 ÷ ΔT）も RI（＝P2/P1）も「P2 という点」の時刻と高さを必要とするため、
  点が無い波形では原理的に計算できない。SDPPG（二次微分）も、下降脚が平滑になるほど
  極値が不安定になる。
・Wu 2010（Atherosclerosis 2010;213:173-7、428 名）は、まさに
  「加齢・動脈硬化例では波形が歪んで従来パラメータの正確な決定が妨げられる」ことを問題にし、
  拡張期ピークが判別できない例でも使える指標として NCT（正規化 crest time）と
  CTR（crest time ratio）を提案した。
・したがって本当の問題は「全員が見えない」ことではなく、
  **見えない 14% だけが解析から系統的に脱落し、選択バイアスを生む**ことである。
出典：6. Cunningham 2023 ／ 22. Elgendi 2012 ／ 36. Wu 2010""")

# ═════════════════════ 67 の直後：① DN-less とは（用語と出典）
old67 = T['7.1  D-N less 信号対策'][0]
s = B.new(old67, '7.1  DN-less 信号とは', 7, cite(38))

add_box(s, 0.75, 1.95, 11.83, 1.45, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.95, 2.16, 11.43, 1.05,
         [[('DN-less signal', {'bold': True, 'color': GOLD, 'size': 30})],
          [('＝ 重複切痕がはっきり見えない PPG 波形', {'bold': True, 'size': 26})]],
         size=26, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=4)

add_text(s, 0.75, 3.62, 11.83, 0.42, 'この言葉が出てくる文献', size=24, bold=True, color=GREY)
add_box(s, 0.75, 4.12, 11.83, 1.90, fill=PALE, line=None)
add_text(s, 0.95, 4.32, 11.43, 1.55,
         [[('38. Pal R, et al.  Comput Methods Programs Biomed', {'bold': True, 'color': INK, 'size': 23})],
          [('2024;254:108283（PMID 38901273）', {'size': 22})],
          [('切痕がはっきりしない信号（DN-less signals）でも検出できるかは', {'size': 23})],
          [('これまで検証されていない、と述べて新しい検出法を提案した論文', {'size': 23})]],
         size=24, color=INK, space_after=3)

add_text(s, 0.75, 6.25, 11.83, 0.85,
         ['・「DN-less」は正式な病名や規格ではなく、信号処理での呼び方',
          '・本デックでも「切痕が見えない波形」という意味でこの語を使う'],
         size=23, bold=True, color=INK, space_after=6)

set_notes(s, """7.1 DN-less 信号とは ― 用語の出典

「DN-less signal（DN-less signals）」という言い方は、重複切痕（dicrotic notch, DN）が
波形上ではっきり見えない PPG／動脈圧波形を指す、信号処理分野の呼称である。

出典は Pal R, et al.
"An algorithm to detect dicrotic notch in arterial blood pressure and photoplethysmography
 waveforms using the iterative envelope mean method."
Comput Methods Programs Biomed 2024;254:108283（PMID 38901273）。
同論文は、既存の切痕検出法について
「ノイズ耐性が十分に検証されていない、あるいは切痕が目立たない信号（DN-less signals）で
 切痕を同定できるかが検討されていない」と述べたうえで、反復包絡平均（iterative envelope
 mean, IEM）法を提案している。
（同内容の医学プレプリントが medRxiv 2024（PMID 38496617）にもある）

対象は周術期大規模データセット MLORD の 17,327 例、
動脈圧波形 1,171,288 拍・PPG 波形 3,424,975 拍。
DN 検出の平均誤差は PPG で IEM 法 0.0046 秒（SD 0.0029）に対し
従来の二次微分法は 0.0968 秒（SD 0.0909）だった。

なお「DN-less」は正式な疾患名でも規格用語でもない。本デックでは
「切痕が見えない波形」という意味で用いる。""")

# ═════════════════════ 67 の直後：② 対策の一覧（平易な表）
s2 = B.new(s, '7.1  対策の一覧', 7, cite(38, 43, 47, 50, 51))

add_text(s2, 0.45, 1.85, 3.55, 0.4, 'やり方', size=23, bold=True, color=GREY)
add_text(s2, 4.15, 1.85, 4.55, 0.4, '何をするか', size=23, bold=True, color=GREY)
add_text(s2, 8.85, 1.85, 2.95, 0.4, '切痕が無くても', size=23, bold=True, color=GREY,
         align=PP_ALIGN.CENTER)
add_text(s2, 11.85, 1.85, 1.20, 0.4, '出典', size=23, bold=True, color=GREY)

rows = [
    ('① 切痕をさがす', ['波形の曲がり角を', '直接みつける'], '×', RED, '22, 38'),
    ('② 切痕を目立たせる', ['ノイズを削って', '切痕を強調する'], '×', RED, '47, 48'),
    ('③ 波を分けてみる', ['山の足し算とみなし', '反射波を推定する'], '○', BLUE, '43'),
    ('④ AI に学ばせる', ['波形の形そのものを', '学習させる'], '○', BLUE, '50'),
    ('⑤ 形を数え上げる', ['波形を図形として', '数値化する'], '○', BLUE, '51'),
]
y = 2.35
for i, (way, what, mark, col, src) in enumerate(rows):
    if i % 2 == 0:
        add_box(s2, 0.40, y - 0.06, 12.48, 0.86, fill=PALE, line=None)
    add_text(s2, 0.45, y + 0.14, 3.55, 0.5, way, size=24, bold=True, color=INK)
    add_text(s2, 4.15, y + 0.03, 4.55, 0.7, what, size=22, color=INK, space_after=1)
    add_text(s2, 8.85, y + 0.10, 2.95, 0.5, mark, size=28, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(s2, 11.85, y + 0.16, 1.20, 0.5, src, size=22, color=GREY)
    y += 0.86

add_box(s2, 0.40, 6.72, 12.48, 0.42, fill=None, line=None)
add_text(s2, 0.45, 6.72, 12.43, 0.4,
         '③〜⑤ は「切痕という目印」に頼らない　→　次に ③ の PDA を見ていく',
         size=24, bold=True, color=INK)

set_notes(s2, """7.1 切痕が無いときの対策（平易版）

前ページの表を、高校卒業程度の言葉に置き換え、5 つの考え方に整理した。
判定欄「○」は切痕が見えなくても値が出せる、「×」は切痕が無いと成立しない、の意味。

① 切痕をさがす（Notch 直接検出）
　　波形の曲がり角を直接みつける。従来法は二次微分（SDPPG）の符号が変わる点を拾う。
　　Pal 2024 の IEM 法は、波の上側と下側の包絡線を繰り返し平均して谷の位置を出す方法で、
　　微分を使わないぶんノイズに強い。PPG の平均誤差は IEM 0.0046 秒に対し
　　二次微分法 0.0968 秒。ただし「切痕そのものを探す」以上、完全に消えていれば使えない。
　　出典：22. Elgendi 2012 ／ 38. Pal 2024

② 切痕を目立たせる（信号強調）
　　MODWT ウェーブレットで切痕の情報が乗っている周波数帯だけを取り出して強調する
　　（Attivissimo 2023、1,080 名・50,000 拍。切痕検出の平均誤差 0.0458 秒、SD 0.0896 秒）。
　　Harmonic-selective Gaussian フィルタは、位相をずらさない（ゼロ位相）フィルタで
　　高調波成分を選び、山谷の時刻をずらさずに形を保つ（Domínguez-Hernández 2026）。
　　これも「切痕がどこかに埋もれている」ことが前提。
　　出典：47. Attivissimo 2023 ／ 48. Domínguez-Hernández 2026

③ 波を分けてみる（PDA＝波形分解）
　　一拍を複数の「山」の足し算とみなし、前進波と反射波に分けて推定する。
　　切痕という点を探さないので、消えていても反射波の位置と高さを数値にできる。
　　出典：43. Fleischhauer 2020（次章で詳述）

④ AI に学ばせる（機械学習・深層学習）
　　目印を探さず、波形の形そのものから統計的に学習する。
　　Shin 2022 は 752 名の PPG から CNN で血管年齢を推定し、実年齢との相関 r=0.61
　　（平均絶対誤差 8.1 年）。ただし、なぜその値になるかは説明しにくい。
　　出典：50. Shin 2022

⑤ 形を数え上げる（グラフ理論）
　　波形の各時点を「点」とし、互いに見通せる点どうしを線で結んでネットワークにする
　　（visibility graph）。谷の有無によらず波形の複雑さを数値化できる。
　　Vargas 2025 は重み付き visibility graph の指標から PWV を推定し R²=0.91。
　　出典：51. Vargas 2025

→ 本デックでは ③ の PDA（pulse decomposition analysis）を次章で扱う。""")

prs.save('deck_s2c.pptx')
print('stage2c ok. slides =', len(prs.slides._sldIdLst))
