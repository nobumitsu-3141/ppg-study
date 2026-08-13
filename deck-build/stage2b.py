# -*- coding: utf-8 -*-
"""Stage 2b: 61 / 62（出典つき対比表）、6章の新設、66・67 の改訂"""
from pptx import Presentation
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from deckkit import (add_text, add_box, add_arrow, add_line, add_freeform,
                     set_notes, set_title, set_nav, set_source, clone_slide,
                     strip_content, slide_index, set_pageno,
                     INK, GOLD, BLUE, VERM, TEAL, RED, GREY, LGREY)
from newpages import Builder, ppg_curve, WAVE_YOUNG, WAVE_MID, WAVE_OLD, WAVE_FLAT
from refs import cite

prs = Presentation('deck_s2a.pptx')
T = {}
for i, s in enumerate(prs.slides, 1):
    for sh in s.shapes:
        if sh.name in ('Title 1', 'タイトル 1') and sh.has_text_frame:
            T.setdefault(sh.text_frame.text.strip(), []).append(s)

S = list(prs.slides)
B = Builder(prs, S[62])          # 5.3 SDPPG の主な知見（= 現 63 枚目）を版面の型に使う
PALE = "F2F2F2"


def comp_table(slide, headL, headR, rows, y0=1.90, rowh=0.80,
               colx=(0.45, 4.15, 7.55, 10.75), colw=(3.55, 3.25, 3.05, 2.10)):
    """左見出し／指標A／指標B／出典 の4列比較表"""
    heads = ['', headL, headR, '出典']
    for i, (hx, hw, ht) in enumerate(zip(colx, colw, heads)):
        if ht:
            add_text(slide, hx, y0, hw, 0.42, ht, size=24, bold=True,
                     color=INK if i < 3 else GREY,
                     align=PP_ALIGN.CENTER if i else PP_ALIGN.LEFT)
    y = y0 + 0.52
    for k, (label, a, b, src) in enumerate(rows):
        if k % 2 == 0:
            add_box(slide, colx[0] - 0.05, y - 0.06, 12.45, rowh, fill=PALE, line=None)
        add_text(slide, colx[0], y + 0.10, colw[0], 0.62, label, size=22, bold=True, color=INK)
        add_text(slide, colx[1], y + (0.02 if isinstance(a, list) and len(a) > 1 else 0.10),
                 colw[1], 0.70, a, size=22, color=INK, align=PP_ALIGN.CENTER, space_after=1)
        add_text(slide, colx[2], y + (0.02 if isinstance(b, list) and len(b) > 1 else 0.10),
                 colw[2], 0.70, b, size=22, color=INK, align=PP_ALIGN.CENTER, space_after=1)
        add_text(slide, colx[3], y + 0.12, colw[3], 0.62, src, size=22, color=GREY)
        y += rowh


# ══════════════════════════════ 61（b/a と SI）の直後：出典を紐づけた改訂版
old61 = T['5.3  b/a と SI'][0]
s = B.new(old61, '5.3  b/a と SI', 5, cite(4, 27, 31, 32, 33))
comp_table(s, 'b/a（SDPPG）', 'SI', [
    ('年齢との相関', '○', '○（R=0.63）', '27, 31'),
    ('PWV との相関', '○', '○（r=0.65）', '4, 31'),
    ('血圧との相関', '×', '△（MAP 依存）', '4, 32'),
    ('利点', ['a・b 波の検出が', '頑健'], ['大動脈弾性と', '相関'], '33, 4'),
    ('批判', ['臓器障害との', '関連は弱い'], ['前向き研究が', 'ない'], '32'),
])
add_text(s, 0.45, 6.60, 12.45, 0.4,
         '※ b/a の相関係数は抄録に記載がなく、原著本文での確認が必要（要出典）',
         size=22, bold=True, color=RED)
set_notes(s, """5.3 b/a と SI ― 文献との整合性チェック（改訂版）

【前ページからの修正点】
1. 出典表記をカタカナ（ミラソー2003 等）から番号つきローマ字に統一した。
2. 「hong2023」は該当文献を同定できなかったため削除した（引用するなら書誌情報が要る）。
3. 「批判：臓器障害との関連がない」→「関連は弱い」に修正。
   根拠は Tabara 2016（J-SHIPP）。SDPTG 指標と頸動脈 IMT 等の関連は
   「弱いか有意でない」であって「無い」ではない（b/a で β=0.069, P=0.002）。
4. 年齢との相関 0.63 は Millasseau **2003**（Am J Hypertens 2003;16:467-72）の値。
   Millasseau **2002** の抄録に SI と年齢の単相関 r は記載がなく、
   載っているのは重回帰 SI=0.63+0.086×age+0.042×MAP（r=0.69）である。
   → 本デックの「5.1 Stiffness Index」ページに書かれている「SI と年齢 r=0.67」は
     Millasseau 2002 からは裏づけられない。要修正。
5. SI と PWVcf の r=0.65（P<0.0001, n=87）は Millasseau 2002 のとおりで正しい。
6. 「血圧との相関 SI あり（特に DBP で 0.68）」は出典を特定できなかったため
   「△（MAP と共変動）」に改めた。Millasseau 2002 の重回帰では MAP が独立に効く。

【未確認のまま残した数値】
・b/a の 年齢 0.51 ／ PWV 0.61 ／ 血圧 0.19 は、von Wowern 2015・Tabara 2016 の
  抄録からは確認できなかった。原著本文で確認するまでは○×の定性表示にとどめてある。""")

# ══════════════════════════════ 62（d/a と RI）の直後：出典を紐づけた改訂版
old62 = T['5.3  d/a と RI'][0]
s = B.new(old62, '5.3  d/a と RI', 5, cite(16, 24, 27, 29, 34, 35))
comp_table(s, 'd/a（SDPPG）', 'RI', [
    ('昇圧薬 AT II', '○（低下）', '○（上昇）', '16, 27'),
    ('降圧薬 GTN', '△（不変）', '○（低下）', '27'),
    ('術中の変化', ['挿管時高血圧で', '低下'], ['導入時低血圧を', '検出'], '34, 29'),
    ('SVR との比較', '報告なし', ['κ=0.33', '（単独）'], '35'),
    ('批判', ['物理的な意味が', '明白でない'], ['HR との交絡', 'がある'], '27, 24'),
])
add_text(s, 0.45, 6.60, 12.45, 0.4,
         '※ 29. Coutrot 2019 が測ったのは Dicpleth であって RI そのものではない',
         size=22, bold=True, color=RED)
set_notes(s, """5.3 d/a と RI ― 文献との整合性チェック（改訂版）

【前ページからの修正点】
1. 「Mdlazim 2023」→ 24. Md Lazim **2020**（Int J Environ Res Public Health 2020;17:2591）。
   2023 年版は存在しない。
2. 「コウジタニ2012」→ 34. Kohjitani 2012（Hypertens Res 2012;35:53-60）。
   セボフルラン単独群で挿管後に d/a が低下（＝末梢血管抵抗上昇）した、が正確な内容。
3. 昇圧薬 AT II に対する反応：Millasseau 2003 では AT II で **RI は 12.6% 増加**、
   d/a は 0.18 低下。前ページは d/a・RI とも「低下」と書いていたが、RI は上昇である。
   → 向きを修正した。
4. 降圧薬 GTN：Millasseau 2003 で **RI は 32% 減少**、d/a は用量依存的に変化しなかった。
   前ページは RI「上昇」と書いていたが、正しくは低下。→ 修正した。
5. 「SVR との比較：RI あり（まずまず）」→ Lee QY 2011 の実測は RI 単独で κ=0.33
   （最良の特徴量セットで 0.57）。「まずまず」は過大評価なので κ の値を明示した。
6. 出典表記をカタカナから番号つきローマ字に統一した。

【注意】
・Coutrot 2019 が追ったのは Dicpleth＝重複切痕の高さ÷脈波振幅であり、
  P2/P1 として定義される RI とは別量である。同一視しないこと。""")

# ══════════════════════════════════ 5.4 の直後：6章の扉と 6.1 本文を新設
old54 = T['5.4  定量化指標まとめ'][0]

# --- 6章扉（4章扉を型に使う）---
div_tpl = T['4. 麻酔中の波形変化'][0]
idx = slide_index(prs, old54)
d = clone_slide(prs, div_tpl, insert_after_idx=idx)
set_title(d, '6. 解析まとめ')
set_nav(d, 6)
set_source(d, None)
set_pageno(d, None)
# 目次本文を差し替え
from stage1helpers import fill_divider
fill_divider(d, 6)

s = B.new(d, '6.1  二軸で読む', 6, cite(4, 14, 16, 20, 29))
add_box(s, 0.55, 1.95, 6.00, 2.35, fill=PALE, line=BLUE, line_w=2.25)
add_text(s, 0.75, 2.15, 5.60, 0.5, '慢性 ＝ スティフネス軸', size=27, bold=True, color=BLUE)
add_text(s, 0.85, 2.80, 5.40, 1.4,
         ['・切痕の平滑さ', '・SI ／ SDPPG の b/a・AGI', '・術前ベースラインの把握に'],
         size=23, color=INK, space_after=5)

add_box(s, 6.78, 1.95, 6.00, 2.35, fill=PALE, line=VERM, line_w=2.25)
add_text(s, 6.98, 2.15, 5.60, 0.5, '急性 ＝ トーヌス軸', size=27, bold=True, color=VERM)
add_text(s, 7.08, 2.80, 5.40, 1.4,
         ['・切痕の位置、反射波の高さ', '・RI ／ SDPPG の d/a', '・術中の一拍ごとの追跡に'],
         size=23, color=INK, space_after=5)

add_box(s, 0.55, 4.55, 12.23, 1.15, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.75, 4.78, 11.83, 0.75,
         ['絶対値の較正は目指さない', '同じ患者・同じ条件での 変化量（Δ）だけを読む'],
         size=26, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=3)

add_text(s, 0.55, 5.95, 12.23, 0.85,
         ['・術前：スティフネス軸でベースラインを取る',
          '・術中：トーヌス軸の Δ で後負荷の動きを追う'],
         size=24, bold=True, color=INK, space_after=5)

set_notes(s, """6.1 解析まとめ ― 二軸で読む

PPG 波形から読める情報は、時間スケールの違う 2 本の軸に整理できる。

【慢性＝スティフネス軸】加齢・動脈硬化で年単位に動く
・切痕の平滑さ（Dawber 1973 の 4 分類、Cunningham 2023 の notch smoothness）
・SI＝身長 ÷ ΔT（Millasseau 2002）
・SDPPG の b/a、血管加齢指数 AGI（Takazawa 1998）
→ 術前のベースライン評価に向く。

【急性＝トーヌス軸】麻酔薬・血管作動薬で分単位に動く
・切痕の位置、反射波の高さ（Murray 1996、Tusman 2019）
・RI＝P2/P1（Chowienczyk 1999）
・SDPPG の d/a（Takazawa 1998）
・Dicpleth＝切痕高 ÷ 振幅（Coutrot 2019、Joachim 2021）
→ 術中の一拍ごとの追跡に向く。

【運用上の結論】
指標はいずれも血圧・心拍数に依存し、絶対値の較正には耐えない。
同一患者・同一条件での相対変化（Δ）に限って解釈する、というのが本デックの立場である。""")

prs.save('deck_s2b.pptx')
print('stage2b ok. slides =', len(prs.slides._sldIdLst))
