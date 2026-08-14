# -*- coding: utf-8 -*-
"""R2 stage2: PDA 章の序盤解説・文献一覧表の追加、用語「山」→「成分波」の統一"""
import math
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.enum.text import PP_ALIGN
from deckkit import (add_text, add_box, add_arrow, add_line, add_freeform,
                     set_notes, set_title, set_source, set_pageno, clone_slide,
                     strip_content, slide_index, move_slide, drop_shapes,
                     INK, GOLD, BLUE, VERM, TEAL, RED, GREY, LGREY)
from r2_common import rewrite_nav, title_text
from refs import cite
from newpages import ppg_curve, comp_curve

prs = Presentation('r2_s1.pptx')
PALE = "F2F2F2"


def idx_of_title(t, nth=0):
    hits = [i for i, s in enumerate(prs.slides, 1) if title_text(s) == t]
    return hits[nth]


def S(n):
    return list(prs.slides)[n - 1]


# ═════════════ 0) 役目を終えた非表示の下書きページを削除
NSP = 'http://schemas.openxmlformats.org/presentationml/2006/main'
RNS = 'http://schemas.openxmlformats.org/officeDocument/2006/relationships'
sldIdLst = prs.slides._sldIdLst
kill = [i for i, sl in enumerate(prs.slides, 1)
        if sl.element.get('show') == '0'
        and (title_text(sl).startswith('7.1  D-N less') or title_text(sl) == '')]
for i in sorted(kill, reverse=True):
    el = list(sldIdLst)[i - 1]
    rid = el.get('{%s}id' % RNS)
    sldIdLst.remove(el)
    prs.part.drop_rel(rid)
print('旧下書きを削除:', kill)

# ═════════════ 1) 用語の統一（山 → 成分波）
TERMS = [
    ('2つの山の足し算とみなす', '2つの成分波の足し算とみなす'),
    ('波形を山の重ね合わせに当てはめ', '波形を成分波の重ね合わせに当てはめ'),
    ('山の形を決める', '成分波の形を決める'),
    ('山の位置・高さ', '成分波の位置・高さ'),
    ('山どうしの', '成分波どうしの'),
    ('どの形の山を何個使うか', 'どの形の成分波を何個使うか'),
    ('対称の山に歪みを', '対称形に歪みを'),
    ('山の足し算とみなし', '成分波の足し算とみなし'),
    ('その山が何かという帰属', 'その成分波が何かという帰属'),
    ('複数の「山」の足し算', '複数の「成分波」の足し算'),
    ('複数の山（基底関数）', '複数の成分波（基底関数）'),
    ('2つの「山」の重ね合わせ', '2つの成分波の重ね合わせ'),
    ('山の位置(時刻)の差', '成分波の位置(時刻)の差'),
    ('山の高さの比', '成分波の高さの比'),
    ('釣鐘型の山', '釣鐘型の成分波'),
    ('複数のガウス山', '複数のガウス成分波'),
]


def sub_all(txt):
    for a, b in TERMS:
        txt = txt.replace(a, b)
    return txt


n_term = 0
for s in prs.slides:
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        for para in sh.text_frame.paragraphs:
            runs = para.runs
            if not runs:
                continue
            whole = ''.join(r.text for r in runs)
            new = sub_all(whole)
            if new != whole:
                runs[0].text = new
                for r in runs[1:]:
                    r.text = ''
                n_term += 1
    if s.has_notes_slide:
        tf = s.notes_slide.notes_text_frame
        nt = sub_all(tf.text)
        if nt != tf.text:
            tf.text = nt
print('用語を置換した段落:', n_term)

# ═════════════ 2) 7.1「PDA とは」の PDA 枠のあふれを直す
s70 = S(idx_of_title('7.1  PDA とは'))
for sh in list(s70.shapes):
    if sh.has_text_frame and sh.text_frame.text.startswith('PDA（Pulse'):
        sh._element.getparent().remove(sh._element)
    if sh.name.startswith('Rounded Rectangle') and sh.top is not None \
            and abs(sh.top / 914400 - 4.20) < 0.05:
        sh.top, sh.height = Emu(int(4.12 * 914400)), Emu(int(2.48 * 914400))
add_text(s70, 6.92, 4.32, 5.86, 2.10,
         [[('PDA', {'bold': True, 'color': GOLD, 'size': 25})],
          [('Pulse Decomposition Analysis', {'size': 21, 'color': GOLD})],
          [('波形を成分波の重ね合わせに当てはめ', {'size': 23})],
          [('反射波の 時刻と高さ を推定する', {'size': 23, 'bold': True})],
          [('→ 切痕が消えていても値が出せる', {'size': 23, 'color': GOLD, 'bold': True})]],
         size=23, color=INK, space_after=3)

# ═════════════ 3) スライドの並べ替え（基底関数の形 → 分解の手順 の順に）
set_title(S(idx_of_title('7.2  基底関数の種類')), '7.2  基底関数の形')
set_title(S(idx_of_title('7.3  分解の手順')), '7.3  当てはめの手順')

i_form = idx_of_title('7.2  基底関数の形')
i_proc = idx_of_title('7.3  当てはめの手順')
if i_form > i_proc:
    move_slide(prs, i_form - 1, i_proc - 1)

# ═════════════ 4) 新規スライドの土台
tpl = S(idx_of_title('5.3  SDPPG の主な知見'))


def new_after(anchor_slide, title, source=None, notes=None):
    i = slide_index(prs, anchor_slide)
    s = clone_slide(prs, tpl, insert_after_idx=i)
    strip_content(s)
    set_title(s, title)
    rewrite_nav(s, 7)
    set_source(s, source)
    set_pageno(s, None)
    if notes:
        set_notes(s, notes)
    return s


# ---------------------------------------------------- 7.2 合成波と成分波
anchor = S(idx_of_title('7.1  PDA とは'))
COMPS = [(0.22, 0.085, 1.00), (0.55, 0.125, 0.46)]

s = new_after(anchor, '7.2  合成波と成分波', cite(39, 41),
              """7.2 合成波と成分波 ― 言葉の整理

本章では次のように言葉を使い分ける。

・合成波（composite wave）＝ モニタに映る、実際に記録された 1 拍の PPG 波形。
・成分波（component wave）＝ 合成波を作っているとみなす個々の波。
　　第 1 成分波 ＝ 心臓から来た前進波
　　第 2 成分波 ＝ 末梢から戻ってきた反射波
・基底関数（basis function）＝ 成分波の形を表す数式のテンプレート。
　　ガウス関数・Gamma 分布・skewed-Gaussian などが使われる。

「山」という言い方は形を思い浮かべやすい一方で、
本デックが一貫して使ってきた「前進波・反射波」という言葉と混ざるため用いない。
波の話は「波」で統一する。

補足：この考え方は音の重ね合わせと同じである。複数の楽器の音が空気中で足し算されて
1 つの波形になるのと同様に、PPG も前進波と反射波が足し算されて 1 つの波形になっている、
と考える。""")
add_text(s, 0.55, 1.85, 5.9, 0.45, '合成波（実際に見える波）', size=24, bold=True, color=INK)
add_freeform(s, ppg_curve(0.60, 2.35, 5.30, 1.45, COMPS), color=INK, width=3.0)
add_line(s, 0.60, 3.82, 5.90, 3.82, color=LGREY, width=1.25)
add_text(s, 0.55, 4.05, 5.9, 0.45, '成分波（分けて考えた波）', size=24, bold=True, color=INK)
add_freeform(s, comp_curve(0.60, 4.55, 5.30, 1.40, COMPS, COMPS[0]), color=BLUE,
             width=2.5, dash='dash')
add_freeform(s, comp_curve(0.60, 4.55, 5.30, 1.40, COMPS, COMPS[1]), color=VERM,
             width=2.5, dash='dash')
add_line(s, 0.60, 5.97, 5.90, 5.97, color=LGREY, width=1.25)
add_text(s, 0.60, 6.08, 2.60, 0.45, '第1成分波', size=22, bold=True, color=BLUE)
add_text(s, 3.30, 6.08, 2.60, 0.45, '第2成分波', size=22, bold=True, color=VERM)

rows = [('合成波', '実際に記録された 1 拍の波形', INK),
        ('成分波', '合成波を作っているとみなす個々の波', BLUE),
        ('基底関数', '成分波の形を表す数式のひな型', TEAL)]
y = 1.90
for name, desc, col in rows:
    add_box(s, 6.70, y, 6.15, 1.18, fill=PALE, line=col, line_w=2.25)
    add_text(s, 6.90, y + 0.12, 5.75, 0.42, name, size=24, bold=True, color=col)
    add_text(s, 6.90, y + 0.62, 5.75, 0.42, desc, size=22, color=INK)
    y += 1.30

add_box(s, 6.70, 5.80, 6.15, 1.05, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 6.88, 5.99, 5.79, 0.70,
         ['第1成分波 ＝ 前進波', '第2成分波 ＝ 反射波 とみなす'],
         size=23, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=3)
anchor = s

# ---------------------------------------------------- 7.2 基底関数とは
s = new_after(anchor, '7.2  基底関数とは', cite(39, 42),
              """7.2 基底関数とは

基底関数（basis function）は、成分波の形を表す数式のテンプレートである。
1 本の成分波は、ふつう次の 3 つのパラメータで決まる。

・振幅 a … その成分波の高さ
・位置 μ … その成分波の頂点が来る時刻
・幅 σ  … その成分波の広がり（持続時間）

たとえばガウス関数なら
　　f(t) = a・exp( −(t − μ)² / 2σ² )
と書ける。合成波はこの成分波を n 本足し合わせたものとしてモデル化される。
　　y(t) ≒ Σ f_k(t)　（k = 1 … n）

■ なぜ「関数」で表すのか
点（切痕・P2）を探す方法は、その点が波形上に存在しないと破綻する。
基底関数を当てはめる方法なら、点が見えなくても
「第 2 成分波の位置 μ₂ と振幅 a₂ はいくつか」という形で答えが出る。
これが PDA の核心である。

■ 使われる基底関数
ガウス関数（対称）、Gamma 分布（非対称）、skewed-Gaussian（非対称）、
Rayleigh 関数、log-normal など。次ページで形を比べる。

出典：39. Rubins 2008（ガウス当てはめの最初期）／42. Tigges 2017（基底関数の総当たり）""")

add_box(s, 0.55, 1.88, 12.25, 1.30, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.75, 2.10, 11.85, 0.90,
         [[('基底関数 ＝ 成分波の形を表す ', {}), ('数式のひな型', {'bold': True, 'color': GOLD})],
          [('1本の成分波は 高さ・位置・幅 の3つの数字で決まる', {'size': 23})]],
         size=25, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=4)

pars = [('高さ  a', '振幅', 'どれだけ大きいか', BLUE),
        ('位置  μ', '時刻', 'いつ頂点が来るか', VERM),
        ('幅  σ', '広がり', 'どれだけ長く続くか', TEAL)]
for i, (sym, nm, desc, col) in enumerate(pars):
    x = 0.55 + i * 4.16
    add_box(s, x, 3.42, 3.86, 1.95, fill=PALE, line=col, line_w=2.25)
    add_text(s, x + 0.15, 3.62, 3.56, 0.5, sym, size=27, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.15, 4.20, 3.56, 0.45, nm, size=23, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.10, 4.72, 3.66, 0.45, desc, size=22, color=INK,
             align=PP_ALIGN.CENTER)

add_box(s, 0.55, 5.62, 12.25, 1.10, fill=PALE, line=None)
add_text(s, 0.75, 5.84, 11.85, 0.70,
         ['合成波 ≒ 成分波1 ＋ 成分波2 ＋ …',
          '点を探すのではなく、この足し算が実測に一番合う数字を探す'],
         size=24, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=3)
anchor = s

# ---------------------------------------------------- 7.2 カーネルの数
anchor = S(idx_of_title('7.2  基底関数の形'))
s = new_after(anchor, '7.2  成分波を何本使うか', cite(42, 43),
              """7.2 成分波を何本使うか（カーネル数）

「カーネル」は、当てはめに使う基底関数 1 本 1 本のこと。
何本使うか（カーネル数）は研究者が決める設計値であり、正解は決まっていない。

・2 本 … 前進波＋反射波。いちばん素直だが、拡張期の細かい形は表せない。
・3〜5 本 … 反射が複数回戻ってくるとみなす。当てはまりは良くなるが、
　　　　　　どの成分波が何に対応するかの解釈が曖昧になる。

■ 本数が増えるとどうなるか
本数を増やせば残差はいくらでも小さくできる（過剰適合）。
一方でパラメータが増えるほど解が一意でなくなり、初期値しだいで答えが動く。
そこで情報量規準（AICc など）で「当てはまりの良さ」と「パラメータの少なさ」を
つり合わせてモデルを選ぶ。

■ 実際の報告
・Tigges 2017：指尖容積脈波 7,805 拍で 4 種の基底関数とモデル次数を総当たりし AICc で選択。
　最頻の最良モデルは Gamma 3 本で 28.1%、次点は Rayleigh 2 本で 14.4%。
　単一の最良モデルは存在しなかった。
・Fleischhauer 2020：カーネル 2 本がノイズ・体動にもっとも頑健で、
　評価指標が 14.09% 改善しつつ、より多いカーネルと同等に形態情報を保った。

→ 「多いほど良い」ではない。2 本が実用上のバランス点という報告が複数ある。

出典：42. Tigges 2017 ／ 43. Fleischhauer 2020""")

add_text(s, 0.55, 1.85, 12.25, 0.45, 'カーネル ＝ 当てはめに使う成分波 1 本 1 本',
         size=24, bold=True, color=GREY)

kk = [('2 本', ['前進波 ＋ 反射波', '解釈がはっきりする', 'ノイズに強い'], BLUE, '○'),
      ('3〜5 本', ['反射が複数回戻ると', 'みなす', '当てはまりは良い'], TEAL, '△'),
      ('もっと多い', ['残差はいくらでも', '小さくできるが', '解が定まらない'], VERM, '×')]
for i, (n, body, col, mark) in enumerate(kk):
    x = 0.55 + i * 4.16
    add_box(s, x, 2.45, 3.86, 2.80, fill=PALE, line=col, line_w=2.25)
    add_text(s, x + 0.15, 2.65, 3.56, 0.5, n, size=27, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.12, 3.28, 3.62, 1.45, body, size=22, color=INK,
             align=PP_ALIGN.CENTER, space_after=4)
    add_text(s, x + 0.15, 4.72, 3.56, 0.45, mark, size=27, bold=True, color=col,
             align=PP_ALIGN.CENTER)

add_box(s, 0.55, 5.52, 12.25, 1.20, fill='FFF3D6', line=GOLD, line_w=2.25)
add_text(s, 0.75, 5.74, 11.85, 0.80,
         ['本数を増やせば残差は小さくなるが、解が一意でなくなる',
          'カーネル 2 本がもっとも頑健という報告が複数ある'],
         size=24, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=3)
anchor = s

# ---------------------------------------------------- 7.3 当てはめの計算
anchor = S(idx_of_title('7.3  当てはめの手順'))
s = new_after(anchor, '7.3  当てはめの計算', cite(41, 44),
              """7.3 当てはめの計算 ― 何をどう決めているのか

■ 目的関数
実測した合成波 y(t) と、モデルが作る合成波 ŷ(t)=Σf_k(t) の差（残差）を、
全時点で二乗して足し合わせた値（残差平方和 RSS）を最小にするパラメータを探す。
　　RSS = Σ_t { y(t) − ŷ(t) }²
サンプル点ごとに重みを変える重み付き最小二乗（WLS）を使う実装もある
（Wang L 2013 は重みを多基準意思決定法で最適化している）。

■ 解き方
パラメータについて非線形なので、解析的には解けない。
Levenberg–Marquardt 法などの反復計算で、初期値から少しずつ動かして最小値を探す。

■ 初期値の問題
反復計算は初期値の近くの谷（局所解）に落ちる。したがって
「同じ波形・同じモデルでも、初期値を変えると別の答えが出る」ことが起こる。
Basso 2024 は skewed-Gaussian モデルがランダム初期値に対して感度が低く
一貫して頑健だったと報告しており、これは裏を返せば
従来モデルでは初期値で答えが動いていたということである。

■ 当てはまりの評価
・残差平方和 RSS、正規化 RMSE、R²
・モデル選択には補正赤池情報量規準 AICc（当てはまりとパラメータ数のつり合い）
Couceiro 2015 の 5 ガウスモデルでは、LVET 絶対誤差 15.41±13.66 ms、ρ=0.78 が得られている。

出典：41. Couceiro 2015 ／ 44. Basso 2024""")

steps2 = [('残差を測る', ['実測 と モデルの差を', '二乗して足し合わせる', '（残差平方和）'], BLUE),
          ('少しずつ動かす', ['初期値から反復計算で', '残差が最小になる', 'パラメータを探す'], TEAL),
          ('当てはまりを見る', ['R² や AICc で', '当てはまりと', 'パラメータ数を比べる'], GOLD)]
for i, (hd, body, col) in enumerate(steps2):
    x = 0.55 + i * 4.16
    add_box(s, x, 1.90, 3.86, 2.65, fill=PALE, line=col, line_w=2.25)
    add_text(s, x + 0.12, 2.10, 3.62, 0.5, hd, size=24, bold=True, color=col,
             align=PP_ALIGN.CENTER)
    add_text(s, x + 0.12, 2.75, 3.62, 1.55, body, size=22, color=INK,
             align=PP_ALIGN.CENTER, space_after=4)
    if i < 2:
        add_arrow(s, x + 3.91, 3.00, 0.18, 0.32, color=GREY)

add_box(s, 0.55, 4.82, 12.25, 1.90, fill=PALE, line=RED, line_w=2.25)
add_text(s, 0.75, 5.02, 11.85, 1.50,
         [[('落とし穴：初期値で答えが変わる', {'bold': True, 'color': RED, 'size': 25})],
          [('反復計算は初期値の近くの谷にはまる（局所解）', {'size': 23})],
          [('同じ波形・同じモデルでも、初期値しだいで別の値が出ることがある', {'size': 23})]],
         size=23, color=INK, space_after=6)
anchor = s

# ---------------------------------------------------- 7.4 得られる指標
s = new_after(anchor, '7.4  得られる指標', cite(40, 41, 11),
              """7.4 分解から何が得られるか

分解が終わると、各成分波の（高さ a・位置 μ・幅 σ）が数値として手に入る。
そこから従来の指標を作り直せる。

・時間差 T1_2 ＝ μ₂ − μ₁（第1成分波と第2成分波の頂点の時刻差）
　　→ 従来の ΔT に相当。SI ＝ 身長 ÷ T1_2 として再定義できる。
・振幅比 R1_2 ＝ a₂ / a₁
　　→ 従来の RI（P2/P1）に相当。
・LVET（左室駆出時間）も、収縮期成分の微分解析から推定できる（Couceiro 2012）。

■ 従来法との決定的な違い
従来法は「切痕・P2 という点」が波形上に存在しないと計算できない。
分解法は点を探さないので、切痕が消えている波形（DN-less）でも
μ₂ と a₂ という数値が返ってくる。これが本章の要点である。

■ ただし
Couceiro 2015 では、振幅比 R1_2 はすべての参照値と低相関だった。
時間側の T1_2 でも総末梢血管抵抗係数との相関は ρ=0.45 にとどまる。
値が返ってくることと、その値が臨床的に意味を持つことは別である。

出典：11. Couceiro 2012 ／ 40. Goswami 2010 ／ 41. Couceiro 2015""")

add_text(s, 0.55, 1.85, 6.0, 0.45, '分解して得られる数値', size=24, bold=True, color=INK)
add_freeform(s, comp_curve(0.60, 2.35, 5.30, 1.55, COMPS, COMPS[0]), color=BLUE,
             width=2.5, dash='dash')
add_freeform(s, comp_curve(0.60, 2.35, 5.30, 1.55, COMPS, COMPS[1]), color=VERM,
             width=2.5, dash='dash')
add_line(s, 0.60, 3.92, 5.90, 3.92, color=LGREY, width=1.25)
add_line(s, 1.77, 2.42, 1.77, 3.92, color=BLUE, width=1.25, dash='sysDash')
add_line(s, 3.52, 2.98, 3.52, 3.92, color=VERM, width=1.25, dash='sysDash')
add_text(s, 1.55, 4.00, 2.0, 0.42, 'μ₁', size=23, bold=True, color=BLUE)
add_text(s, 3.30, 4.00, 2.0, 0.42, 'μ₂', size=23, bold=True, color=VERM)
add_text(s, 0.60, 4.55, 5.4, 0.42, 'a₁ ＝ 第1成分波の高さ', size=22, color=BLUE)
add_text(s, 0.60, 5.05, 5.4, 0.42, 'a₂ ＝ 第2成分波の高さ', size=22, color=VERM)

idxs = [('T1_2 ＝ μ₂ − μ₁', '時間差', 'SI ＝ 身長 ÷ T1_2', BLUE),
        ('R1_2 ＝ a₂ ／ a₁', '振幅比', 'RI に相当', VERM),
        ('LVET', '左室駆出時間', '収縮期成分から推定', TEAL)]
y = 1.90
for f, nm, use, col in idxs:
    add_box(s, 6.70, y, 6.15, 1.18, fill=PALE, line=col, line_w=2.25)
    add_text(s, 6.90, y + 0.12, 5.75, 0.42, f, size=24, bold=True, color=col)
    add_text(s, 6.90, y + 0.62, 5.75, 0.42, f'{nm} → {use}', size=22, color=INK)
    y += 1.30

add_box(s, 6.70, 5.80, 6.15, 1.05, fill=PALE, line=RED, line_w=2.25)
add_text(s, 6.88, 5.99, 5.79, 0.70,
         ['値は返ってくる。', 'ただし臨床的な意味は別問題'],
         size=23, bold=True, color=INK, align=PP_ALIGN.CENTER, space_after=3)
anchor = s

# ---------------------------------------------------- 7.5 重要文献の全体像
s = new_after(anchor, '7.5  重要文献の全体像', None,
              """7.5 重要文献の全体像

PDA をめぐる 9 本を、役割で 4 群に分けて並べた。次ページ以降で 1 本ずつ扱う。

【A 群：指標を作った】
・39. Rubins 2008 … ガウス当てはめから RI を出した最初期（健常者 40 名）
・40. Goswami 2010 … Rayleigh 2 波合成モデルで SI と RI を同時に導出（113 信号）
・11. Couceiro 2012 … 多ガウス分解から LVET を推定（誤差 15.84±13.56 ms）
・41. Couceiro 2015 … 5 ガウス分解を血圧・総末梢血管抵抗係数まで対比した到達点

【B 群：モデルの選び方】
・42. Tigges 2017 … 7,805 拍の総当たり。最頻モデルでも 28.1%
・43. Fleischhauer 2020 … Gamma＋Gaussian・カーネル 2 本が最も頑健（14.09% 改善）
・44. Basso 2024 … skewed-Gaussian。初期値に頑健（MIMIC-III 8,000 拍）

【C 群：臨床応用】
・45. Baruch 2014 … PDA 由来の P2P1・T13 で中心動脈圧を推定（63 名、R²=0.92）

【D 群：比較対象】
・38. Pal 2024 … 分解ではなく切痕を直接検出する IEM 法。DN-less という語の出典

【読みどころ】
A 群を時系列に読むと、「点を探すのをやめて波を当てはめる」という発想の転換も、
そこから SI・RI を作り直すことも、血圧・血管抵抗と対比することも、
2008〜2015 年にすべて実施済みだとわかる。しかも成績は芳しくない。
B 群は、その手法がいまだに一意に解けていないことを示す。""")

add_text(s, 0.45, 1.80, 2.30, 0.4, '群', size=22, bold=True, color=GREY)
add_text(s, 3.00, 1.80, 3.30, 0.4, '文献', size=22, bold=True, color=GREY)
add_text(s, 6.55, 1.80, 6.30, 0.4, '何を示したか', size=22, bold=True, color=GREY)

lit = [
    ('A 指標を作った', '39. Rubins 2008', 'ガウス当てはめで RI を算出', BLUE),
    ('', '40. Goswami 2010', 'Rayleigh 2 波合成で SI と RI', BLUE),
    ('', '11. Couceiro 2012', '多ガウス分解で LVET を推定', BLUE),
    ('', '41. Couceiro 2015', '血圧・血管抵抗まで対比した到達点', BLUE),
    ('B モデル選択', '42. Tigges 2017', '最良モデルは決まらない（28.1%）', TEAL),
    ('', '43. Fleischhauer 2020', 'カーネル 2 本が最も頑健', TEAL),
    ('', '44. Basso 2024', 'skewed-Gaussian は初期値に頑健', TEAL),
    ('C 臨床応用', '45. Baruch 2014', '中心動脈圧の推定に実装済み', GOLD),
    ('D 比較対象', '38. Pal 2024', '分解ではなく切痕を直接検出', VERM),
]
y = 2.22
for i, (grp, name, what, col) in enumerate(lit):
    if i % 2 == 0:
        add_box(s, 0.40, y - 0.04, 12.53, 0.50, fill=PALE, line=None)
    if grp:
        add_text(s, 0.45, y + 0.02, 2.45, 0.42, grp, size=22, bold=True, color=col)
    add_text(s, 2.95, y + 0.02, 3.50, 0.42, name, size=22, bold=True, color=INK)
    add_text(s, 6.55, y + 0.02, 6.35, 0.42, what, size=22, color=INK)
    y += 0.50

add_text(s, 0.45, 6.82, 12.53, 0.4,
         '→ A群を時系列に読むと、この設問は 2008〜2015 年に一通り済んでいる',
         size=22, bold=True, color=INK)

prs.save('r2_s2a.pptx')
print('r2 stage2a ok:', len(prs.slides._sldIdLst), 'slides')
